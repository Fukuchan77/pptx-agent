"""Performance profiling tests for QA pass operations.

Tests QA pass performance on presentations of varying sizes to ensure
the system meets the <5 seconds target for 10-slide presentations.
"""

import time
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.engine import QAEngine
from pptx_agent.qa.rules.register_defaults import register_default_rules


@pytest.fixture
def qa_engine() -> QAEngine:
    """Create QA engine with all default rules registered."""
    engine = QAEngine()
    register_default_rules()
    return engine


@pytest.fixture
def small_presentation(tmp_path: Path) -> Path:
    """Create a small presentation (5 slides) for performance testing."""
    prs = Presentation()

    # Add 5 slides with content
    for i in range(5):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i + 1}"

        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            if hasattr(content, "text_frame"):
                tf = content.text_frame  # type: ignore[attr-defined]
                tf.text = f"Content for slide {i + 1}"

                # Add bullet points
                for j in range(3):
                    p = tf.add_paragraph()
                    p.text = f"Bullet point {j + 1}"
                    p.level = 0

    output_path = tmp_path / "small_presentation.pptx"
    prs.save(str(output_path))
    return output_path


@pytest.fixture
def medium_presentation(tmp_path: Path) -> Path:
    """Create a medium presentation (10 slides) for performance testing."""
    prs = Presentation()

    # Add 10 slides with content
    for i in range(10):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i + 1}"

        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            if hasattr(content, "text_frame"):
                tf = content.text_frame  # type: ignore[attr-defined]
                tf.text = f"Content for slide {i + 1}"

                # Add bullet points
                for j in range(5):
                    p = tf.add_paragraph()
                    p.text = f"Bullet point {j + 1} with some additional text to make it longer"
                    p.level = 0

    output_path = tmp_path / "medium_presentation.pptx"
    prs.save(str(output_path))
    return output_path


@pytest.fixture
def large_presentation(tmp_path: Path) -> Path:
    """Create a large presentation (30 slides) for performance testing."""
    prs = Presentation()

    # Add 30 slides with content
    for i in range(30):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i + 1}: Performance Testing"

        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            if hasattr(content, "text_frame"):
                tf = content.text_frame  # type: ignore[attr-defined]
                tf.text = f"Content for slide {i + 1}"

                # Add many bullet points
                for j in range(8):
                    p = tf.add_paragraph()
                    p.text = f"Bullet point {j + 1} with substantial text content to simulate real presentations"
                    p.level = 0

    output_path = tmp_path / "large_presentation.pptx"
    prs.save(str(output_path))
    return output_path


def test_qa_pass_performance_small(
    qa_engine: QAEngine,
    small_presentation: Path,
) -> None:
    """Profile QA pass performance on small presentation (5 slides).

    Expected: Should complete in <2 seconds.
    """
    wrapper = PresentationWrapper()
    wrapper.load_template(str(small_presentation))

    start_time = time.perf_counter()
    report = qa_engine.validate(wrapper)
    elapsed = time.perf_counter() - start_time

    # Log performance metrics
    print("\n=== Small Presentation (5 slides) ===")
    print(f"QA Pass Time: {elapsed:.3f}s")
    print(f"Issues Found: {len(report.issues)}")
    print(f"Rules Executed: {len(qa_engine.registry.get_all_rules())}")

    # Performance assertion: should be fast for small presentations
    assert elapsed < 2.0, f"QA pass took {elapsed:.3f}s, expected <2s"


def test_qa_pass_performance_medium(
    qa_engine: QAEngine,
    medium_presentation: Path,
) -> None:
    """Profile QA pass performance on medium presentation (10 slides).

    Expected: Should complete in <5 seconds (constitutional requirement).
    """
    wrapper = PresentationWrapper()
    wrapper.load_template(str(medium_presentation))

    start_time = time.perf_counter()
    report = qa_engine.validate(wrapper)
    elapsed = time.perf_counter() - start_time

    # Log performance metrics
    print("\n=== Medium Presentation (10 slides) ===")
    print(f"QA Pass Time: {elapsed:.3f}s")
    print(f"Issues Found: {len(report.issues)}")
    print(f"Rules Executed: {len(qa_engine.registry.get_all_rules())}")
    print(f"Time per slide: {elapsed / 10:.3f}s")

    # Performance assertion: constitutional requirement
    assert elapsed < 5.0, f"QA pass took {elapsed:.3f}s, expected <5s for 10 slides"


def test_qa_pass_performance_large(
    qa_engine: QAEngine,
    large_presentation: Path,
) -> None:
    """Profile QA pass performance on large presentation (30 slides).

    Expected: Should scale linearly, approximately <15 seconds.
    """
    wrapper = PresentationWrapper()
    wrapper.load_template(str(large_presentation))

    start_time = time.perf_counter()
    report = qa_engine.validate(wrapper)
    elapsed = time.perf_counter() - start_time

    # Log performance metrics
    print("\n=== Large Presentation (30 slides) ===")
    print(f"QA Pass Time: {elapsed:.3f}s")
    print(f"Issues Found: {len(report.issues)}")
    print(f"Rules Executed: {len(qa_engine.registry.get_all_rules())}")
    print(f"Time per slide: {elapsed / 30:.3f}s")

    # Performance assertion: should scale linearly
    # 30 slides should take roughly 3x the time of 10 slides
    assert elapsed < 15.0, f"QA pass took {elapsed:.3f}s, expected <15s for 30 slides"


def test_qa_pass_performance_by_category(
    qa_engine: QAEngine,
    medium_presentation: Path,
) -> None:
    """Profile QA pass performance by rule category.

    Measures time spent in each category (layout, content, style) to identify
    performance bottlenecks.
    """
    wrapper = PresentationWrapper()
    wrapper.load_template(str(medium_presentation))

    categories = ["layout", "content", "style"]
    timings: dict[str, float] = {}

    for category in categories:
        start_time = time.perf_counter()
        qa_engine.validate(wrapper, categories=[category])
        elapsed = time.perf_counter() - start_time
        timings[category] = elapsed

    # Log category timings
    print("\n=== QA Pass Performance by Category ===")
    for category, elapsed in timings.items():
        print(f"{category.capitalize()}: {elapsed:.3f}s")

    total_time = sum(timings.values())
    print(f"Total (sequential): {total_time:.3f}s")

    # Each category should be reasonably fast
    for category, elapsed in timings.items():
        assert elapsed < 3.0, f"{category} checks took {elapsed:.3f}s, expected <3s"


def test_qa_pass_memory_usage(
    qa_engine: QAEngine,
    large_presentation: Path,
) -> None:
    """Profile memory usage during QA pass on large presentation.

    Ensures QA operations don't consume excessive memory.
    """
    import tracemalloc

    wrapper = PresentationWrapper()
    wrapper.load_template(str(large_presentation))

    # Start memory tracking
    tracemalloc.start()

    # Run QA pass
    report = qa_engine.validate(wrapper)

    # Get memory statistics
    current, peak = tracemalloc.get_traced_memory()
    tracemalloc.stop()

    # Log memory metrics
    print("\n=== QA Pass Memory Usage ===")
    print(f"Current Memory: {current / 1024 / 1024:.2f} MB")
    print(f"Peak Memory: {peak / 1024 / 1024:.2f} MB")
    print(f"Issues Found: {len(report.issues)}")

    # Memory assertion: should not exceed 100 MB for QA operations
    assert peak < 100 * 1024 * 1024, f"Peak memory {peak / 1024 / 1024:.2f} MB exceeded 100 MB"


@pytest.mark.parametrize("slide_count", [5, 10, 20, 30])
def test_qa_pass_scalability(
    qa_engine: QAEngine,
    tmp_path: Path,
    slide_count: int,
) -> None:
    """Test QA pass scalability across different presentation sizes.

    Verifies that performance scales linearly with slide count.
    """
    # Create presentation with specified slide count
    prs = Presentation()
    for i in range(slide_count):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i + 1}"

    output_path = tmp_path / f"presentation_{slide_count}slides.pptx"
    prs.save(str(output_path))

    # Profile QA pass
    wrapper = PresentationWrapper()
    wrapper.load_template(str(output_path))
    start_time = time.perf_counter()
    report = qa_engine.validate(wrapper)
    elapsed = time.perf_counter() - start_time

    # Log scalability metrics
    time_per_slide = elapsed / slide_count
    print(f"\n=== Scalability Test ({slide_count} slides) ===")
    print(f"Total Time: {elapsed:.3f}s")
    print(f"Time per Slide: {time_per_slide:.3f}s")
    print(f"Issues Found: {len(report.issues)}")

    # Performance should scale linearly: ~0.5s per slide max
    assert time_per_slide < 0.5, f"Time per slide {time_per_slide:.3f}s exceeded 0.5s"


# Made with Bob
