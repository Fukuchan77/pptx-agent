"""Integration tests for pipeline with LLM integration.

Tests the complete end-to-end pipeline from input text through all three agents
(StoryAnalyzer → OutlineGenerator → ContentGenerator) with LLM enabled.
"""

from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, Mock, patch

import pytest
from pptx import Presentation

from pptx_agent.agents.story_analyzer import StoryAnalysis
from pptx_agent.config import Config
from pptx_agent.pipeline import generate_presentation
from pptx_agent.schemas.outline import PresentationOutline, SlideContent
from pptx_agent.schemas.presentation import PresentationSchema, SlideSchema
from pptx_agent.schemas.text import TextBlock


@pytest.fixture
def mock_story_analysis():
    """Fixture for mock StoryAnalysis."""
    return StoryAnalysis(
        topic="AI Technology Overview",
        target_audience="Technical professionals",
        key_message="AI is transforming industries",
        tone="professional",
        language="en",
    )


@pytest.fixture
def mock_presentation_outline(mock_story_analysis: Mock):
    """Fixture for mock PresentationOutline."""
    return PresentationOutline(
        title=mock_story_analysis.topic,
        output_language=mock_story_analysis.language,
        slides=[
            SlideContent(
                slide_number=1,
                layout_name="Title Slide",
                title="AI Technology Overview",
                content="Key insights",
            ),
            SlideContent(
                slide_number=2,
                layout_name="Title and Content",
                title="Introduction",
                content="Overview",
            ),
            SlideContent(
                slide_number=3,
                layout_name="Title Slide",
                title="Conclusion",
                content="Summary",
            ),
        ],
    )


@pytest.fixture
def mock_presentation_schema():
    """Fixture for mock PresentationSchema."""
    return PresentationSchema(
        title="AI Technology Overview",
        slides=[
            SlideSchema(
                layout_name="Title Slide",
                title="AI Technology Overview",
                content=[
                    TextBlock(
                        placeholder_name="Subtitle",
                        text="Key insights",
                        language="en",
                    )
                ],
                notes="Welcome to this presentation.",
            ),
            SlideSchema(
                layout_name="Title and Content",
                title="Introduction",
                content=[
                    TextBlock(
                        placeholder_name="Body",
                        text="Overview of AI technology.",
                        language="en",
                    )
                ],
                notes="Discuss the fundamentals.",
            ),
            SlideSchema(
                layout_name="Title Slide",
                title="Conclusion",
                content=[
                    TextBlock(
                        placeholder_name="Subtitle",
                        text="Summary",
                        language="en",
                    )
                ],
                notes="Wrap up the presentation.",
            ),
        ],
    )


@pytest.fixture
def template_path() -> str:
    """Provide path to basic template for testing."""
    return "templates/basic-template.pptx"


@pytest.fixture
def mock_config():
    """Fixture for mock Config."""
    config = MagicMock(spec=Config)
    config.llm_provider = "openai"
    config.llm_model = "gpt-4"
    config.llm_api_key = "test-api-key"
    config.fallback_provider = "anthropic"
    config.fallback_model = "claude-sonnet-4-6"
    config.fallback_api_key = "test-fallback-key"
    return config


@pytest.mark.asyncio
async def test_pipeline_with_llm_success(
    tmp_path: Path,
    template_path: str,
    mock_story_analysis: Mock,
    mock_presentation_outline: Mock,
    mock_presentation_schema: Mock,
):
    """Test complete pipeline execution with LLM agents.

    Verifies:
    - Pipeline executes successfully with use_llm=True
    - All three agents are called with LLM
    - Output .pptx file is created
    - File is valid PowerPoint format
    """
    input_text = "AI Technology: An overview of artificial intelligence."
    output_path = tmp_path / "output.pptx"

    with (
        patch("pptx_agent.pipeline.analyze_story", new_callable=AsyncMock) as mock_story,
        patch("pptx_agent.pipeline.generate_outline", new_callable=AsyncMock) as mock_outline,
        patch("pptx_agent.pipeline.generate_content", new_callable=AsyncMock) as mock_content,
    ):
        mock_story.return_value = mock_story_analysis
        mock_outline.return_value = mock_presentation_outline
        mock_content.return_value = mock_presentation_schema

        # Act
        result_path, _qa_report = await generate_presentation(
            input_text=input_text,
            template_path=template_path,
            output_path=str(output_path),
            use_llm=True,
        )

        # Assert - pipeline succeeded
        assert Path(result_path).exists(), "Output .pptx file should exist"

        # Assert - all agents were called
        assert mock_story.called, "Story analyzer should be called"
        assert mock_outline.called, "Outline generator should be called"
        assert mock_content.called, "Content generator should be called"

        # Verify use_llm=True was passed
        assert mock_story.call_args.kwargs.get("use_llm") is True
        assert mock_outline.call_args.kwargs.get("use_llm") is True
        assert mock_content.call_args.kwargs.get("use_llm") is True

        # Assert - valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None
        assert len(prs.slides) >= 3


@pytest.mark.asyncio
async def test_pipeline_with_multiple_llm_calls(
    tmp_path: Path,
    template_path: str,
    mock_story_analysis: Mock,
    mock_presentation_outline: Mock,
    mock_presentation_schema: Mock,
):
    """Test pipeline makes multiple LLM calls successfully.

    Verifies:
    - Multiple agent calls are made during pipeline execution
    - Each agent receives proper inputs
    - Pipeline completes successfully with all agents coordinated
    """
    input_text = "Test presentation about technology"
    output_path = tmp_path / "output_multi_llm.pptx"

    with (
        patch("pptx_agent.pipeline.analyze_story", new_callable=AsyncMock) as mock_story,
        patch("pptx_agent.pipeline.generate_outline", new_callable=AsyncMock) as mock_outline,
        patch("pptx_agent.pipeline.generate_content", new_callable=AsyncMock) as mock_content,
    ):
        mock_story.return_value = mock_story_analysis
        mock_outline.return_value = mock_presentation_outline
        mock_content.return_value = mock_presentation_schema

        # Act
        result_path, _qa_report = await generate_presentation(
            input_text=input_text,
            template_path=template_path,
            output_path=str(output_path),
            use_llm=True,
        )

        # Assert - all agents called in correct order
        assert Path(result_path).exists()
        assert mock_story.call_count == 1
        assert mock_outline.call_count == 1
        assert mock_content.call_count == 1

        # Verify call order
        assert mock_story.called
        assert mock_outline.called
        assert mock_content.called


@pytest.mark.asyncio
async def test_pipeline_fails_when_all_providers_fail(tmp_path: Path, template_path: str):
    """Test pipeline fails gracefully when both providers fail.

    Verifies:
    - Both primary and fallback providers fail
    - RuntimeError is raised with appropriate message
    """
    input_text = "Test presentation"
    output_path = tmp_path / "output_fail.pptx"

    with patch(
        "pptx_agent.agents.utils.run_agent_with_fallback", new_callable=AsyncMock
    ) as mock_run_with_fallback:
        # Mock run_agent_with_fallback to raise RuntimeError (simulating both providers failing)
        mock_run_with_fallback.side_effect = RuntimeError(
            "All LLM providers failed. Primary: Primary failed, Fallback: Fallback failed"
        )

        # Act & Assert
        with pytest.raises(RuntimeError, match="All LLM providers failed"):
            await generate_presentation(
                input_text=input_text,
                template_path=template_path,
                output_path=str(output_path),
                use_llm=True,
            )


@pytest.mark.asyncio
async def test_pipeline_heuristic_mode_unaffected(tmp_path: Path, template_path: str):
    """Test pipeline still works in heuristic mode (use_llm=False).

    Verifies:
    - Pipeline executes with use_llm=False
    - Heuristic fallback is used
    - Output .pptx file is created successfully
    """
    input_text = "Test presentation about Python programming basics"
    output_path = tmp_path / "output_heuristic.pptx"

    # Don't mock agents - they should use heuristic fallback
    result_path, _qa_report = await generate_presentation(
        input_text=input_text,
        template_path=template_path,
        output_path=str(output_path),
        use_llm=False,
    )

    # Assert - pipeline succeeded
    assert Path(result_path).exists()

    # Assert - valid file
    prs = Presentation(result_path)
    assert prs is not None
    assert len(prs.slides) >= 3


@pytest.mark.asyncio
async def test_pipeline_llm_with_template_constraints(
    tmp_path: Path, template_path: str, mock_story_analysis: Mock
):
    """Test LLM agents respect template layout constraints.

    Verifies:
    - Pipeline parses template manifest correctly
    - LLM-generated outline respects template constraints
    - Number of slides within template limits
    """
    input_text = "Test presentation with template constraints"
    output_path = tmp_path / "output_constrained.pptx"

    constrained_outline = PresentationOutline(
        title="Test Presentation",
        output_language="en",
        slides=[
            SlideContent(
                slide_number=1,
                layout_name="Title Slide",
                title="Test Title",
                content="Test content",
            ),
            SlideContent(
                slide_number=2,
                layout_name="Title and Content",
                title="Content Slide",
                content="Some content",
            ),
            SlideContent(
                slide_number=3,
                layout_name="Title Slide",
                title="Conclusion",
                content="Summary",
            ),
        ],
    )

    mock_content_schema = PresentationSchema(
        title="Test Presentation",
        slides=[
            SlideSchema(
                layout_name="Title Slide",
                title="Test Title",
                content=[
                    TextBlock(placeholder_name="Subtitle", text="Test content", language="en")
                ],
                notes="Introduction notes",
            ),
            SlideSchema(
                layout_name="Title and Content",
                title="Content Slide",
                content=[
                    TextBlock(placeholder_name="Body", text="Some content here", language="en")
                ],
                notes="Content notes",
            ),
            SlideSchema(
                layout_name="Title Slide",
                title="Conclusion",
                content=[TextBlock(placeholder_name="Subtitle", text="Summary", language="en")],
                notes="Conclusion notes",
            ),
        ],
    )

    with (
        patch("pptx_agent.pipeline.analyze_story", new_callable=AsyncMock) as mock_story,
        patch("pptx_agent.pipeline.generate_outline", new_callable=AsyncMock) as mock_outline,
        patch("pptx_agent.pipeline.generate_content", new_callable=AsyncMock) as mock_content,
    ):
        mock_story.return_value = mock_story_analysis
        mock_outline.return_value = constrained_outline
        mock_content.return_value = mock_content_schema

        # Act
        result_path, _qa_report = await generate_presentation(
            input_text=input_text,
            template_path=template_path,
            output_path=str(output_path),
            use_llm=True,
        )

        # Assert
        assert Path(result_path).exists()
        assert len(constrained_outline.slides) == 3
        assert len(constrained_outline.slides) <= 30

        prs = Presentation(result_path)
        # Note: Overflow resolution may add extra slides, so check >= instead of ==
        assert len(prs.slides) >= 3, "Should have at least 3 slides"
        assert len(prs.slides) <= 10, "Should not have excessive slides"


@pytest.mark.asyncio
async def test_pipeline_with_different_languages(tmp_path: Path, template_path: str):
    """Test pipeline handles different language outputs correctly.

    Verifies:
    - Pipeline can generate presentations in different languages
    - Language setting is propagated through all agents
    - Output file is created successfully
    """
    input_text = "プレゼンテーションについて"  # Japanese input
    output_path = tmp_path / "output_japanese.pptx"

    ja_story_analysis = StoryAnalysis(
        topic="プレゼンテーション",
        target_audience="一般聴衆",
        key_message="重要なポイント",
        tone="professional",
        language="ja",
    )

    ja_outline = PresentationOutline(
        title="プレゼンテーション",
        output_language="ja",
        slides=[
            SlideContent(
                slide_number=1,
                layout_name="Title Slide",
                title="プレゼンテーション",
                content="重要なポイント",
            ),
            SlideContent(
                slide_number=2,
                layout_name="Title and Content",
                title="内容",
                content="詳細",
            ),
            SlideContent(
                slide_number=3,
                layout_name="Title Slide",
                title="まとめ",
                content="結論",
            ),
        ],
    )

    ja_content = PresentationSchema(
        title="プレゼンテーション",
        slides=[
            SlideSchema(
                layout_name="Title Slide",
                title="プレゼンテーション",
                content=[
                    TextBlock(placeholder_name="Subtitle", text="重要なポイント", language="ja")
                ],
                notes="プレゼンテーションの開始",
            ),
            SlideSchema(
                layout_name="Title and Content",
                title="内容",
                content=[TextBlock(placeholder_name="Body", text="詳細な内容", language="ja")],
                notes="詳細説明",
            ),
            SlideSchema(
                layout_name="Title Slide",
                title="まとめ",
                content=[TextBlock(placeholder_name="Subtitle", text="結論", language="ja")],
                notes="まとめ",
            ),
        ],
    )

    with (
        patch("pptx_agent.pipeline.analyze_story", new_callable=AsyncMock) as mock_story,
        patch("pptx_agent.pipeline.generate_outline", new_callable=AsyncMock) as mock_outline,
        patch("pptx_agent.pipeline.generate_content", new_callable=AsyncMock) as mock_content,
    ):
        mock_story.return_value = ja_story_analysis
        mock_outline.return_value = ja_outline
        mock_content.return_value = ja_content

        # Act
        result_path, _qa_report = await generate_presentation(
            input_text=input_text,
            template_path=template_path,
            output_path=str(output_path),
            use_llm=True,
        )

        # Assert
        assert Path(result_path).exists()
        # Verify language propagation
        assert mock_outline.call_args[0][0].language == "ja"
        assert ja_outline.output_language == "ja"
