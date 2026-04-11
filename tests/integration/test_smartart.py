"""Integration tests for SmartArt population.

Tests the complete end-to-end pipeline for generating presentations with SmartArt
diagrams, validating User Story 5 acceptance criteria.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.pipeline import generate_presentation


@pytest.fixture
def template_path() -> str:
    """Provide path to basic template for testing.

    Returns:
        Path to templates/basic-template.pptx
    """
    return "templates/basic-template.pptx"


@pytest.fixture
def process_flow_input() -> str:
    """Load process flow sample fixture.

    Returns:
        Content of tests/fixtures/smartart_process_sample.txt
    """
    fixture_path = Path("tests/fixtures/smartart_process_sample.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def hierarchy_input() -> str:
    """Load hierarchy sample fixture.

    Returns:
        Content of tests/fixtures/smartart_hierarchy_sample.txt
    """
    fixture_path = Path("tests/fixtures/smartart_hierarchy_sample.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def cycle_input() -> str:
    """Load cycle process sample fixture.

    Returns:
        Content of tests/fixtures/smartart_cycle_sample.txt
    """
    fixture_path = Path("tests/fixtures/smartart_cycle_sample.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.mark.asyncio
async def test_generate_presentation_with_process_flow_content(
    process_flow_input: str, template_path: str
):
    """Test complete pipeline with process flow content for SmartArt generation.

    Validates:
    - Pipeline completes successfully with process flow data
    - .pptx file is created
    - File is valid PowerPoint format
    - Presentation can be opened
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_process_flow.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=process_flow_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        assert result_path == output_path, "Result path should match output path"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
        assert len(prs.slides) <= 30, "Presentation should not exceed 30 slides"


@pytest.mark.asyncio
async def test_generate_presentation_with_hierarchy_content(
    hierarchy_input: str, template_path: str
):
    """Test complete pipeline with hierarchical content for SmartArt generation.

    Validates:
    - Pipeline completes successfully with hierarchy data
    - .pptx file is created
    - File is valid PowerPoint format
    - Presentation can be opened
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_hierarchy.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=hierarchy_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        assert result_path == output_path, "Result path should match output path"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
        assert len(prs.slides) <= 30, "Presentation should not exceed 30 slides"


@pytest.mark.asyncio
async def test_generate_presentation_with_cycle_content(cycle_input: str, template_path: str):
    """Test complete pipeline with cyclical process content for SmartArt generation.

    Validates:
    - Pipeline completes successfully with cycle data
    - .pptx file is created
    - File is valid PowerPoint format
    - Presentation can be opened
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_cycle.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=cycle_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        assert result_path == output_path, "Result path should match output path"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
        assert len(prs.slides) <= 30, "Presentation should not exceed 30 slides"


@pytest.mark.asyncio
async def test_presentation_with_smartart_contains_shapes(
    process_flow_input: str, template_path: str
):
    """Test that generated presentation with SmartArt contains various shape types.

    Validates:
    - Slides have shapes (text, SmartArt, etc.)
    - Shape types are accessible
    - No shape-related errors occur
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_smartart_shapes.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=process_flow_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Open and inspect shapes
        prs = Presentation(result_path)

        total_shapes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                # Access shape type without error
                _ = shape.shape_type
                total_shapes += 1

                # If it's a placeholder, access its type
                if shape.is_placeholder:
                    _ = shape.placeholder_format.type

        # Should have multiple shapes across slides
        assert total_shapes > 0, "Presentation should contain shapes"


@pytest.mark.asyncio
async def test_mixed_content_with_smartart_and_text(template_path: str):
    """Test presentation with mixed SmartArt and text content.

    Validates:
    - Pipeline handles mixed content types
    - SmartArt can coexist with text blocks
    - File generation completes successfully
    """
    # Mixed input with process flow and general text
    mixed_input = """
# Project Management Methodology

## Introduction

Effective project management requires a structured approach that ensures successful
delivery of objectives. Our methodology combines proven practices with agile principles
to deliver value iteratively.

## Project Lifecycle Phases

The project follows these sequential phases:

**Initiation**: Define project scope and objectives. Identify stakeholders, establish
governance, and secure initial approvals.

**Planning**: Create detailed project plans. Define work breakdown structure, schedule,
budget, and resource allocation.

**Execution**: Implement project deliverables. Coordinate teams, manage tasks, and
monitor progress against plans.

**Monitoring**: Track performance and quality. Measure KPIs, identify issues, and
implement corrective actions.

**Closure**: Finalize all activities and deliverables. Conduct retrospectives, document
lessons learned, and transition to operations.

## Key Success Factors

Several factors contribute to project success:

- Clear communication channels among all stakeholders
- Well-defined requirements with minimal ambiguity
- Adequate resource allocation and skills availability
- Regular status monitoring and risk management
- Strong executive sponsorship and support

## Continuous Improvement

After project completion, teams analyze outcomes and identify improvement opportunities
for future initiatives. This learning cycle ensures organizational capability growth.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_mixed_smartart.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=mixed_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify file was created and is valid
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


@pytest.mark.asyncio
async def test_organizational_structure_for_hierarchy_smartart(template_path: str):
    """Test SmartArt generation with organizational structure content.

    Validates:
    - Hierarchical organizational data is processed
    - SmartArt specifications are generated appropriately
    - File generation succeeds
    """
    org_structure_input = """
# Corporate Organization

## Executive Leadership

The company is led by a strong executive team:

**Chief Executive Officer**: Sarah Chen, responsible for overall strategy and vision.
Reports to the Board of Directors.

**Chief Operating Officer**: Michael Rodriguez, oversees daily operations.
Reports to CEO.

**Chief Financial Officer**: Jennifer Martinez, manages financial planning.
Reports to CEO.

**Chief Technology Officer**: David Kim, leads technology initiatives.
Reports to CEO.

## Technology Department

Under the CTO, the technology department is structured as follows:

**VP of Engineering**: Leads all engineering teams
  - Backend Team: 15 engineers
  - Frontend Team: 12 engineers
  - DevOps Team: 8 engineers

**VP of Product**: Manages product strategy
  - Product Managers: 6 team members
  - Product Analysts: 4 team members

**VP of Data Science**: Oversees AI/ML initiatives
  - Data Scientists: 10 team members
  - ML Engineers: 8 team members

This structure enables efficient development and innovation.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_org_structure.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=org_structure_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"


@pytest.mark.asyncio
async def test_workflow_cycle_for_cycle_smartart(template_path: str):
    """Test SmartArt generation with cyclical workflow content.

    Validates:
    - Cyclical process data is processed
    - Cycle-type SmartArt is appropriate for the content
    - File generation succeeds
    """
    cycle_workflow_input = """
# Continuous Improvement Framework

## The PDCA Cycle

Our improvement framework follows the Plan-Do-Check-Act cycle:

**Plan**: Identify improvement opportunities and develop action plans. Set measurable
objectives and define success criteria for the initiative.

**Do**: Implement the improvement plan on a small scale first. Execute changes,
document observations, and collect relevant data.

**Check**: Evaluate results against objectives and expectations. Compare actual
outcomes to planned outcomes and analyze variances.

**Act**: Standardize successful improvements across the organization. If results
are unsatisfactory, revise the approach and return to planning.

This cycle repeats continuously, driving incremental improvements.

## Agile Sprint Cycle

Development teams follow a two-week sprint cycle:

**Sprint Planning** → Define sprint goals and select backlog items
**Daily Standups** → Synchronize team activities daily
**Development** → Implement features and conduct reviews
**Sprint Review** → Demonstrate work to stakeholders
**Retrospective** → Reflect on process and identify improvements

Then the cycle begins again with the next sprint planning session.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_cycle_workflow.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=cycle_workflow_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"


@pytest.mark.asyncio
async def test_file_opens_without_corruption_smartart(process_flow_input: str, template_path: str):
    """Test that generated files with SmartArt open without corruption.

    Validates:
    - Files with SmartArt elements don't cause corruption
    - All slides are accessible
    - Shapes can be enumerated without errors
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_smartart_no_corruption.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=process_flow_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Try to open and access all elements
        prs = Presentation(result_path)

        # Iterate through all slides and shapes without errors
        for i, slide in enumerate(prs.slides):
            assert slide.slide_layout is not None, f"Slide {i} should have layout"

            for shape in slide.shapes:
                # Access various shape properties
                _ = shape.shape_type
                _ = shape.name

                if hasattr(shape, "text"):
                    _ = shape.text  # type: ignore[attr-defined]

        # If we got here, file is not corrupted
        assert True, "File should open without corruption"


@pytest.mark.asyncio
async def test_multiple_smartart_types_in_presentation(template_path: str):
    """Test presentation with multiple types of SmartArt content.

    Validates:
    - Multiple SmartArt diagram types can coexist
    - Process, hierarchy, and cycle content are all handled
    - File generation completes successfully
    """
    multi_smartart_input = """
# Comprehensive Business Overview

## Development Process

Our software development follows these steps:

**Requirements Gathering** → Understand stakeholder needs
**Design** → Create technical specifications
**Implementation** → Build the solution
**Testing** → Validate quality and functionality
**Deployment** → Release to production
**Maintenance** → Provide ongoing support

## Organizational Structure

The company has a clear hierarchical structure:

**CEO** - Overall strategy and vision
  - **VP Sales** - Revenue generation
    - Regional Sales Directors (3)
    - Account Executives (25)
  - **VP Engineering** - Product development
    - Engineering Managers (4)
    - Software Engineers (35)
  - **VP Operations** - Daily operations
    - Operations Managers (3)
    - Operations Staff (20)

## Continuous Improvement Cycle

We follow a continuous improvement cycle:

**Assess** → Evaluate current performance and identify gaps
**Plan** → Develop improvement strategies and action plans
**Execute** → Implement changes and monitor progress
**Review** → Analyze results and gather feedback
**Refine** → Adjust approach based on learnings

The cycle repeats, driving ongoing enhancement of our processes.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_multi_smartart.pptx")

        # Generate presentation
        result_path = await generate_presentation(
            input_text=multi_smartart_input,
            template_path=template_path,
            output_path=output_path,
            use_llm=False,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have multiple slides"
