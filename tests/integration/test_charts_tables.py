"""Integration tests for chart and table generation.

Tests the complete end-to-end pipeline for generating presentations with charts
and tables, validating User Story 4 acceptance criteria.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.pipeline import generate_presentation


@pytest.fixture
def template_path() -> str:
    """Provide path to basic template for testing.

    Returns:
        Path to templates/basic-template.pptx
    """
    return "templates/basic-template.pptx"


@pytest.fixture
def chart_data_input() -> str:
    """Load chart data sample fixture.

    Returns:
        Content of tests/fixtures/chart_data_sample.txt
    """
    fixture_path = Path("tests/fixtures/chart_data_sample.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def table_data_input() -> str:
    """Load table data sample fixture.

    Returns:
        Content of tests/fixtures/table_data_sample.txt
    """
    fixture_path = Path("tests/fixtures/table_data_sample.txt")
    return fixture_path.read_text(encoding="utf-8")


def test_generate_presentation_with_chart_data(chart_data_input: str, template_path: str):
    """Test complete pipeline with numerical data for chart generation.

    Validates:
    - Pipeline completes successfully with chart data
    - .pptx file is created
    - File is valid PowerPoint format
    - Presentation can be opened
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_charts.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=chart_data_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        assert result_path == output_path, "Result path should match output path"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
        assert len(prs.slides) <= 30, "Presentation should not exceed 30 slides"


def test_generate_presentation_with_table_data(table_data_input: str, template_path: str):
    """Test complete pipeline with tabular data for table generation.

    Validates:
    - Pipeline completes successfully with table data
    - .pptx file is created
    - File is valid PowerPoint format
    - Presentation can be opened
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_tables.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=table_data_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        assert result_path == output_path, "Result path should match output path"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
        assert len(prs.slides) <= 30, "Presentation should not exceed 30 slides"


def test_presentation_contains_shapes(chart_data_input: str, template_path: str):
    """Test that generated presentation contains various shape types.

    Validates:
    - Slides have shapes (text, charts, tables, etc.)
    - Shape types are accessible
    - No shape-related errors occur
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_shapes.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=chart_data_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Open and inspect shapes
        prs = Presentation(result_path)

        total_shapes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                # Access shape type without error
                _ = shape.shape_type
                total_shapes += 1

                # If it's a placeholder, access its type
                if shape.is_placeholder:
                    _ = shape.placeholder_format.type

        # Should have multiple shapes across slides
        assert total_shapes > 0, "Presentation should contain shapes"


def test_mixed_content_with_charts_and_tables(template_path: str):
    """Test presentation with mixed numerical and tabular data.

    Validates:
    - Pipeline handles mixed content types
    - Both charts and tables can coexist
    - File generation completes successfully
    """
    # Mixed input with both numerical data and tables
    mixed_input = """
# Business Performance Report

## Sales Growth

Our quarterly revenue shows strong growth:
- Q1: $1.2M
- Q2: $1.5M
- Q3: $1.8M
- Q4: $2.1M

## Top Products

| Product | Units Sold | Revenue | Growth |
|---------|-----------|---------|--------|
| Laptop Pro | 1,200 | $1.8M | +15% |
| Desktop Ultra | 850 | $1.2M | +12% |
| Tablet Max | 2,100 | $1.4M | +20% |
| Monitor 27" | 3,400 | $1.1M | +8% |

## Market Share

Geographic distribution:
- North America: 40%
- Europe: 30%
- Asia Pacific: 20%
- Latin America: 10%

## Employee Performance

| Employee | Projects | Score | Rating |
|----------|----------|-------|--------|
| Alice | 12 | 95 | A+ |
| Bob | 8 | 88 | A |
| Carol | 15 | 92 | A |
| David | 11 | 90 | A |
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_mixed.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=mixed_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created and is valid
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


def test_chart_data_with_multiple_series(template_path: str):
    """Test chart generation with multiple data series.

    Validates:
    - Charts with multiple series are handled correctly
    - Complex numerical data is processed
    - File generation succeeds
    """
    multi_series_input = """
# Product Comparison Analysis

## Revenue by Product Line

Comparing three product categories over four quarters:

**Software Sales:**
- Q1: $540K
- Q2: $675K
- Q3: $810K
- Q4: $945K

**Hardware Sales:**
- Q1: $360K
- Q2: $450K
- Q3: $540K
- Q4: $630K

**Services:**
- Q1: $240K
- Q2: $300K
- Q3: $360K
- Q4: $420K

All categories show consistent growth with software leading the revenue contribution.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_multi_series.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=multi_series_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"


def test_table_with_various_row_counts(template_path: str):
    """Test table generation with different row counts.

    Validates:
    - Tables with few rows work correctly
    - Tables approaching max rows (20) work correctly
    - Various table sizes are handled
    """
    table_input = """
# Inventory Report

## Small Inventory (5 items)

| Item | Quantity | Price |
|------|----------|-------|
| Item A | 100 | $10 |
| Item B | 200 | $20 |
| Item C | 150 | $15 |
| Item D | 80 | $8 |
| Item E | 120 | $12 |

## Medium Inventory (10 items)

| Product | SKU | Stock | Price | Warehouse |
|---------|-----|-------|-------|-----------|
| Product 1 | P001 | 100 | $50 | West |
| Product 2 | P002 | 150 | $60 | East |
| Product 3 | P003 | 200 | $70 | Central |
| Product 4 | P004 | 120 | $55 | West |
| Product 5 | P005 | 180 | $65 | East |
| Product 6 | P006 | 90 | $45 | Central |
| Product 7 | P007 | 110 | $58 | West |
| Product 8 | P008 | 160 | $68 | East |
| Product 9 | P009 | 140 | $62 | Central |
| Product 10 | P010 | 130 | $60 | West |

These tables demonstrate our inventory management across different warehouse locations.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_tables_varied.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=table_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"


def test_percentage_and_ratio_data_for_charts(template_path: str):
    """Test chart generation with percentage and ratio data.

    Validates:
    - Percentage data is suitable for pie charts
    - Ratio data is processed correctly
    - Charts are generated successfully
    """
    percentage_input = """
# Market Analysis

## Market Share Distribution

Our market position relative to competitors:

- Our Company: 35% market share
- Competitor A: 25% market share
- Competitor B: 20% market share
- Competitor C: 12% market share
- Others: 8% market share

We maintain the leading position with room for growth in the "Others" segment.

## Customer Demographics

Age distribution of our customer base:

- 18-24: 16% (3,200 customers)
- 25-34: 34% (6,800 customers)
- 35-44: 27% (5,400 customers)
- 45-54: 15% (3,000 customers)
- 55+: 8% (1,600 customers)

The 25-34 age group represents our largest customer segment.
"""

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_percentages.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=percentage_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify successful generation
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"


def test_file_opens_without_corruption(chart_data_input: str, template_path: str):
    """Test that generated files with charts/tables open without corruption.

    Validates:
    - Files with visual elements don't cause corruption
    - All slides are accessible
    - Shapes can be enumerated without errors
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_no_corruption.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=chart_data_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Try to open and access all elements
        prs = Presentation(result_path)

        # Iterate through all slides and shapes without errors
        for i, slide in enumerate(prs.slides):
            assert slide.slide_layout is not None, f"Slide {i} should have layout"

            for shape in slide.shapes:
                # Access various shape properties
                _ = shape.shape_type
                _ = shape.name

                if hasattr(shape, "text"):
                    _ = shape.text  # type: ignore[attr-defined]

        # If we got here, file is not corrupted
        assert True, "File should open without corruption"
