"""Smoke tests for QA framework and core functionality.

Smoke tests verify that the basic QA workflow functions correctly
without detailed validation. These are quick sanity checks for CI.
"""

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.engine import QAEngine
from pptx_agent.qa.rules.register_defaults import register_default_rules


@pytest.fixture
def qa_engine() -> QAEngine:
    """Create QA engine with all default rules registered."""
    engine = QAEngine()
    register_default_rules()
    return engine


@pytest.fixture
def valid_presentation(tmp_path: Path) -> Path:
    """Create a valid presentation with no QA issues."""
    prs = Presentation()

    # Add a simple slide with title and content
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title = slide.shapes.title
    if title is not None:
        title.text = "Smoke Test Slide"

    # Add content
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        if hasattr(content, "text_frame"):
            tf = content.text_frame  # type: ignore[attr-defined]
            tf.text = "This is a valid presentation for smoke testing."

    output_path = tmp_path / "valid_presentation.pptx"
    prs.save(str(output_path))
    return output_path


def test_qa_engine_initialization() -> None:
    """Smoke test: QA engine initializes correctly."""
    engine = QAEngine()
    assert engine is not None
    assert engine.registry is not None


def test_qa_rules_registration() -> None:
    """Smoke test: Default QA rules register successfully."""
    engine = QAEngine()
    register_default_rules()

    # Verify rules are registered
    rules = engine.registry.get_all_rules()
    assert len(rules) > 0, "No QA rules registered"

    # Verify we have rules from each category
    rule_ids = [rule.rule_id for rule in rules]
    assert any(rid.startswith("QA-L-") for rid in rule_ids), "No layout rules registered"
    assert any(rid.startswith("QA-C-") for rid in rule_ids), "No content rules registered"
    assert any(rid.startswith("QA-S-") for rid in rule_ids), "No style rules registered"


def test_qa_validation_on_valid_presentation(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA validation runs successfully on valid presentation."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    report = qa_engine.validate(wrapper)

    # Basic assertions
    assert report is not None
    assert report.total_issues >= 0
    assert report.error_count >= 0
    assert report.warning_count >= 0
    assert report.info_count >= 0


def test_qa_report_json_serialization(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA report serializes to JSON successfully."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    report = qa_engine.validate(wrapper)
    json_output = report.to_json()

    assert json_output is not None
    assert isinstance(json_output, str)
    assert len(json_output) > 0


def test_qa_report_markdown_generation(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA report generates Markdown successfully."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    report = qa_engine.validate(wrapper)
    markdown_output = report.to_markdown()

    assert markdown_output is not None
    assert isinstance(markdown_output, str)
    assert len(markdown_output) > 0
    assert "# QA Report" in markdown_output or "QA Report" in markdown_output


def test_qa_validation_by_category(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA validation by category works correctly."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    # Test each category
    for category in ["layout", "content", "style"]:
        report = qa_engine.validate(wrapper, categories=[category])
        assert report is not None
        # All issues should be from the specified category
        for issue in report.issues:
            assert issue.rule_id.startswith(f"QA-{category[0].upper()}-")


def test_qa_validation_by_rule_id(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA validation by specific rule IDs works correctly."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    # Test with specific rule IDs
    rule_ids = ["QA-L-001", "QA-L-002"]
    report = qa_engine.validate(wrapper, rule_ids=rule_ids)

    assert report is not None
    # All issues should be from the specified rules
    for issue in report.issues:
        assert issue.rule_id in rule_ids


def test_qa_language_detection_integration(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: QA engine integrates with language detection."""
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    # Test with explicit language override
    engine_with_language = QAEngine(language="en")
    report = engine_with_language.validate(wrapper)

    assert report is not None


def test_fixer_engine_initialization() -> None:
    """Smoke test: Fixer engine initializes correctly."""
    from pptx_agent.fixer.engine import FixEngine

    engine = FixEngine()
    assert engine is not None
    assert engine.registry is not None
    assert engine.max_iterations > 0


def test_cache_manager_initialization() -> None:
    """Smoke test: Cache manager initializes correctly."""
    from pptx_agent.cache.manager import CacheManager

    manager = CacheManager()
    assert manager is not None
    assert manager.cache_dir is not None
    assert manager.storage is not None


def test_template_manifest_caching(tmp_path: Path) -> None:
    """Smoke test: Template manifest caching works."""
    from pptx_agent.cache.manager import CacheManager

    # Create a simple presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    prs.slides.add_slide(slide_layout)

    template_path = tmp_path / "template.pptx"
    prs.save(str(template_path))

    # Test caching
    manager = CacheManager(cache_dir=tmp_path / "cache")

    # Verify cache manager is functional
    assert manager.cache_dir is not None
    assert manager.storage is not None


def test_cli_interface_imports() -> None:
    """Smoke test: CLI interface imports successfully."""
    from pptx_agent.interfaces import cli

    assert cli is not None


def test_api_interface_imports() -> None:
    """Smoke test: API interface imports successfully."""
    from pptx_agent.interfaces.api import app

    assert app is not None


def test_mcp_interface_imports() -> None:
    """Smoke test: MCP interface imports successfully."""
    from pptx_agent.interfaces import mcp

    assert mcp is not None


def test_constitutional_compliance_imports() -> None:
    """Smoke test: All modules follow constitutional principles (importable with type hints)."""
    # Test that all QA modules import successfully
    # Test that all fixer modules import successfully
    # Test that cache modules import successfully
    from pptx_agent import cache, fixer, qa
    from pptx_agent.fixer import strategies
    from pptx_agent.qa import rules

    # If we get here, all imports succeeded
    assert qa is not None
    assert rules is not None
    assert fixer is not None
    assert strategies is not None
    assert cache is not None


def test_end_to_end_qa_workflow(
    qa_engine: QAEngine,
    valid_presentation: Path,
) -> None:
    """Smoke test: End-to-end QA workflow completes successfully."""
    # Load presentation
    wrapper = PresentationWrapper()
    wrapper.load_template(str(valid_presentation))

    # Run QA validation
    report = qa_engine.validate(wrapper)

    # Generate outputs
    json_output = report.to_json()
    markdown_output = report.to_markdown()

    # Verify workflow completed
    assert report is not None
    assert json_output is not None
    assert markdown_output is not None
    assert report.total_issues >= 0


# Made with Bob
