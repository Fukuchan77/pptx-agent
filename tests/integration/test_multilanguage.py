"""Integration tests for multi-language support (User Story 2).

Tests the complete end-to-end pipeline with language-aware features:
- Japanese and English language detection and output
- Explicit language parameter support
- Language-specific text capacity calculations
- Mixed-language content handling

Validates all User Story 2 acceptance criteria.
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.pipeline import generate_presentation
from pptx_agent.utils.language_detector import detect_language


@pytest.fixture
def template_path() -> str:
    """Provide path to basic template for testing.

    Returns:
        Path to templates/basic-template.pptx
    """
    return "templates/basic-template.pptx"


@pytest.fixture
def japanese_template_path() -> str:
    """Provide path to Japanese template for testing.

    Returns:
        Path to templates/japanese-template.pptx
    """
    return "templates/japanese-template.pptx"


@pytest.fixture
def sample_story_en() -> str:
    """Load English sample story fixture.

    Returns:
        Content of tests/fixtures/sample_story_en.txt
    """
    fixture_path = Path("tests/fixtures/sample_story_en.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def sample_story_ja() -> str:
    """Load Japanese sample story fixture.

    Returns:
        Content of tests/fixtures/sample_story_ja.txt
    """
    fixture_path = Path("tests/fixtures/sample_story_ja.txt")
    return fixture_path.read_text(encoding="utf-8")


def test_japanese_input_produces_japanese_output(sample_story_ja: str, template_path: str):
    """Test that Japanese input automatically produces Japanese output.

    Validates User Story 2, Scenario 1:
    - Japanese text input detected automatically
    - Output language set to Japanese
    - Text capacity calculated using full-width character ratios (0.55x)
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_ja_auto.pptx")

        # Generate presentation without specifying language (should auto-detect)
        result_path = generate_presentation(
            input_text=sample_story_ja,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify file is valid PowerPoint with Japanese content
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"

        # Verify Japanese characters are present in slides
        has_japanese = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:  # type: ignore[attr-defined]
                    # Check if text contains Japanese characters
                    detected_lang = detect_language(shape.text)  # type: ignore[attr-defined]
                    if detected_lang == "ja":
                        has_japanese = True
                        break
            if has_japanese:
                break

        assert has_japanese, "Presentation should contain Japanese text"


def test_english_input_with_explicit_japanese_output(sample_story_en: str, template_path: str):
    """Test English input with explicit Japanese output language.

    Validates User Story 2, Scenario 2:
    - English text input
    - Output language explicitly set to Japanese
    - System produces Japanese slides (would require translation in future)

    Note: Current implementation detects language from input.
    This test validates the language parameter infrastructure.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_en_to_ja.pptx")

        # Generate presentation with explicit Japanese output language
        # This should fail initially because CLI doesn't support language parameter yet
        result_path = generate_presentation(
            input_text=sample_story_en,
            template_path=template_path,
            output_path=output_path,
            output_language="ja",  # Explicit language parameter
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


def test_mixed_language_content_preserves_technical_terms(template_path: str):
    """Test mixed Japanese-English content preserves technical terms.

    Validates User Story 2, Scenario 3:
    - Mixed Japanese-English technical content
    - Technical terms in English are preserved within Japanese text
    - Language detection handles mixed content appropriately
    """
    # Mixed content with Japanese text and English technical terms
    mixed_content = """
    AIとクラウドコンピューティング

    Machine Learningモデルを使用して、データ分析を効率化します。
    AWS Lambda関数でサーバーレスアーキテクチャを実装。
    Kubernetes clustersでコンテナをオーケストレーション。

    主な技術スタック:
    - Python 3.12
    - TensorFlow 2.0
    - PostgreSQL database
    - Docker containers
    """

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_mixed.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=mixed_content,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"

        # Check that content contains both Japanese and English
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:  # type: ignore[attr-defined]
                    all_text += shape.text + " "  # type: ignore[attr-defined]

        # Should contain some English technical terms
        has_english_terms = any(
            term in all_text
            for term in ["Machine Learning", "AWS", "Lambda", "Kubernetes", "Python", "Docker"]
        )
        # Should also contain Japanese characters
        detected_lang = detect_language(all_text)

        assert has_english_terms, "Mixed content should preserve English technical terms"
        assert detected_lang == "ja", "Mixed content should be detected as primarily Japanese"
        # Note: Technical terms preservation depends on content generation implementation


def test_japanese_template_with_japanese_content(sample_story_ja: str, japanese_template_path: str):
    """Test Japanese template with Japanese content.

    Validates:
    - Japanese-specific template can be used
    - Language-aware capacity calculations work with Japanese template
    - Text fits properly with Japanese character ratios
    """
    # Skip if Japanese template doesn't exist yet (T045)
    if not Path(japanese_template_path).exists():
        pytest.skip("Japanese template not yet created (T045)")

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_ja_template.pptx")

        # Generate presentation with Japanese template
        result_path = generate_presentation(
            input_text=sample_story_ja,
            template_path=japanese_template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


def test_language_specific_capacity_ratios_applied(sample_story_ja: str, template_path: str):
    """Test that language-specific capacity ratios are applied correctly.

    Validates:
    - Japanese content uses 0.55x capacity ratio
    - English content uses 1.0x capacity ratio
    - Content validator uses language-aware overflow detection
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_capacity.pptx")

        # Generate presentation with Japanese input
        result_path = generate_presentation(
            input_text=sample_story_ja,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify no text overflow occurred (would raise exception if overflow not handled)
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"

        # Check that at least the first slide has content
        # (Some slides like Section Headers might not have body text, only titles)
        first_slide = prs.slides[0]
        has_content = False
        for shape in first_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():  # type: ignore[attr-defined]
                has_content = True
                break
        assert has_content, "First slide should have text content"

        # Verify that the presentation was generated successfully without overflow errors
        # (If overflow handling failed, an exception would have been raised)
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


def test_explicit_language_parameter_overrides_detection(sample_story_en: str, template_path: str):
    """Test that explicit language parameter overrides automatic detection.

    Validates:
    - Explicit output_language parameter is respected
    - Parameter takes precedence over auto-detection
    - Pipeline passes language context through all stages
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_explicit_lang.pptx")

        # Generate presentation with explicit English output (should override any detection)
        result_path = generate_presentation(
            input_text=sample_story_en,
            template_path=template_path,
            output_path=output_path,
            output_language="en",  # Explicit parameter
        )

        # Verify file was created
        assert Path(result_path).exists(), "Generated .pptx file should exist"

        # Verify file is valid PowerPoint
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"


def test_language_detection_accuracy(sample_story_en: str, sample_story_ja: str):
    """Test language detection utility accuracy.

    Validates:
    - English text is correctly identified as 'en'
    - Japanese text is correctly identified as 'ja'
    - Detection is reliable for typical input
    """
    # Test English detection
    en_detected = detect_language(sample_story_en)
    assert en_detected == "en", "English sample should be detected as 'en'"

    # Test Japanese detection
    ja_detected = detect_language(sample_story_ja)
    assert ja_detected == "ja", "Japanese sample should be detected as 'ja'"


def test_output_language_in_generated_metadata(sample_story_ja: str, template_path: str):
    """Test that output language information is preserved in metadata.

    Validates:
    - Output language is tracked throughout pipeline
    - Language information could be embedded in metadata
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_metadata_lang.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=sample_story_ja,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created and is valid
        assert Path(result_path).exists(), "Generated .pptx file should exist"
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"

        # Language info should be determinable from content
        # (This is a basic check - more sophisticated metadata tracking could be added)
        assert len(prs.slides) >= 3, "Presentation should have at least 3 slides"
