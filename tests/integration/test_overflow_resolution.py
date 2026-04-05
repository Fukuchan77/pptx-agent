"""Integration tests for overflow resolution.

Tests the complete end-to-end pipeline with intentionally long content that
triggers overflow detection and resolution strategies.
"""

import logging
import tempfile
from pathlib import Path

import pytest
from _pytest.logging import LogCaptureFixture
from pptx import Presentation

from pptx_agent.pipeline import generate_presentation
from pptx_agent.schemas.template_manifest import LayoutInfo, PlaceholderInfo, TemplateManifest


@pytest.fixture
def template_path() -> str:
    """Provide path to basic template for testing.

    Returns:
        Path to templates/basic-template.pptx
    """
    return "templates/basic-template.pptx"


@pytest.fixture
def overflow_minor() -> str:
    """Load minor overflow fixture (~600 chars).

    Returns:
        Content of tests/fixtures/overflow_minor.txt
    """
    fixture_path = Path("tests/fixtures/overflow_minor.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def overflow_moderate() -> str:
    """Load moderate overflow fixture (~1,200 chars).

    Returns:
        Content of tests/fixtures/overflow_moderate.txt
    """
    fixture_path = Path("tests/fixtures/overflow_moderate.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def overflow_large() -> str:
    """Load large overflow fixture (~2,800 chars).

    Returns:
        Content of tests/fixtures/overflow_large.txt
    """
    fixture_path = Path("tests/fixtures/overflow_large.txt")
    return fixture_path.read_text(encoding="utf-8")


@pytest.fixture
def overflow_extreme() -> str:
    """Load extreme overflow fixture (~5,900 chars).

    Returns:
        Content of tests/fixtures/overflow_extreme.txt
    """
    fixture_path = Path("tests/fixtures/overflow_extreme.txt")
    return fixture_path.read_text(encoding="utf-8")


def test_overflow_resolution_stage_executes_with_manifest(
    overflow_moderate: str, template_path: str, caplog: LogCaptureFixture
) -> None:
    """Test that overflow resolution stage executes when manifest is provided.

    Validates:
    - Stage 5.5 (Overflow Resolution) runs when manifest is provided
    - Stage execution time is logged
    - Pipeline completes successfully
    """
    caplog.set_level(logging.INFO)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_resolution.pptx")

        # Create a mock manifest with common layout names for testing
        title_placeholder = PlaceholderInfo(
            name="Title",
            type="TITLE",
            max_chars=100,
            language_ratio=1.0,
        )
        content_placeholder = PlaceholderInfo(
            name="Content",
            type="BODY",
            max_chars=500,
            language_ratio=1.0,
        )

        # Add common layouts that LLM might generate
        layouts = [
            LayoutInfo(name="Title Slide", placeholders=[title_placeholder]),
            LayoutInfo(
                name="Title and Content", placeholders=[title_placeholder, content_placeholder]
            ),
            LayoutInfo(name="Section Header", placeholders=[title_placeholder]),
            LayoutInfo(name="Two Content", placeholders=[title_placeholder, content_placeholder]),
            LayoutInfo(name="Content Only", placeholders=[content_placeholder]),
        ]

        manifest = TemplateManifest(
            template_name="Test Template",
            layouts=layouts,
            default_language="en",
        )

        # Generate presentation with manifest
        result_path = generate_presentation(
            input_text=overflow_moderate,
            template_path=template_path,
            output_path=output_path,
            template_manifest=manifest,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Presentation should be generated"

        # Check that overflow resolution stage executed
        overflow_stage_logs = [
            record for record in caplog.records if "Overflow Resolution completed" in record.message
        ]

        assert len(overflow_stage_logs) > 0, (
            "Overflow Resolution stage should execute when manifest is provided"
        )


def test_overflow_resolution_without_manifest(
    overflow_minor: str, template_path: str, caplog: LogCaptureFixture
) -> None:
    """Test that pipeline works when manifest is not provided.

    Validates:
    - Overflow resolution stage is skipped when manifest is None
    - Pipeline still completes successfully
    - Presentation is generated
    """
    caplog.set_level(logging.INFO)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_no_manifest.pptx")

        # Generate presentation WITHOUT manifest
        result_path = generate_presentation(
            input_text=overflow_minor,
            template_path=template_path,
            output_path=output_path,
            template_manifest=None,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Presentation should be generated"

        # Verify overflow resolution stage did NOT execute
        overflow_stage_logs = [
            record for record in caplog.records if "Overflow Resolution completed" in record.message
        ]

        assert len(overflow_stage_logs) == 0, (
            "Overflow Resolution stage should be skipped when manifest is None"
        )


def test_pipeline_generates_valid_presentation_with_large_content(
    overflow_large: str, template_path: str
) -> None:
    """Test that pipeline generates valid presentations with large content.

    Validates:
    - Large content doesn't break pipeline
    - Generated file is valid
    - Slides are accessible
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_large_overflow.pptx")

        # Generate presentation with large content
        result_path = generate_presentation(
            input_text=overflow_large,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created and is valid
        assert Path(result_path).exists(), "Presentation should be generated"

        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have minimum slides"

        # Verify all slides are accessible
        for i, slide in enumerate(prs.slides):
            assert slide.slide_layout is not None, f"Slide {i} should have layout"


def test_extreme_content_handling(overflow_extreme: str, template_path: str) -> None:
    """Test that pipeline handles extremely long content gracefully.

    Validates:
    - Extreme content doesn't crash pipeline
    - Presentation is generated
    - Content is present
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_extreme_overflow.pptx")

        # Generate presentation with extreme content
        result_path = generate_presentation(
            input_text=overflow_extreme,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Presentation should be generated"

        # Verify file is valid
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have slides"

        # Verify first slide has content
        first_slide = prs.slides[0]
        has_text = False
        for shape in first_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():  # type: ignore[attr-defined]
                has_text = True
                break

        assert has_text, "Presentation should have text content"


def test_multiple_content_lengths_in_single_presentation(template_path: str) -> None:
    """Test pipeline with content that generates multiple slides with varying lengths.

    Validates:
    - Pipeline handles slides with different content lengths
    - All slides are processed
    - Presentation is generated successfully
    """
    # Create input with multiple topics (will generate multiple slides)
    mixed_input = """
    Machine Learning Platform

    Our machine learning platform provides comprehensive capabilities for data scientists
    and engineers. It includes data preprocessing, model training, deployment, and monitoring.
    The system supports multiple frameworks and scales automatically based on workload.

    Key Features

    Advanced data pipelines with automated preprocessing and feature engineering capabilities.
    Distributed training across multiple GPUs and nodes for faster model development.
    Automated hyperparameter tuning using Bayesian optimization and grid search methods.
    Model versioning with Git integration for experiment tracking and reproducibility.
    Real-time inference API with low latency and high throughput characteristics.

    Security and Compliance

    Role-based access control with fine-grained permissions for data and model access.
    Data encryption at rest and in transit using industry-standard protocols.
    Audit logging for all system operations and model predictions for compliance.
    Integration with enterprise authentication systems including LDAP and OAuth.
    """

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = str(Path(tmpdir) / "output_mixed.pptx")

        # Generate presentation
        result_path = generate_presentation(
            input_text=mixed_input,
            template_path=template_path,
            output_path=output_path,
        )

        # Verify file was created
        assert Path(result_path).exists(), "Presentation should be generated"

        # Verify file is valid
        prs = Presentation(result_path)
        assert prs is not None, "Presentation should be readable"
        assert len(prs.slides) >= 3, "Presentation should have multiple slides"
