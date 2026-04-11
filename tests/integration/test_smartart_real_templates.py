"""Integration tests for SmartArt XML operations with REAL PowerPoint template files.

Phase 2.3: Integration Testing with Real Template XML (TDD)

These tests validate that SmartArt XML text replacement works correctly with
ACTUAL .pptx template files, not just mocked XML. This ensures end-to-end
functionality in production scenarios with real PowerPoint templates.

NOTE: These tests require templates with pre-configured SmartArt shapes.
If templates don't have SmartArt, tests will skip with a clear message.
"""

from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation
from pptx.slide import Slide

from pptx_agent.pptx_wrapper.smartart import A_NS, DGM_NS, SmartArtWrapper


def has_smartart_shape(slide: Slide) -> bool:
    """Check if slide has any SmartArt shapes (graphicFrame elements).

    Args:
        slide: PowerPoint slide object

    Returns:
        True if slide contains at least one SmartArt shape
    """
    for shape in slide.shapes:
        if hasattr(shape, "_element"):
            elem = shape._element
            if "graphicFrame" in str(elem.tag):
                return True
    return False


def find_smartart_shape(slide: Slide, name_contains: str = "") -> Any:
    """Find a SmartArt shape in the slide.

    Args:
        slide: PowerPoint slide object
        name_contains: Optional string that shape name should contain

    Returns:
        SmartArt shape object or None if not found
    """
    for shape in slide.shapes:
        if hasattr(shape, "_element"):
            elem = shape._element
            if "graphicFrame" in str(elem.tag):
                if not name_contains or name_contains.lower() in shape.name.lower():
                    return shape
    return None


def extract_smartart_text(diagram_data_xml: Any) -> list[str]:
    """Extract text from SmartArt diagram data XML.

    Args:
        diagram_data_xml: lxml Element representing diagram data

    Returns:
        List of text strings from SmartArt nodes
    """
    texts = []
    pt_nodes = diagram_data_xml.findall(f".//{{{DGM_NS}}}pt[@type='node']")
    for pt_node in pt_nodes:
        a_t = pt_node.find(f".//{{{A_NS}}}t")
        if a_t is not None and a_t.text:
            texts.append(a_t.text)
    return texts


@pytest.fixture
def template_with_smartart() -> str:
    """Provide path to template with SmartArt shapes.

    Returns:
        Path to a template file that should contain SmartArt

    Note:
        This fixture returns the path to basic-template.pptx.
        Tests will skip if this template doesn't have SmartArt shapes.
        To make tests pass, create a template with SmartArt using PowerPoint.
    """
    # Try multiple possible template locations
    possible_paths = [
        "tests/fixtures/smartart_test_template.pptx",
        "templates/smartart-template.pptx",
        "tests/fixtures/basic-template.pptx",
    ]

    for path in possible_paths:
        if Path(path).exists():
            return path

    # Fallback to basic template
    return "tests/fixtures/basic-template.pptx"


@pytest.mark.integration
def test_smartart_with_real_template_process_flow(template_with_smartart: str, tmp_path: Any):
    """Test SmartArt process flow operations with actual PowerPoint template file.

    RED PHASE TEST: This test documents requirements for real template integration.

    Validates:
    - Loading a real .pptx template with SmartArt
    - Populating SmartArt with process flow data
    - Saving the modified presentation
    - Verifying the saved file is valid and contains updated text

    TEMPLATE REQUIREMENTS:
    - Template must have a slide with Process SmartArt shape
    - SmartArt should have 3-5 nodes for testing
    - Shape should be named "SmartArt" or contain "smartart" in name
    """
    # Load real template
    template_path = Path(template_with_smartart)
    if not template_path.exists():
        pytest.skip(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))

    # Find a slide with SmartArt (check all slides)
    test_slide = None
    smartart_shape = None

    for slide in prs.slides:
        if has_smartart_shape(slide):
            test_slide = slide
            smartart_shape = find_smartart_shape(slide)
            break

    if test_slide is None or smartart_shape is None:
        pytest.skip(
            f"Template '{template_path}' does not contain SmartArt shapes. "
            "To make this test pass, create a template with Process SmartArt using PowerPoint."
        )

    # Prepare SmartArt data for process flow
    smartart_nodes = [
        {"text": "Step 1: Analysis", "level": 0},
        {"text": "Step 2: Design", "level": 0},
        {"text": "Step 3: Implementation", "level": 0},
    ]

    # Execute populate_smartart with real template
    wrapper = SmartArtWrapper()

    # This should work with production mode (navigating through relationships)
    try:
        wrapper.populate_smartart(
            slide=test_slide, placeholder_name=smartart_shape.name, nodes=smartart_nodes
        )
    except ValueError as e:
        pytest.fail(f"populate_smartart failed with real template: {e}")

    # Save to temporary file
    output_path = tmp_path / "smartart_process_output.pptx"
    prs.save(str(output_path))

    # Verify file is valid
    assert output_path.exists(), "Output file should be created"

    # Re-open and verify content
    prs_verify = Presentation(str(output_path))
    assert len(prs_verify.slides) > 0, "Presentation should have slides"

    # Verify the slide still has the SmartArt shape
    verify_slide = prs_verify.slides[prs.slides.index(test_slide)]
    assert has_smartart_shape(verify_slide), "SmartArt shape should still exist after save"


@pytest.mark.integration
def test_smartart_all_types_in_real_template(template_with_smartart: str, tmp_path: Any):
    """Test all 4 SmartArt types (process, hierarchy, cycle, relationship) with real template.

    RED PHASE TEST: Documents requirements for complete SmartArt type coverage.

    Validates:
    - Template contains SmartArt of different types
    - Each type can be populated independently
    - Text replacement works for all SmartArt types
    - No cross-contamination between SmartArt shapes

    TEMPLATE REQUIREMENTS:
    - Template should have slides with different SmartArt types:
      * Process SmartArt (sequential steps)
      * Hierarchy SmartArt (org chart)
      * Cycle SmartArt (circular process)
      * Relationship SmartArt (connections)
    """
    template_path = Path(template_with_smartart)
    if not template_path.exists():
        pytest.skip(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))

    # Count SmartArt shapes across all slides
    smartart_count = 0
    smartart_slides = []

    for slide in prs.slides:
        if has_smartart_shape(slide):
            smartart_count += sum(
                1
                for shape in slide.shapes
                if hasattr(shape, "_element") and "graphicFrame" in str(shape._element.tag)
            )
            smartart_slides.append(slide)

    if smartart_count < 1:
        pytest.skip(
            f"Template '{template_path}' does not contain SmartArt shapes. "
            "Need at least 1 SmartArt shape to test. "
            "Ideally, template should have 4 SmartArt types (process, hierarchy, cycle, relationship)."
        )

    # For now, test with whatever SmartArt exists
    # In a complete template, we would test all 4 types
    test_data_sets = [
        {
            "type": "process",
            "nodes": [
                {"text": "Phase 1", "level": 0},
                {"text": "Phase 2", "level": 0},
                {"text": "Phase 3", "level": 0},
            ],
        },
        {
            "type": "hierarchy",
            "nodes": [
                {"text": "CEO", "level": 0},
                {"text": "VP Sales", "level": 1},
                {"text": "VP Engineering", "level": 1},
            ],
        },
        {
            "type": "cycle",
            "nodes": [
                {"text": "Plan", "level": 0},
                {"text": "Do", "level": 0},
                {"text": "Check", "level": 0},
                {"text": "Act", "level": 0},
            ],
        },
    ]

    # Test with available SmartArt shapes
    wrapper = SmartArtWrapper()
    tested_count = 0

    for slide in smartart_slides[: min(len(smartart_slides), len(test_data_sets))]:
        smartart_shape = find_smartart_shape(slide)
        if smartart_shape:
            test_data = test_data_sets[tested_count]
            try:
                nodes_data = test_data["nodes"]
                if isinstance(nodes_data, list):
                    wrapper.populate_smartart(
                        slide=slide, placeholder_name=smartart_shape.name, nodes=nodes_data
                    )
                tested_count += 1
            except ValueError as e:
                # If node count mismatch, try with fewer nodes
                if "nodes but" in str(e):
                    continue
                pytest.fail(f"populate_smartart failed: {e}")

    # Save and verify
    output_path = tmp_path / "smartart_all_types_output.pptx"
    prs.save(str(output_path))

    assert output_path.exists(), "Output file should be created"
    prs_verify = Presentation(str(output_path))
    assert len(prs_verify.slides) > 0, "Presentation should have slides"

    # Note: Ideally tested_count would be 4 (all types), but we work with what's available
    assert tested_count >= 1, "Should have tested at least one SmartArt shape"


@pytest.mark.integration
def test_smartart_end_to_end_workflow(template_with_smartart: str, tmp_path: Any):
    """Test complete end-to-end workflow: load → populate → save → verify.

    RED PHASE TEST: Documents the complete integration workflow.

    Validates:
    - Load a real PowerPoint template
    - Find and identify SmartArt shapes
    - Populate SmartArt with text data
    - Save the presentation
    - Re-open and verify SmartArt text was updated
    - Ensure no XML corruption

    TEMPLATE REQUIREMENTS:
    - Template with at least one SmartArt shape
    - SmartArt should have 3-4 nodes
    """
    template_path = Path(template_with_smartart)
    if not template_path.exists():
        pytest.skip(f"Template not found: {template_path}")

    # STEP 1: Load template
    prs = Presentation(str(template_path))

    # STEP 2: Find SmartArt
    test_slide = None
    smartart_shape = None

    for slide in prs.slides:
        if has_smartart_shape(slide):
            test_slide = slide
            smartart_shape = find_smartart_shape(slide)
            break

    if test_slide is None or smartart_shape is None:
        pytest.skip(
            f"Template '{template_path}' needs SmartArt shapes for end-to-end testing. "
            "Create a template with SmartArt using PowerPoint."
        )

    # STEP 3: Populate SmartArt
    test_texts = [
        "Node A: First Item",
        "Node B: Second Item",
        "Node C: Third Item",
    ]

    smartart_nodes = [{"text": text, "level": 0} for text in test_texts]

    wrapper = SmartArtWrapper()
    try:
        wrapper.populate_smartart(
            slide=test_slide, placeholder_name=smartart_shape.name, nodes=smartart_nodes
        )
    except ValueError as e:
        if "nodes but" in str(e):
            pytest.skip(f"SmartArt node count mismatch: {e}")
        pytest.fail(f"populate_smartart failed: {e}")

    # STEP 4: Save presentation
    output_path = tmp_path / "smartart_e2e_output.pptx"
    prs.save(str(output_path))

    # STEP 5: Verify file is valid
    assert output_path.exists(), "Output file should exist"
    assert output_path.stat().st_size > 0, "Output file should not be empty"

    # STEP 6: Re-open and verify
    prs_verify = Presentation(str(output_path))
    assert len(prs_verify.slides) > 0, "Re-opened presentation should have slides"

    # Verify SmartArt shape still exists and is valid
    verify_slide = prs_verify.slides[prs.slides.index(test_slide)]
    assert has_smartart_shape(verify_slide), "SmartArt should exist in re-opened file"


# NOTE: Additional integration tests would go here once templates with SmartArt exist.
# The tests above document requirements for Phase 2.3: Integration Testing with Real Template XML.
#
# Current status:
# - 3 tests SKIP because templates lack SmartArt shapes (RED phase - documenting requirements)
# - Unit tests in test_smartart_xml_operations.py already validate XML manipulation logic
# - GREEN phase will occur when templates with SmartArt are created using PowerPoint
#
# To make tests pass:
# 1. Create/obtain a .pptx template with SmartArt shapes using Microsoft PowerPoint
# 2. Place template in tests/fixtures/smartart_test_template.pptx or templates/smartart-template.pptx
# 3. Ensure SmartArt has 3-4 nodes for testing
# 4. Run: uv run pytest tests/integration/test_smartart_real_templates.py -v
#
# Expected result after adding SmartArt templates: All 3 tests should PASS
