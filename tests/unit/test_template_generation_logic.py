"""Comprehensive unit tests for template generation logic.

Tests for generate_basic_template() and generate_japanese_template()
functions in src/pptx_agent/config.py.

Testing Strategy:
- Normal cases: successful generation with valid paths
- File validation: generated files are valid PPTX
- Content validation: generated files can be opened and have expected structure
- Error cases: invalid paths, permission errors, etc.
- Integration: both templates work with validate_templates()
"""

from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from pptx_agent.config import (
    generate_basic_template,
    generate_japanese_template,
    validate_templates,
)


class TestGenerateBasicTemplate:
    """Test suite for generate_basic_template() function."""

    def test_generates_file_at_valid_path(self, tmp_path: Path) -> None:
        """Should successfully generate template at valid path."""
        output_path = tmp_path / "test-basic.pptx"

        generate_basic_template(output_path)

        assert output_path.exists(), "Template file should be created"
        assert output_path.is_file(), "Output should be a file"
        assert output_path.stat().st_size > 0, "Generated file should not be empty"

    def test_generated_file_is_valid_pptx(self, tmp_path: Path) -> None:
        """Generated file should be a valid PPTX file."""
        output_path = tmp_path / "test-basic.pptx"

        generate_basic_template(output_path)

        # Should be able to open with python-pptx without errors
        try:
            prs = Presentation(str(output_path))
            assert prs is not None, "Should successfully create Presentation object"
        except Exception as e:
            pytest.fail(f"Generated file is not a valid PPTX: {e}")

    def test_generated_file_has_layouts(self, tmp_path: Path) -> None:
        """Generated file should contain at least one slide layout."""
        output_path = tmp_path / "test-basic.pptx"

        generate_basic_template(output_path)

        prs = Presentation(str(output_path))
        assert len(prs.slide_layouts) > 0, "Template should have at least one layout"

    def test_generated_file_has_master(self, tmp_path: Path) -> None:
        """Generated file should have a slide master."""
        output_path = tmp_path / "test-basic.pptx"

        generate_basic_template(output_path)

        prs = Presentation(str(output_path))
        assert len(prs.slide_masters) > 0, "Template should have at least one slide master"

    def test_creates_parent_directories(self, tmp_path: Path) -> None:
        """Should create parent directories if they don't exist."""
        output_path = tmp_path / "nested" / "dirs" / "test-basic.pptx"
        assert not output_path.parent.exists(), "Parent directory should not exist initially"

        generate_basic_template(output_path)

        assert output_path.exists(), "Template file should be created"
        assert output_path.parent.exists(), "Parent directories should be created"

    def test_overwrites_existing_file(self, tmp_path: Path) -> None:
        """Should overwrite existing file without error."""
        output_path = tmp_path / "test-basic.pptx"
        # Create a dummy file first
        output_path.write_bytes(b"dummy content")
        original_size = output_path.stat().st_size

        generate_basic_template(output_path)

        assert output_path.exists(), "File should still exist"
        new_size = output_path.stat().st_size
        assert new_size != original_size, "File should be overwritten with different content"

    def test_raises_runtime_error_on_invalid_path(self, tmp_path: Path) -> None:
        """Should raise RuntimeError when path is invalid."""
        # Use a path that cannot be written to (e.g., a file as parent)
        dummy_file = tmp_path / "dummy.txt"
        dummy_file.write_text("test")
        invalid_path = dummy_file / "test-basic.pptx"  # File as parent directory

        with pytest.raises(RuntimeError) as exc_info:
            generate_basic_template(invalid_path)

        error_msg = str(exc_info.value)
        assert "Failed to generate basic template" in error_msg
        assert str(invalid_path) in error_msg

    def test_raises_runtime_error_on_permission_denied(self, tmp_path: Path) -> None:
        """Should raise RuntimeError when write permission is denied."""
        if not hasattr(tmp_path, "chmod"):
            pytest.skip("chmod not available on this platform")

        # Create a read-only directory
        readonly_dir = tmp_path / "readonly"
        readonly_dir.mkdir()
        output_path = readonly_dir / "test-basic.pptx"

        try:
            # Make directory read-only
            readonly_dir.chmod(0o444)

            with pytest.raises(RuntimeError) as exc_info:
                generate_basic_template(output_path)

            error_msg = str(exc_info.value)
            assert "Failed to generate basic template" in error_msg
        finally:
            # Restore permissions for cleanup
            readonly_dir.chmod(0o755)

    def test_logs_success_message(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        """Should log success message after generation."""
        import logging

        output_path = tmp_path / "test-basic.pptx"

        # Ensure we capture INFO level logs
        caplog.set_level(logging.INFO, logger="pptx_agent.templates")

        generate_basic_template(output_path)

        # Check if success message was logged
        assert any("Generated basic template" in record.message for record in caplog.records), (
            "Should log success message"
        )

    def test_logs_error_on_failure(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        """Should log error message on generation failure."""
        dummy_file = tmp_path / "dummy.txt"
        dummy_file.write_text("test")
        invalid_path = dummy_file / "test-basic.pptx"

        with pytest.raises(RuntimeError):
            generate_basic_template(invalid_path)

        assert any(
            "Failed to generate basic template" in record.message for record in caplog.records
        ), "Should log error message on failure"

    def test_wraps_exception_in_runtime_error(self, tmp_path: Path) -> None:
        """Should wrap underlying exceptions in RuntimeError."""
        output_path = tmp_path / "test-basic.pptx"

        # Mock Presentation to raise an exception
        with patch("pptx_agent.templates.Presentation") as mock_prs:
            mock_prs.side_effect = OSError("Disk full")

            with pytest.raises(RuntimeError) as exc_info:
                generate_basic_template(output_path)

            # Should have the original exception as __cause__
            assert exc_info.value.__cause__ is not None
            assert isinstance(exc_info.value.__cause__, OSError)


class TestGenerateJapaneseTemplate:
    """Test suite for generate_japanese_template() function."""

    def test_generates_file_at_valid_path(self, tmp_path: Path) -> None:
        """Should successfully generate template at valid path."""
        output_path = tmp_path / "test-japanese.pptx"

        generate_japanese_template(output_path)

        assert output_path.exists(), "Template file should be created"
        assert output_path.is_file(), "Output should be a file"
        assert output_path.stat().st_size > 0, "Generated file should not be empty"

    def test_generated_file_is_valid_pptx(self, tmp_path: Path) -> None:
        """Generated file should be a valid PPTX file."""
        output_path = tmp_path / "test-japanese.pptx"

        generate_japanese_template(output_path)

        # Should be able to open with python-pptx without errors
        try:
            prs = Presentation(str(output_path))
            assert prs is not None, "Should successfully create Presentation object"
        except Exception as e:
            pytest.fail(f"Generated file is not a valid PPTX: {e}")

    def test_generated_file_has_layouts(self, tmp_path: Path) -> None:
        """Generated file should contain at least one slide layout."""
        output_path = tmp_path / "test-japanese.pptx"

        generate_japanese_template(output_path)

        prs = Presentation(str(output_path))
        assert len(prs.slide_layouts) > 0, "Template should have at least one layout"

    def test_creates_parent_directories(self, tmp_path: Path) -> None:
        """Should create parent directories if they don't exist."""
        output_path = tmp_path / "nested" / "dirs" / "test-japanese.pptx"

        generate_japanese_template(output_path)

        assert output_path.exists(), "Template file should be created"
        assert output_path.parent.exists(), "Parent directories should be created"

    def test_overwrites_existing_file(self, tmp_path: Path) -> None:
        """Should overwrite existing file without error."""
        output_path = tmp_path / "test-japanese.pptx"
        output_path.write_bytes(b"dummy content")

        generate_japanese_template(output_path)

        assert output_path.exists(), "File should still exist"
        # Verify it's a valid PPTX now
        prs = Presentation(str(output_path))
        assert prs is not None

    def test_raises_runtime_error_on_invalid_path(self, tmp_path: Path) -> None:
        """Should raise RuntimeError when path is invalid."""
        dummy_file = tmp_path / "dummy.txt"
        dummy_file.write_text("test")
        invalid_path = dummy_file / "test-japanese.pptx"

        with pytest.raises(RuntimeError) as exc_info:
            generate_japanese_template(invalid_path)

        error_msg = str(exc_info.value)
        assert "Failed to generate japanese template" in error_msg
        assert str(invalid_path) in error_msg

    def test_logs_success_message(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        """Should log success message after generation."""
        import logging

        output_path = tmp_path / "test-japanese.pptx"

        # Ensure we capture INFO level logs
        caplog.set_level(logging.INFO, logger="pptx_agent.templates")

        generate_japanese_template(output_path)

        # Check if success message was logged
        assert any("Generated japanese template" in record.message for record in caplog.records), (
            "Should log success message"
        )

    def test_wraps_exception_in_runtime_error(self, tmp_path: Path) -> None:
        """Should wrap underlying exceptions in RuntimeError."""
        output_path = tmp_path / "test-japanese.pptx"

        with patch("pptx_agent.templates.Presentation") as mock_prs:
            mock_prs.side_effect = OSError("Disk full")

            with pytest.raises(RuntimeError) as exc_info:
                generate_japanese_template(output_path)

            assert exc_info.value.__cause__ is not None


class TestTemplateGenerationIntegration:
    """Integration tests for template generation with validate_templates()."""

    def test_can_generate_both_templates_simultaneously(self, tmp_path: Path) -> None:
        """Should be able to generate both templates without conflict."""
        basic_path = tmp_path / "basic-template.pptx"
        japanese_path = tmp_path / "japanese-template.pptx"

        generate_basic_template(basic_path)
        generate_japanese_template(japanese_path)

        assert basic_path.exists(), "Basic template should be created"
        assert japanese_path.exists(), "Japanese template should be created"

        # Both should be valid PPTX files
        Presentation(str(basic_path))
        Presentation(str(japanese_path))

    def test_generated_templates_pass_validation(self, tmp_path: Path) -> None:
        """Generated templates should pass validate_templates() check."""
        template_dir = tmp_path / "templates"
        template_dir.mkdir()

        basic_path = template_dir / "basic-template.pptx"
        japanese_path = template_dir / "japanese-template.pptx"

        generate_basic_template(basic_path)
        generate_japanese_template(japanese_path)

        # Should not raise FileNotFoundError
        validate_templates(templates_dir=template_dir, auto_generate=False)

    def test_generated_templates_can_create_presentations(self, tmp_path: Path) -> None:
        """Should be able to create presentations from generated templates."""
        template_path = tmp_path / "test-template.pptx"

        generate_basic_template(template_path)

        # Load template and add a slide
        prs = Presentation(str(template_path))
        layout = prs.slide_layouts[0]
        _ = prs.slides.add_slide(layout)

        # Should be able to save modified presentation
        output_path = tmp_path / "output.pptx"
        prs.save(str(output_path))

        assert output_path.exists(), "Modified presentation should be saved"

    def test_generated_template_has_same_structure_as_blank(self, tmp_path: Path) -> None:
        """Generated template should have same structure as blank Presentation()."""
        generated_path = tmp_path / "generated.pptx"
        reference_path = tmp_path / "reference.pptx"

        # Generate our template
        generate_basic_template(generated_path)

        # Create reference blank presentation
        blank_prs = Presentation()
        blank_prs.save(str(reference_path))

        # Load both and compare basic structure
        generated_prs = Presentation(str(generated_path))
        reference_prs = Presentation(str(reference_path))

        # Both should have same number of slide masters and layouts
        assert len(generated_prs.slide_masters) == len(reference_prs.slide_masters), (
            "Should have same number of slide masters"
        )
        assert len(generated_prs.slide_layouts) == len(reference_prs.slide_layouts), (
            "Should have same number of slide layouts"
        )


class TestPerformanceCharacteristics:
    """Test performance characteristics of template generation."""

    def test_generation_completes_within_reasonable_time(self, tmp_path: Path) -> None:
        """Template generation should complete within 2 seconds."""
        import time

        output_path = tmp_path / "test-perf.pptx"

        start_time = time.time()
        generate_basic_template(output_path)
        elapsed_time = time.time() - start_time

        assert elapsed_time < 2.0, f"Template generation took {elapsed_time:.2f}s, should be < 2.0s"

    def test_multiple_generations_do_not_leak_memory(self, tmp_path: Path) -> None:
        """Multiple template generations should not leak memory significantly."""
        # Generate multiple templates to detect obvious memory leaks
        for i in range(10):
            output_path = tmp_path / f"template_{i}.pptx"
            generate_basic_template(output_path)
            assert output_path.exists(), f"Template {i} should be created"

        # If we got here without MemoryError, we're good
        # (This is a basic smoke test, not a comprehensive memory test)


class TestEdgeCases:
    """Test edge cases and boundary conditions."""

    def test_handles_unicode_in_path(self, tmp_path: Path) -> None:
        """Should handle Unicode characters in file path."""
        output_path = tmp_path / "テンプレート_日本語.pptx"

        generate_basic_template(output_path)

        assert output_path.exists(), "Should handle Unicode in path"
        # Should be readable
        prs = Presentation(str(output_path))
        assert prs is not None

    def test_handles_spaces_in_path(self, tmp_path: Path) -> None:
        """Should handle spaces in file path."""
        output_path = tmp_path / "template with spaces.pptx"

        generate_basic_template(output_path)

        assert output_path.exists(), "Should handle spaces in path"

    def test_handles_long_path_names(self, tmp_path: Path) -> None:
        """Should handle reasonably long path names."""
        # Create a long but valid path
        long_dir = tmp_path / ("a" * 50) / ("b" * 50)
        output_path = long_dir / "template.pptx"

        generate_basic_template(output_path)

        assert output_path.exists(), "Should handle long path names"

    def test_handles_dot_in_directory_name(self, tmp_path: Path) -> None:
        """Should handle dots in directory names."""
        output_path = tmp_path / "my.templates" / "test.pptx"

        generate_basic_template(output_path)

        assert output_path.exists(), "Should handle dots in directory names"
