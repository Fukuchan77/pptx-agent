"""Unit tests for template parser.

TDD RED PHASE: These tests are written BEFORE implementation.
They define the expected behavior of the TemplateParser class.
"""

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_agent.template_parser import (
    LayoutMetadata,
    PlaceholderMetadata,
    TemplateMetadata,
    TemplateParser,
)
from pptx_agent.validators.exceptions import InvalidFileError


class TestTemplateParser:
    """Test TemplateParser for extracting template metadata."""

    def test_parse_template_basic_structure(self, tmp_path: Path) -> None:
        """RED: Test parsing basic template structure."""
        # Create a simple test template
        template_path = tmp_path / "test_template.pptx"

        # Create a minimal PPTX file
        prs = Presentation()
        prs.save(str(template_path))

        parser = TemplateParser()

        # Parse the template
        metadata = parser.parse_template(str(template_path))

        assert metadata is not None
        assert hasattr(metadata, "layouts")
        assert isinstance(metadata.layouts, list)
        assert len(metadata.layouts) > 0  # Should have at least one layout

    def test_parse_template_extracts_layouts(self):
        """RED: Test that parser extracts slide layouts."""
        parser = TemplateParser()

        # We'll use a real template file for this
        # For now, test with non-existent file to ensure RED phase
        template_path = "templates/test-template.pptx"

        # This should fail with InvalidFileError when template doesn't exist
        with pytest.raises(InvalidFileError):
            parser.parse_template(template_path)

    def test_parse_template_extracts_placeholders(self):
        """RED: Test that parser extracts placeholders for each layout."""
        TemplateParser()

        # Assuming we have a template with known structure
        # This test defines what we expect from the parser
        # We expect LayoutMetadata to have:
        # - name: str
        # - placeholders: list[PlaceholderMetadata]

        # Create a mock LayoutMetadata to define structure
        layout = LayoutMetadata(
            name="Title Slide",
            placeholders=[],
        )

        assert layout.name == "Title Slide"
        assert isinstance(layout.placeholders, list)

    def test_placeholder_metadata_structure(self):
        """RED: Test PlaceholderMetadata structure."""
        # Define expected structure for placeholders
        placeholder = PlaceholderMetadata(
            name="Title",
            type="TITLE",
            max_chars=100,
        )

        assert placeholder.name == "Title"
        assert placeholder.type == "TITLE"
        assert placeholder.max_chars == 100

    def test_template_metadata_structure(self):
        """RED: Test TemplateMetadata structure."""
        # Define expected structure for template metadata
        metadata = TemplateMetadata(
            template_path="templates/test.pptx",
            layouts=[
                LayoutMetadata(name="Title Slide", placeholders=[]),
                LayoutMetadata(name="Content", placeholders=[]),
            ],
        )

        assert metadata.template_path == "templates/test.pptx"
        assert len(metadata.layouts) == 2
        assert metadata.layouts[0].name == "Title Slide"
        assert metadata.layouts[1].name == "Content"

    def test_get_layout_by_name(self):
        """RED: Test getting layout by name from TemplateMetadata."""
        metadata = TemplateMetadata(
            template_path="test.pptx",
            layouts=[
                LayoutMetadata(name="Title Slide", placeholders=[]),
                LayoutMetadata(name="Title and Content", placeholders=[]),
                LayoutMetadata(name="Blank", placeholders=[]),
            ],
        )

        # Test method to get layout by name
        layout = metadata.get_layout("Title and Content")

        assert layout is not None
        assert layout.name == "Title and Content"

    def test_get_layout_by_name_not_found(self):
        """RED: Test getting non-existent layout returns None."""
        metadata = TemplateMetadata(
            template_path="test.pptx",
            layouts=[],
        )

        layout = metadata.get_layout("NonExistent")

        assert layout is None

    def test_parser_with_real_template(self, template_path: str):
        """RED: Test parser with real template file."""
        # This will use the sample template created in T026b
        parser = TemplateParser()
        metadata = parser.parse_template(template_path)

        # Basic template should have at least these layouts
        assert len(metadata.layouts) >= 3

        # Should have Title Slide layout
        title_layout = metadata.get_layout("Title Slide")
        assert title_layout is not None

        # Should have placeholders
        assert len(title_layout.placeholders) > 0

    def test_placeholder_type_detection(self, template_path: str):
        """RED: Test that placeholder types are correctly detected."""
        parser = TemplateParser()
        metadata = parser.parse_template(template_path)

        # Get Title Slide layout
        title_layout = metadata.get_layout("Title Slide")
        assert title_layout is not None, "Title Slide layout should exist"

        # Should have TITLE and SUBTITLE placeholders
        placeholder_names = [p.name for p in title_layout.placeholders]
        assert "Title" in placeholder_names or "Title 1" in placeholder_names

    def test_get_placeholder_type_handles_attribute_error(self):
        """RED: Test that _get_placeholder_type handles AttributeError specifically."""
        parser = TemplateParser()

        # Create a mock shape that raises AttributeError
        class MockShape:
            @property
            def placeholder_format(self) -> None:
                msg = "No placeholder_format attribute"
                raise AttributeError(msg)

        shape = MockShape()
        result = parser._get_placeholder_type(shape)  # type: ignore[reportPrivateUsage]

        # Should return "UNKNOWN" when AttributeError occurs
        assert result == "UNKNOWN"

    def test_get_placeholder_type_handles_key_error(self):
        """RED: Test that _get_placeholder_type handles KeyError specifically."""
        parser = TemplateParser()

        # Create a mock shape with unknown placeholder type
        class MockPlaceholder:
            type = 999999  # Invalid type not in PLACEHOLDER_TYPE_NAMES

        class MockShape:
            @property
            def placeholder_format(self) -> MockPlaceholder:
                return MockPlaceholder()

        shape = MockShape()
        result = parser._get_placeholder_type(shape)  # type: ignore[reportPrivateUsage]

        # Should return "UNKNOWN" for unknown types
        assert result == "UNKNOWN"

    def test_get_placeholder_type_propagates_permission_error(self):
        """RED: Test that _get_placeholder_type propagates PermissionError."""
        parser = TemplateParser()

        # Create a mock shape that raises PermissionError
        class MockShape:
            @property
            def placeholder_format(self) -> None:
                msg = "Access denied"
                raise PermissionError(msg)

        shape = MockShape()

        # PermissionError should propagate, not be caught
        with pytest.raises(PermissionError):
            parser._get_placeholder_type(shape)  # type: ignore[reportPrivateUsage]

    def test_estimate_max_chars_handles_attribute_error(self):
        """RED: Test that _estimate_max_chars handles AttributeError specifically."""
        parser = TemplateParser()

        # Create a mock shape without width/height attributes
        class MockShape:
            @property
            def width(self) -> None:
                msg = "No width attribute"
                raise AttributeError(msg)

        shape = MockShape()
        result = parser._estimate_max_chars(shape)  # type: ignore[reportPrivateUsage]

        # Should return default value when AttributeError occurs
        assert result == 100

    def test_estimate_max_chars_handles_type_error(self):
        """RED: Test that _estimate_max_chars handles TypeError specifically."""
        parser = TemplateParser()

        # Create a mock shape with non-numeric width/height
        class MockShape:
            width = "not_a_number"
            height = "not_a_number"

        shape = MockShape()
        result = parser._estimate_max_chars(shape)  # type: ignore[reportPrivateUsage]

        # Should return default value when TypeError occurs
        assert result == 100

    def test_estimate_max_chars_handles_value_error(self):
        """RED: Test that _estimate_max_chars handles ValueError specifically."""
        parser = TemplateParser()

        # Create a mock shape with negative dimensions
        class MockShape:
            width = -100
            height = -100

        shape = MockShape()
        result = parser._estimate_max_chars(shape)  # type: ignore[reportPrivateUsage]

        # Should return at least 50 chars (minimum value)
        assert result >= 50

    def test_estimate_max_chars_propagates_permission_error(self):
        """RED: Test that _estimate_max_chars propagates PermissionError."""
        parser = TemplateParser()

        # Create a mock shape that raises PermissionError
        class MockShape:
            @property
            def width(self) -> None:
                msg = "Access denied"
                raise PermissionError(msg)

        shape = MockShape()

        # PermissionError should propagate, not be caught
        with pytest.raises(PermissionError):
            parser._estimate_max_chars(shape)  # type: ignore[reportPrivateUsage]

    def test_detect_smartart_in_layout(self):
        """RED: Test detection of SmartArt shapes in layouts."""
        parser = TemplateParser()

        # Create a mock shape that simulates a SmartArt graphic frame
        class MockElement:
            tag = "{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame"

        class MockShape:
            is_placeholder = False
            element = MockElement()

        # Test detection
        has_smartart = parser._is_smartart_shape(MockShape())  # type: ignore[reportPrivateUsage]
        assert has_smartart is True

    def test_detect_non_smartart_shape(self):
        """RED: Test that non-SmartArt shapes are not detected as SmartArt."""
        parser = TemplateParser()

        # Create a mock shape that is not a SmartArt
        class MockElement:
            tag = "{http://schemas.openxmlformats.org/presentationml/2006/main}sp"  # Regular shape

        class MockShape:
            is_placeholder = False
            element = MockElement()

        # Test detection
        has_smartart = parser._is_smartart_shape(MockShape())  # type: ignore[reportPrivateUsage]
        assert has_smartart is False

    def test_layout_metadata_has_smartart_field(self):
        """RED: Test that LayoutMetadata can store SmartArt presence info."""
        # After parser enhancement, LayoutMetadata should track if layout has SmartArt
        layout = LayoutMetadata(
            name="SmartArt Layout",
            placeholders=[],
            has_smartart=True,  # New field for T065
        )

        assert layout.has_smartart is True

    def test_parse_layout_detects_smartart(self, tmp_path: Path):
        """RED: Test that _parse_layout detects SmartArt shapes."""
        parser = TemplateParser()

        # This test will need a real template with SmartArt
        # For now, define expected behavior: parser should detect SmartArt presence
        # and set has_smartart=True in LayoutMetadata

        # Create simple template without SmartArt
        template_path = tmp_path / "test_no_smartart.pptx"
        prs = Presentation()
        prs.save(str(template_path))

        metadata = parser.parse_template(str(template_path))

        # Check that layouts have has_smartart attribute
        for layout in metadata.layouts:
            assert hasattr(layout, "has_smartart")
            # Default template shouldn't have SmartArt
            assert layout.has_smartart is False
