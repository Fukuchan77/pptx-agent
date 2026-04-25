"""Unit tests for QA layout check rules."""

from unittest.mock import Mock

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.util import Inches

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.rules.layout_checks import (
    EmptyPlaceholderRule,
    OverlappingObjectsRule,
    TextOverflowRule,
)
from pptx_agent.qa.schemas import Severity


@pytest.fixture
def mock_presentation() -> Mock:
    """Create mock presentation for testing.

    Returns:
        Mock PresentationWrapper instance
    """
    return Mock(spec=PresentationWrapper)


@pytest.fixture
def mock_slide_with_overflow() -> Mock:
    """Create mock slide with text overflow.

    Returns:
        Mock slide with overflowing text frame
    """
    slide = Mock()

    # Create shape with text frame that has overflow
    shape = Mock()
    shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
    shape.has_text_frame = True

    # Mock text frame with overflow
    text_frame = Mock()
    text_frame.text = "A" * 1000  # Long text

    # Mock text_frame.fit_text() to indicate overflow
    # In python-pptx, when text doesn't fit, fit_text() returns False
    # or we can check via auto_size property
    text_frame.auto_size = None  # No auto-sizing

    # Create a mock that simulates overflow by having text exceed bounds
    shape.text_frame = text_frame
    shape.width = Inches(3)
    shape.height = Inches(1)

    # Add method to detect overflow (simplified)
    def has_overflow() -> bool:
        # Simulate overflow detection: if text is very long, it overflows
        return len(text_frame.text) > 500

    shape.has_overflow = has_overflow

    slide.shapes = [shape]
    return slide


@pytest.fixture
def mock_slide_with_empty_title() -> Mock:
    """Create mock slide with empty title placeholder.

    Returns:
        Mock slide with empty title
    """
    slide = Mock()

    # Create title shape that is empty
    title_shape = Mock()
    title_shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
    title_shape.placeholder_format = Mock()
    title_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.TITLE
    title_shape.has_text_frame = True
    title_shape.text_frame = Mock()
    title_shape.text_frame.text = ""  # Empty title

    slide.shapes = Mock()
    slide.shapes.title = title_shape
    slide.shapes.__iter__ = Mock(return_value=iter([title_shape]))

    return slide


@pytest.fixture
def mock_slide_with_overlapping_shapes() -> Mock:
    """Create mock slide with overlapping shapes.

    Returns:
        Mock slide with two overlapping shapes
    """
    slide = Mock()

    # Create two shapes that overlap
    shape1 = Mock()
    shape1.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
    shape1.left = Inches(1)
    shape1.top = Inches(1)
    shape1.width = Inches(3)
    shape1.height = Inches(2)

    shape2 = Mock()
    shape2.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
    shape2.left = Inches(2)  # Overlaps with shape1
    shape2.top = Inches(1.5)  # Overlaps with shape1
    shape2.width = Inches(3)
    shape2.height = Inches(2)

    slide.shapes = [shape1, shape2]
    return slide


class TestTextOverflowRule:
    """Tests for text overflow detection rule (QA-L-001)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = TextOverflowRule()

        assert rule.rule_id == "QA-L-001"
        assert "overflow" in rule.description.lower()
        assert rule.auto_fixable is True

    def test_detects_text_overflow(
        self, mock_presentation: Mock, mock_slide_with_overflow: Mock
    ) -> None:
        """Test detection of text overflow in text frames."""
        # Setup presentation with one slide containing overflow
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_overflow]
        mock_presentation.prs = mock_prs

        rule = TextOverflowRule()
        issues = rule.validate(mock_presentation)

        # Should detect overflow
        assert len(issues) > 0
        assert issues[0].rule_id == "QA-L-001"
        assert issues[0].severity == Severity.ERROR
        assert issues[0].auto_fixable is True
        assert "overflow" in issues[0].message.lower()

    def test_no_issues_when_text_fits(self, mock_presentation: Mock) -> None:
        """Test no issues when text fits within bounds."""
        # Create slide with text that fits
        slide = Mock()
        shape = Mock()
        shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
        shape.has_text_frame = True
        shape.text_frame = Mock()
        shape.text_frame.text = "Short text"
        shape.width = Inches(5)
        shape.height = Inches(3)
        shape.has_overflow = lambda: False

        slide.shapes = [shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = TextOverflowRule()
        issues = rule.validate(mock_presentation)

        # Should not detect any issues
        assert len(issues) == 0

    def test_handles_slides_without_text_frames(self, mock_presentation: Mock) -> None:
        """Test graceful handling of slides without text frames."""
        # Create slide with shape that has no text frame
        slide = Mock()
        shape = Mock()
        shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        shape.has_text_frame = False

        slide.shapes = [shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = TextOverflowRule()
        issues = rule.validate(mock_presentation)

        # Should handle gracefully without errors
        assert len(issues) == 0


class TestEmptyPlaceholderRule:
    """Tests for empty placeholder detection rule (QA-L-002)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = EmptyPlaceholderRule()

        assert rule.rule_id == "QA-L-002"
        assert "empty" in rule.description.lower() or "title" in rule.description.lower()
        assert rule.auto_fixable is True

    def test_detects_empty_title(
        self, mock_presentation: Mock, mock_slide_with_empty_title: Mock
    ) -> None:
        """Test detection of empty title placeholders."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_empty_title]
        mock_presentation.prs = mock_prs

        rule = EmptyPlaceholderRule()
        issues = rule.validate(mock_presentation)

        # Should detect empty title
        assert len(issues) > 0
        assert issues[0].rule_id == "QA-L-002"
        assert issues[0].severity == Severity.ERROR
        assert issues[0].auto_fixable is True
        assert "title" in issues[0].message.lower() or "empty" in issues[0].message.lower()

    def test_no_issues_when_title_populated(self, mock_presentation: Mock) -> None:
        """Test no issues when title is populated."""
        slide = Mock()

        title_shape = Mock()
        title_shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
        title_shape.placeholder_format = Mock()
        title_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.TITLE
        title_shape.has_text_frame = True
        title_shape.text_frame = Mock()
        title_shape.text_frame.text = "Valid Title"

        slide.shapes = Mock()
        slide.shapes.title = title_shape
        slide.shapes.__iter__ = Mock(return_value=iter([title_shape]))

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = EmptyPlaceholderRule()
        issues = rule.validate(mock_presentation)

        # Should not detect any issues
        assert len(issues) == 0

    def test_handles_slides_without_title(self, mock_presentation: Mock) -> None:
        """Test graceful handling of slides without title placeholder."""
        slide = Mock()
        slide.shapes = Mock()
        slide.shapes.title = None
        slide.shapes.__iter__ = Mock(return_value=iter([]))

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = EmptyPlaceholderRule()
        issues = rule.validate(mock_presentation)

        # Should handle gracefully - no title means no empty title issue
        assert len(issues) == 0


class TestOverlappingObjectsRule:
    """Tests for overlapping objects detection rule (QA-L-004)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = OverlappingObjectsRule()

        assert rule.rule_id == "QA-L-004"
        assert "overlap" in rule.description.lower()
        assert rule.auto_fixable is False  # Overlapping is report-only

    def test_detects_overlapping_shapes(
        self, mock_presentation: Mock, mock_slide_with_overlapping_shapes: Mock
    ) -> None:
        """Test detection of overlapping shapes."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_overlapping_shapes]
        mock_presentation.prs = mock_prs

        rule = OverlappingObjectsRule()
        issues = rule.validate(mock_presentation)

        # Should detect overlapping shapes
        assert len(issues) > 0
        assert issues[0].rule_id == "QA-L-004"
        assert issues[0].severity == Severity.WARNING
        assert issues[0].auto_fixable is False
        assert "overlap" in issues[0].message.lower()

    def test_no_issues_when_shapes_dont_overlap(self, mock_presentation: Mock) -> None:
        """Test no issues when shapes don't overlap."""
        slide = Mock()

        # Create two shapes that don't overlap
        shape1 = Mock()
        shape1.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        shape1.left = Inches(1)
        shape1.top = Inches(1)
        shape1.width = Inches(2)
        shape1.height = Inches(1)

        shape2 = Mock()
        shape2.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        shape2.left = Inches(4)  # Far from shape1
        shape2.top = Inches(1)
        shape2.width = Inches(2)
        shape2.height = Inches(1)

        slide.shapes = [shape1, shape2]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = OverlappingObjectsRule()
        issues = rule.validate(mock_presentation)

        # Should not detect any issues
        assert len(issues) == 0

    def test_handles_single_shape_slide(self, mock_presentation: Mock) -> None:
        """Test graceful handling of slides with single shape."""
        slide = Mock()

        shape = Mock()
        shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        shape.left = Inches(1)
        shape.top = Inches(1)
        shape.width = Inches(2)
        shape.height = Inches(1)

        slide.shapes = [shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = OverlappingObjectsRule()
        issues = rule.validate(mock_presentation)

        # Should handle gracefully - can't overlap with only one shape
        assert len(issues) == 0


# Made with Bob
