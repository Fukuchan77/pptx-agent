"""Unit tests for QA content check rules."""

from unittest.mock import Mock

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.rules.content_checks import (
    BulletLengthRule,
    DuplicateTitleRule,
    MissingChartDataRule,
    PathologicalTableDimensionRule,
    SpeakerNotesVerificationRule,
    UnpopulatedImagePlaceholderRule,
)
from pptx_agent.qa.schemas import Severity


@pytest.fixture
def mock_presentation() -> Mock:
    """Create mock presentation for testing.

    Returns:
        Mock PresentationWrapper instance
    """
    pres = Mock(spec=PresentationWrapper)
    pres.is_loaded = True
    return pres


@pytest.fixture
def mock_slide_with_long_bullets() -> Mock:
    """Create mock slide with overly long bullet points.

    Returns:
        Mock slide with long bullets
    """
    slide = Mock()

    # Create shape with text frame containing long bullets
    shape = Mock()
    shape.has_text_frame = True

    text_frame = Mock()

    # Create paragraphs with long text
    para1 = Mock()
    para1.level = 0
    para1.text = "A" * 150  # Exceeds default 120 char limit

    para2 = Mock()
    para2.level = 1
    para2.text = "B" * 130  # Also exceeds limit

    text_frame.paragraphs = [para1, para2]
    shape.text_frame = text_frame

    slide.shapes = [shape]
    return slide


@pytest.fixture
def mock_slides_with_duplicate_titles() -> list[Mock]:
    """Create mock slides with duplicate titles.

    Returns:
        List of mock slides with duplicate titles
    """
    slides = []

    for i in range(3):
        slide = Mock()
        title_shape = Mock()
        title_shape.is_placeholder = True
        title_shape.placeholder_format = Mock()
        title_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.TITLE
        title_shape.has_text_frame = True
        title_shape.text_frame = Mock()

        # First two slides have same title
        if i < 2:
            title_shape.text_frame.text = "Duplicate Title"
        else:
            title_shape.text_frame.text = "Unique Title"

        slide.shapes = [title_shape]
        slides.append(slide)

    return slides


@pytest.fixture
def mock_slide_with_empty_image_placeholder() -> Mock:
    """Create mock slide with unpopulated image placeholder.

    Returns:
        Mock slide with empty image placeholder
    """
    slide = Mock()

    # Create image placeholder that is not populated
    image_shape = Mock()
    image_shape.is_placeholder = True
    image_shape.placeholder_format = Mock()
    image_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.PICTURE
    image_shape.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER  # Still a placeholder, not filled

    slide.shapes = [image_shape]
    return slide


@pytest.fixture
def mock_slide_with_pathological_table() -> Mock:
    """Create mock slide with 1x1 table.

    Returns:
        Mock slide with pathological table
    """
    slide = Mock()

    # Create 1x1 table
    table_shape = Mock()
    table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

    table = Mock()

    # Create 1x1 table structure
    row = Mock()
    cell = Mock()
    cell.text = "Single cell"
    row.cells = [cell]

    table.rows = [row]
    table.columns = [Mock()]  # 1 column

    table_shape.table = table

    slide.shapes = [table_shape]
    return slide


@pytest.fixture
def mock_slide_with_large_table() -> Mock:
    """Create mock slide with very large table.

    Returns:
        Mock slide with large table (12x10)
    """
    slide = Mock()

    # Create large table
    table_shape = Mock()
    table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

    table = Mock()

    # Create 12 rows
    rows = []
    for _ in range(12):
        row = Mock()
        cells = []
        for _ in range(10):
            cell = Mock()
            cell.text = "Data"
            cells.append(cell)
        row.cells = cells
        rows.append(row)

    table.rows = rows
    table.columns = [Mock() for _ in range(10)]  # 10 columns

    table_shape.table = table

    slide.shapes = [table_shape]
    return slide


@pytest.fixture
def mock_slide_with_empty_chart() -> Mock:
    """Create mock slide with chart that has no data.

    Returns:
        Mock slide with empty chart
    """
    slide = Mock()

    # Create chart with no data
    chart_shape = Mock()
    chart_shape.shape_type = MSO_SHAPE_TYPE.CHART

    chart = Mock()
    chart.plots = []  # No plots
    chart.series = []

    chart_shape.chart = chart

    slide.shapes = [chart_shape]
    return slide


@pytest.fixture
def mock_slide_with_zero_value_chart() -> Mock:
    """Create mock slide with chart that has all zero values.

    Returns:
        Mock slide with zero-value chart
    """
    slide = Mock()

    # Create chart with zero values
    chart_shape = Mock()
    chart_shape.shape_type = MSO_SHAPE_TYPE.CHART

    chart = Mock()
    chart.plots = [Mock()]  # Has plots

    # Create series with zero values
    series = Mock()
    series.values = [0, 0, 0, 0]
    chart.series = [series]

    chart_shape.chart = chart

    slide.shapes = [chart_shape]
    return slide


@pytest.fixture
def mock_slide_without_notes() -> Mock:
    """Create mock slide without speaker notes.

    Returns:
        Mock slide without notes
    """
    slide = Mock()
    slide.has_notes_slide = False
    slide.shapes = []
    return slide


@pytest.fixture
def mock_slide_with_empty_notes() -> Mock:
    """Create mock slide with empty speaker notes.

    Returns:
        Mock slide with empty notes
    """
    slide = Mock()
    slide.has_notes_slide = True

    notes_slide = Mock()
    notes_slide.notes_text_frame = Mock()
    notes_slide.notes_text_frame.text = ""

    slide.notes_slide = notes_slide
    slide.shapes = []
    return slide


class TestBulletLengthRule:
    """Tests for bullet length detection rule (QA-C-001)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = BulletLengthRule()

        assert rule.rule_id == "QA-C-001"
        assert "bullet" in rule.description.lower()
        assert rule.auto_fixable is False

    def test_detects_long_bullets(
        self, mock_presentation: Mock, mock_slide_with_long_bullets: Mock
    ) -> None:
        """Test detection of overly long bullet points."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_long_bullets]
        mock_presentation.prs = mock_prs

        rule = BulletLengthRule(max_length=120)
        issues = rule.validate(mock_presentation)

        # Should detect both long bullets
        assert len(issues) == 2
        assert all(issue.rule_id == "QA-C-001" for issue in issues)
        assert all(issue.severity == Severity.WARNING for issue in issues)
        assert all("exceeds" in issue.message.lower() for issue in issues)

    def test_no_issues_with_short_bullets(self, mock_presentation: Mock) -> None:
        """Test no issues when bullets are within limit."""
        slide = Mock()
        shape = Mock()
        shape.has_text_frame = True

        text_frame = Mock()
        para = Mock()
        para.level = 1
        para.text = "Short bullet"
        text_frame.paragraphs = [para]
        shape.text_frame = text_frame

        slide.shapes = [shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = BulletLengthRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0

    def test_custom_max_length(
        self, mock_presentation: Mock, mock_slide_with_long_bullets: Mock
    ) -> None:
        """Test custom max length threshold."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_long_bullets]
        mock_presentation.prs = mock_prs

        # Set very high threshold
        rule = BulletLengthRule(max_length=200)
        issues = rule.validate(mock_presentation)

        # Should detect fewer issues with higher threshold
        assert len(issues) < 2


class TestDuplicateTitleRule:
    """Tests for duplicate title detection rule (QA-C-002)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = DuplicateTitleRule()

        assert rule.rule_id == "QA-C-002"
        assert "duplicate" in rule.description.lower()
        assert rule.auto_fixable is False

    def test_detects_duplicate_titles(
        self, mock_presentation: Mock, mock_slides_with_duplicate_titles: list[Mock]
    ) -> None:
        """Test detection of duplicate slide titles."""
        mock_prs = Mock()
        mock_prs.slides = mock_slides_with_duplicate_titles
        mock_presentation.prs = mock_prs

        rule = DuplicateTitleRule()
        issues = rule.validate(mock_presentation)

        # Should detect duplicates on first two slides
        assert len(issues) == 2
        assert all(issue.rule_id == "QA-C-002" for issue in issues)
        assert all(issue.severity == Severity.INFO for issue in issues)
        assert all("duplicate" in issue.message.lower() for issue in issues)

    def test_no_issues_with_unique_titles(self, mock_presentation: Mock) -> None:
        """Test no issues when all titles are unique."""
        slides = []
        for i in range(3):
            slide = Mock()
            title_shape = Mock()
            title_shape.is_placeholder = True
            title_shape.placeholder_format = Mock()
            title_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.TITLE
            title_shape.has_text_frame = True
            title_shape.text_frame = Mock()
            title_shape.text_frame.text = f"Unique Title {i}"

            slide.shapes = [title_shape]
            slides.append(slide)

        mock_prs = Mock()
        mock_prs.slides = slides
        mock_presentation.prs = mock_prs

        rule = DuplicateTitleRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0


class TestUnpopulatedImagePlaceholderRule:
    """Tests for unpopulated image placeholder detection rule (QA-C-003)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = UnpopulatedImagePlaceholderRule()

        assert rule.rule_id == "QA-C-003"
        assert "image" in rule.description.lower()
        assert rule.auto_fixable is True

    def test_detects_empty_image_placeholder(
        self, mock_presentation: Mock, mock_slide_with_empty_image_placeholder: Mock
    ) -> None:
        """Test detection of unpopulated image placeholders."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_empty_image_placeholder]
        mock_presentation.prs = mock_prs

        rule = UnpopulatedImagePlaceholderRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-003"
        assert issues[0].severity == Severity.WARNING
        assert issues[0].auto_fixable is True
        assert "not populated" in issues[0].message.lower()

    def test_no_issues_with_populated_image(self, mock_presentation: Mock) -> None:
        """Test no issues when image placeholder is populated."""
        slide = Mock()
        image_shape = Mock()
        image_shape.is_placeholder = True
        image_shape.placeholder_format = Mock()
        image_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.PICTURE
        image_shape.shape_type = MSO_SHAPE_TYPE.PICTURE  # Populated with image

        slide.shapes = [image_shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = UnpopulatedImagePlaceholderRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0


class TestPathologicalTableDimensionRule:
    """Tests for pathological table dimension detection rule (QA-C-004)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = PathologicalTableDimensionRule()

        assert rule.rule_id == "QA-C-004"
        assert "table" in rule.description.lower()
        assert rule.auto_fixable is False

    def test_detects_1x1_table(
        self, mock_presentation: Mock, mock_slide_with_pathological_table: Mock
    ) -> None:
        """Test detection of 1x1 tables."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_pathological_table]
        mock_presentation.prs = mock_prs

        rule = PathologicalTableDimensionRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-004"
        assert issues[0].severity == Severity.WARNING
        assert "1x1" in issues[0].message

    def test_detects_large_table(
        self, mock_presentation: Mock, mock_slide_with_large_table: Mock
    ) -> None:
        """Test detection of very large tables."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_large_table]
        mock_presentation.prs = mock_prs

        rule = PathologicalTableDimensionRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-004"
        assert issues[0].severity == Severity.WARNING
        assert "12x10" in issues[0].message

    def test_no_issues_with_normal_table(self, mock_presentation: Mock) -> None:
        """Test no issues with normal-sized tables."""
        slide = Mock()
        table_shape = Mock()
        table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

        table = Mock()

        # Create 3x4 table
        rows = []
        for _ in range(3):
            row = Mock()
            cells = []
            for _ in range(4):
                cell = Mock()
                cell.text = "Data"
                cells.append(cell)
            row.cells = cells
            rows.append(row)

        table.rows = rows
        table.columns = [Mock() for _ in range(4)]

        table_shape.table = table
        slide.shapes = [table_shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = PathologicalTableDimensionRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0


class TestMissingChartDataRule:
    """Tests for missing chart data detection rule (QA-C-005)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = MissingChartDataRule()

        assert rule.rule_id == "QA-C-005"
        assert "chart" in rule.description.lower()
        assert rule.auto_fixable is True

    def test_detects_empty_chart(
        self, mock_presentation: Mock, mock_slide_with_empty_chart: Mock
    ) -> None:
        """Test detection of charts with no data."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_empty_chart]
        mock_presentation.prs = mock_prs

        rule = MissingChartDataRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-005"
        assert issues[0].severity == Severity.ERROR
        assert issues[0].auto_fixable is True
        assert "no data" in issues[0].message.lower()

    def test_detects_zero_value_chart(
        self, mock_presentation: Mock, mock_slide_with_zero_value_chart: Mock
    ) -> None:
        """Test detection of charts with all zero values."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_zero_value_chart]
        mock_presentation.prs = mock_prs

        rule = MissingChartDataRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-005"
        assert issues[0].severity == Severity.WARNING
        assert "zero values" in issues[0].message.lower()

    def test_no_issues_with_valid_chart(self, mock_presentation: Mock) -> None:
        """Test no issues with charts containing valid data."""
        slide = Mock()
        chart_shape = Mock()
        chart_shape.shape_type = MSO_SHAPE_TYPE.CHART

        chart = Mock()
        chart.plots = [Mock()]

        series = Mock()
        series.values = [10, 20, 30, 40]
        chart.series = [series]

        chart_shape.chart = chart
        slide.shapes = [chart_shape]

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = MissingChartDataRule()
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0


class TestSpeakerNotesVerificationRule:
    """Tests for speaker notes verification rule (QA-C-006)."""

    def test_rule_properties(self) -> None:
        """Test rule has correct properties."""
        rule = SpeakerNotesVerificationRule()

        assert rule.rule_id == "QA-C-006"
        assert "notes" in rule.description.lower()
        assert rule.auto_fixable is True

    def test_detects_missing_notes_when_required(
        self, mock_presentation: Mock, mock_slide_without_notes: Mock
    ) -> None:
        """Test detection of missing speaker notes when required."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_without_notes]
        mock_presentation.prs = mock_prs

        rule = SpeakerNotesVerificationRule(require_notes=True)
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-006"
        assert issues[0].severity == Severity.INFO
        assert issues[0].auto_fixable is True
        assert "missing" in issues[0].message.lower()

    def test_detects_empty_notes_when_required(
        self, mock_presentation: Mock, mock_slide_with_empty_notes: Mock
    ) -> None:
        """Test detection of empty speaker notes when required."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_with_empty_notes]
        mock_presentation.prs = mock_prs

        rule = SpeakerNotesVerificationRule(require_notes=True)
        issues = rule.validate(mock_presentation)

        assert len(issues) == 1
        assert issues[0].rule_id == "QA-C-006"
        assert issues[0].severity == Severity.INFO
        assert "empty" in issues[0].message.lower()

    def test_no_issues_when_notes_not_required(
        self, mock_presentation: Mock, mock_slide_without_notes: Mock
    ) -> None:
        """Test no issues when speaker notes are not required."""
        mock_prs = Mock()
        mock_prs.slides = [mock_slide_without_notes]
        mock_presentation.prs = mock_prs

        rule = SpeakerNotesVerificationRule(require_notes=False)
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0

    def test_no_issues_with_valid_notes(self, mock_presentation: Mock) -> None:
        """Test no issues when slides have valid speaker notes."""
        slide = Mock()
        slide.has_notes_slide = True

        notes_slide = Mock()
        notes_slide.notes_text_frame = Mock()
        notes_slide.notes_text_frame.text = "Valid speaker notes content"

        slide.notes_slide = notes_slide
        slide.shapes = []

        mock_prs = Mock()
        mock_prs.slides = [slide]
        mock_presentation.prs = mock_prs

        rule = SpeakerNotesVerificationRule(require_notes=True)
        issues = rule.validate(mock_presentation)

        assert len(issues) == 0


# Made with Bob
