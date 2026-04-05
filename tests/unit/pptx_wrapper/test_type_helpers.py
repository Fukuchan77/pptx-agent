"""Tests for type ignore helper functions.

These helpers encapsulate operations that require type: ignore comments
due to incomplete type stubs in python-pptx library.
"""

from unittest.mock import MagicMock, Mock

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx_agent.pptx_wrapper.type_helpers import (
    add_chart_with_bounds,
    add_picture_with_bounds,
    add_table_with_bounds,
    set_text_frame_text,
)


class TestAddChartWithBounds:
    """Test add_chart_with_bounds helper."""

    def test_adds_chart_with_position_bounds(self):
        """Test that add_chart_with_bounds calls add_chart with correct parameters."""
        # Arrange
        mock_slide = Mock()
        mock_shapes = Mock()
        mock_slide.shapes = mock_shapes

        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        left = Inches(1)
        top = Inches(2)
        width = Inches(5)
        height = Inches(4)
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Series1", [1, 2])

        # Act
        add_chart_with_bounds(mock_slide, chart_type, left, top, width, height, chart_data)

        # Assert
        mock_shapes.add_chart.assert_called_once_with(
            chart_type, left, top, width, height, chart_data
        )

    def test_handles_all_chart_types(self):
        """Test that helper works with different chart types."""
        # Arrange
        mock_slide = Mock()
        mock_slide.shapes = Mock()

        chart_data = CategoryChartData()
        chart_data.categories = ["X", "Y"]
        chart_data.add_series("Test", [10, 20])

        # Act - test different chart types
        for chart_type in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.PIE, XL_CHART_TYPE.BAR_CLUSTERED]:
            add_chart_with_bounds(
                mock_slide,
                chart_type,
                Inches(0),
                Inches(0),
                Inches(6),
                Inches(4),
                chart_data,
            )

        # Assert
        assert mock_slide.shapes.add_chart.call_count == 3


class TestAddTableWithBounds:
    """Test add_table_with_bounds helper."""

    def test_adds_table_with_position_bounds(self):
        """Test that add_table_with_bounds calls add_table with correct parameters."""
        # Arrange
        mock_slide = Mock()
        mock_shapes = Mock()
        mock_table = Mock()
        mock_graphic_frame = Mock()
        mock_graphic_frame.table = mock_table
        mock_shapes.add_table.return_value = mock_graphic_frame
        mock_slide.shapes = mock_shapes

        rows = 3
        cols = 4
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(3)

        # Act
        result = add_table_with_bounds(mock_slide, rows, cols, left, top, width, height)

        # Assert
        mock_shapes.add_table.assert_called_once_with(rows, cols, left, top, width, height)
        assert result == mock_table

    def test_returns_table_shape(self):
        """Test that helper returns the table shape object."""
        # Arrange
        mock_slide = Mock()
        mock_table = MagicMock()
        mock_graphic_frame = Mock()
        mock_graphic_frame.table = mock_table
        mock_slide.shapes.add_table.return_value = mock_graphic_frame

        # Act
        result = add_table_with_bounds(mock_slide, 2, 3, Inches(0), Inches(0), Inches(5), Inches(3))

        # Assert
        assert result is mock_table


class TestAddPictureWithBounds:
    """Test add_picture_with_bounds helper."""

    def test_adds_picture_with_position_bounds(self):
        """Test that add_picture_with_bounds calls add_picture with correct parameters."""
        # Arrange
        mock_slide = Mock()
        mock_shapes = Mock()
        mock_pic = Mock()
        mock_shapes.add_picture.return_value = mock_pic
        mock_slide.shapes = mock_shapes

        image_path = "/path/to/image.png"
        left = Inches(2)
        top = Inches(1)
        width = Inches(4)
        height = Inches(3)

        # Act
        result = add_picture_with_bounds(mock_slide, image_path, left, top, width, height)

        # Assert
        mock_shapes.add_picture.assert_called_once_with(
            image_path, left, top, width=width, height=height
        )
        assert result == mock_pic

    def test_returns_picture_shape(self):
        """Test that helper returns the picture shape object."""
        # Arrange
        mock_slide = Mock()
        mock_pic = MagicMock()
        mock_slide.shapes.add_picture.return_value = mock_pic

        # Act
        result = add_picture_with_bounds(
            mock_slide, "test.jpg", Inches(0), Inches(0), Inches(3), Inches(2)
        )

        # Assert
        assert result is mock_pic


class TestSetTextFrameText:
    """Test set_text_frame_text helper."""

    def test_sets_text_on_text_frame(self):
        """Test that set_text_frame_text sets text on placeholder text_frame."""
        # Arrange
        mock_placeholder = Mock()
        mock_text_frame = Mock()
        mock_placeholder.text_frame = mock_text_frame
        text = "Test content"

        # Act
        set_text_frame_text(mock_placeholder, text)

        # Assert
        assert mock_text_frame.text == text

    def test_works_with_multiline_text(self):
        """Test that helper works with multiline text."""
        # Arrange
        mock_placeholder = Mock()
        mock_text_frame = Mock()
        mock_placeholder.text_frame = mock_text_frame
        text = "Line 1\nLine 2\nLine 3"

        # Act
        set_text_frame_text(mock_placeholder, text)

        # Assert
        assert mock_text_frame.text == text

    def test_works_with_empty_text(self):
        """Test that helper handles empty text."""
        # Arrange
        mock_placeholder = Mock()
        mock_text_frame = Mock()
        mock_placeholder.text_frame = mock_text_frame

        # Act
        set_text_frame_text(mock_placeholder, "")

        # Assert
        assert mock_text_frame.text == ""
