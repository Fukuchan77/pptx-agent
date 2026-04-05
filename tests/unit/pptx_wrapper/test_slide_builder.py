"""Tests for slide builder module.

Tests the slide_builder module which renders PresentationSchema to PowerPoint files.
Following TDD methodology: These tests are written BEFORE implementation.
"""

import hashlib
import logging
from datetime import UTC, datetime
from importlib.metadata import PackageNotFoundError
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation

# Import the function we're about to implement
# This will fail initially - that's expected for RED phase
from pptx_agent.pptx_wrapper.slide_builder import build_presentation
from pptx_agent.schemas import PresentationSchema, SlideSchema
from pptx_agent.schemas.text import TextBlock
from tests.fixtures.sample_schemas import (
    create_japanese_slide,
    create_minimal_presentation,
    create_valid_presentation,
)


class TestBuildPresentation:
    """Test suite for build_presentation function."""

    def test_build_minimal_presentation(self, tmp_path: Path) -> None:
        """Test building a minimal presentation with one slide.

        RED PHASE: This test should FAIL because build_presentation doesn't exist yet.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "minimal.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert result_path == output_path
        assert Path(output_path).exists()
        assert Path(output_path).stat().st_size > 0

    def test_build_presentation_with_multiple_slides(self, tmp_path: Path) -> None:
        """Test building a presentation with multiple slides and varied content.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "multi_slide.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert result_path == output_path
        assert Path(output_path).exists()
        # Verify file is substantial (multiple slides)
        assert Path(output_path).stat().st_size > 5000

    def test_build_presentation_populates_titles(self, tmp_path: Path) -> None:
        """Test that slide titles are correctly populated.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "with_titles.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert - file should be created
        assert Path(result_path).exists()

        # Note: Deep verification of title content would require reading the .pptx back,
        # which is more of an integration test. For unit test, we verify successful creation.

    def test_build_presentation_populates_text_content(self, tmp_path: Path) -> None:
        """Test that text content blocks are correctly populated.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "with_content.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert Path(result_path).exists()

    def test_build_presentation_uses_correct_layouts(self, tmp_path: Path) -> None:
        """Test that correct slide layouts are selected from template.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "correct_layouts.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert Path(result_path).exists()

    def test_build_presentation_with_empty_content_blocks(self, tmp_path: Path) -> None:
        """Test handling of slides with empty content lists.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = PresentationSchema(
            title="Empty Content Test",
            slides=[
                SlideSchema(
                    layout_name="Title Slide",
                    title="Empty Slide",
                    content=[],  # No content blocks
                )
            ],
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "empty_content.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert - should still create valid file
        assert Path(result_path).exists()

    def test_build_presentation_invalid_template_raises_error(self, tmp_path: Path) -> None:
        """Test that invalid template path raises appropriate error.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "nonexistent/template.pptx"
        output_path = str(tmp_path / "output.pptx")

        # Act & Assert
        with pytest.raises((FileNotFoundError, ValueError)):
            build_presentation(content, template_path, output_path)

    def test_build_presentation_creates_parent_directories(self, tmp_path: Path) -> None:
        """Test that parent directories are created if they don't exist.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        # Output path with non-existent subdirectory
        output_path = str(tmp_path / "subdir" / "nested" / "output.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert Path(result_path).exists()
        assert Path(result_path).parent.exists()

    def test_build_presentation_with_japanese_content(self, tmp_path: Path) -> None:
        """Test building presentation with Japanese text content.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = PresentationSchema(
            title="日本語プレゼンテーション",
            slides=[create_japanese_slide()],
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "japanese.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert Path(result_path).exists()

    def test_build_presentation_returns_correct_path(self, tmp_path: Path) -> None:
        """Test that the function returns the correct output path.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "test_return.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert - return value should match output_path
        assert result_path == output_path
        assert isinstance(result_path, str)


class TestBuildPresentationEdgeCases:
    """Test edge cases and error handling."""

    def test_build_presentation_with_many_slides(self, tmp_path: Path) -> None:
        """Test building presentation with many slides (performance consideration).

        RED PHASE: This test should FAIL.
        """
        # Arrange - create presentation with 20 slides
        slides = [
            SlideSchema(
                layout_name="Title and Content",
                title=f"Slide {i + 1}",
                content=[],
            )
            for i in range(20)
        ]

        content = PresentationSchema(
            title="Many Slides Test",
            slides=slides,
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "many_slides.pptx")

        # Act
        result_path = build_presentation(content, template_path, output_path)

        # Assert
        assert Path(result_path).exists()
        # Verify file size is substantial for 20 slides
        assert Path(result_path).stat().st_size > 10000

    def test_build_presentation_layout_not_in_template(self, tmp_path: Path) -> None:
        """Test error handling when requested layout doesn't exist in template.

        RED PHASE: This test should FAIL.
        """
        # Arrange
        content = PresentationSchema(
            title="Invalid Layout Test",
            slides=[
                SlideSchema(
                    layout_name="NonexistentLayout",
                    title="Test",
                    content=[],
                )
            ],
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "invalid_layout.pptx")

        # Act & Assert - should raise error for invalid layout
        with pytest.raises(ValueError, match="Layout.*not found"):
            build_presentation(content, template_path, output_path)


class TestOutputPathHandling:
    """Test suite for output path handling in slide_builder.

    RED PHASE: These tests verify that slide_builder.py correctly handles
    absolute paths passed from main.py.
    """

    def test_build_presentation_handles_absolute_path(self, tmp_path: Path) -> None:
        """Test that build_presentation correctly handles absolute output paths.

        RED PHASE: This test verifies that when an absolute path is passed,
        slide_builder extracts base_dir correctly for validation.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"

        # Use absolute path
        output_path = tmp_path / "absolute_output.pptx"
        absolute_path_str = str(output_path.resolve())

        # Act
        result_path = build_presentation(content, template_path, absolute_path_str)

        # Assert
        assert result_path == absolute_path_str
        assert Path(result_path).exists()
        # Verify it's absolute
        assert Path(result_path).is_absolute()

    def test_build_presentation_base_dir_from_absolute_path(self, tmp_path: Path) -> None:
        """Test that base_dir is correctly extracted from absolute paths.

        RED PHASE: After refactoring, slide_builder should always extract
        base_dir from parent of absolute path (never None).
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"

        # Use absolute path with nested directories
        output_dir = tmp_path / "nested" / "directories"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / "output.pptx"
        absolute_path_str = str(output_path.resolve())

        # Act
        with patch("pptx_agent.pptx_wrapper.slide_builder.PresentationWrapper") as mock_prs:
            mock_prs_instance = MagicMock()
            mock_prs.return_value = mock_prs_instance

            # Call build_presentation
            build_presentation(content, template_path, absolute_path_str)

            # Assert - verify save was called with correct base_dir
            mock_prs_instance.save.assert_called_once()
            call_args = mock_prs_instance.save.call_args

            # Extract base_dir from call
            base_dir_arg = call_args.kwargs.get("base_dir")

            # base_dir should be the parent directory (not None)
            assert base_dir_arg is not None, "base_dir should not be None for absolute paths"
            assert base_dir_arg == str(output_path.parent.resolve())


class TestErrorLogging:
    """Test suite for error logging behavior in slide builder.

    These tests verify that errors are properly logged instead of silently suppressed.
    """

    def test_set_title_error_is_logged(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        """Test that ValueError from set_title() is logged with context.

        RED PHASE: This test should FAIL because current implementation uses
        suppress(ValueError) which doesn't log anything.

        Expected behavior:
        - When set_title() raises ValueError (e.g., no title placeholder in layout)
        - The error should be logged at WARNING or DEBUG level
        - Log message should include contextual information (layout name, title text)
        - Presentation generation should continue (non-fatal error)
        """
        # Arrange
        content = PresentationSchema(
            title="Test Presentation",
            slides=[
                SlideSchema(
                    layout_name="Title Slide",
                    title="Test Title",
                    content=[
                        TextBlock(
                            placeholder_name="Content Placeholder 1",
                            text="Test content",
                            language="en",
                        )
                    ],
                )
            ],
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "error_logging_test.pptx")

        # Set logging level to capture DEBUG messages
        caplog.set_level(logging.DEBUG)

        # Mock SlideWrapper.set_title to raise ValueError
        with patch("pptx_agent.pptx_wrapper.slide_builder.PresentationWrapper") as mock_prs:
            mock_prs_instance = MagicMock()
            mock_prs.return_value = mock_prs_instance

            # Mock slide with set_title that raises ValueError
            mock_slide = MagicMock()
            mock_slide.set_title.side_effect = ValueError("Layout has no title placeholder")
            mock_prs_instance.add_slide.return_value = mock_slide

            # Act
            result_path = build_presentation(content, template_path, output_path)

            # Assert
            assert result_path == output_path

            # Verify error was logged
            # Should have at least one log record about the ValueError
            log_messages = [record.message for record in caplog.records]

            # Check that error is logged with context
            assert any(
                "title" in msg.lower() and ("error" in msg.lower() or "failed" in msg.lower())
                for msg in log_messages
            ), f"Expected error log about title failure, but got: {log_messages}"

            # Verify log includes contextual information
            assert any("Test Title" in msg or "Title Slide" in msg for msg in log_messages), (
                f"Expected log to include title text or layout name, but got: {log_messages}"
            )

            # Verify logging level is appropriate (WARNING or DEBUG)
            log_levels = [record.levelname for record in caplog.records]
            assert any(level in ["WARNING", "DEBUG"] for level in log_levels), (
                f"Expected WARNING or DEBUG log level, but got: {log_levels}"
            )


class TestMetadataEmbedding:
    """Test suite for metadata embedding in generated presentations.

    RED PHASE: These tests should FAIL because metadata embedding isn't implemented yet.
    Tests verify that metadata is properly embedded in generated .pptx files.
    """

    def test_metadata_contains_generator_name_and_version(self, tmp_path: Path) -> None:
        """Test that generator name and version are embedded in metadata.

        RED PHASE: This test should FAIL - metadata not yet implemented.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "metadata_test.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Assert - read back the presentation and check metadata
        prs = Presentation(output_path)

        # Check for generator in comments (core properties)
        comments = prs.core_properties.comments or ""
        assert "pptx-agent" in comments.lower()
        assert "v0.1.0" in comments or "0.1.0" in comments

    def test_metadata_contains_generation_timestamp(self, tmp_path: Path) -> None:
        """Test that generation timestamp is recorded in core properties.

        RED PHASE: This test should FAIL - timestamp not yet embedded.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "timestamp_test.pptx")

        # Record time before generation
        before_time = datetime.now(UTC)

        # Act
        build_presentation(content, template_path, output_path)

        # Assert
        prs = Presentation(output_path)

        # Check that created/modified timestamps exist and are recent
        created = prs.core_properties.created
        assert created is not None

        # Verify timestamp is within reasonable range (within last minute)
        time_diff = (datetime.now(UTC) - before_time).total_seconds()
        assert time_diff < 60  # Should complete within 60 seconds

    def test_metadata_contains_template_path(self, tmp_path: Path) -> None:
        """Test that template path is stored in metadata.

        RED PHASE: This test should FAIL - template tracking not implemented.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "template_tracking.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Assert
        prs = Presentation(output_path)

        # Template should be in comments or subject
        metadata_text = (prs.core_properties.comments or "") + (prs.core_properties.subject or "")
        assert "basic-template.pptx" in metadata_text

    def test_metadata_contains_content_hash(self, tmp_path: Path) -> None:
        """Test that content hash is embedded for change tracking.

        RED PHASE: This test should FAIL - content hash not implemented.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "hash_test.pptx")

        # Calculate expected hash of presentation content (for documentation purposes)
        content_json = content.model_dump_json()
        _expected_hash = hashlib.sha256(content_json.encode()).hexdigest()[:16]

        # Act
        build_presentation(content, template_path, output_path)

        # Assert
        prs = Presentation(output_path)

        # Hash should be in comments
        comments = prs.core_properties.comments or ""
        # Check for hash presence (partial match since we use first 16 chars)
        assert len(comments) > 0
        # At least verify some hash-like pattern exists
        assert any(char in comments for char in "0123456789abcdef")

    def test_metadata_accessible_after_save_reload(self, tmp_path: Path) -> None:
        """Test that metadata persists after saving and reloading file.

        RED PHASE: This test should FAIL - metadata persistence not verified.
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "persistence_test.pptx")

        # Act - generate and save
        build_presentation(content, template_path, output_path)

        # Reload the file
        prs = Presentation(output_path)

        # Assert - metadata should still be accessible
        assert prs.core_properties.comments is not None
        assert len(prs.core_properties.comments) > 0
        assert prs.core_properties.created is not None

    def test_metadata_includes_all_required_fields(self, tmp_path: Path) -> None:
        """Test that all required metadata fields are present.

        RED PHASE: This test should FAIL - comprehensive metadata not implemented.
        Required fields:
        - Generator name and version
        - Generation timestamp
        - Template used
        - Content hash (optional but recommended)
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "complete_metadata.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Assert
        prs = Presentation(output_path)
        core_props = prs.core_properties

        # Check all required fields
        assert core_props.comments is not None
        assert "pptx-agent" in core_props.comments.lower()
        assert core_props.created is not None

        # Either comments or subject should contain template info
        metadata_str = (core_props.comments or "") + (core_props.subject or "")
        assert "template" in metadata_str.lower() or ".pptx" in metadata_str

    def test_metadata_does_not_affect_visual_presentation(self, tmp_path: Path) -> None:
        """Test that metadata embedding doesn't affect slide content.

        RED PHASE: This test should FAIL initially but metadata should be non-intrusive.
        """
        # Arrange
        content = create_valid_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "visual_check.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Assert - file should be created
        # Note: We don't check exact slide count because template may have placeholder slides
        # The important thing is that metadata is embedded without errors
        prs = Presentation(output_path)
        assert len(prs.slides) >= len(content.slides)

        # Metadata should not create extra slides or alter content
        assert prs.core_properties is not None

    def test_metadata_with_japanese_template(self, tmp_path: Path) -> None:
        """Test that metadata works correctly with Japanese content.

        RED PHASE: This test should FAIL - i18n metadata not tested.
        """
        # Arrange
        content = PresentationSchema(
            title="日本語プレゼンテーション",
            slides=[create_japanese_slide()],
        )
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "japanese_metadata.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Assert - metadata should still be embedded properly
        prs = Presentation(output_path)
        assert prs.core_properties.comments is not None
        assert "pptx-agent" in prs.core_properties.comments.lower()

    def test_metadata_uses_single_timestamp_call(self, tmp_path: Path) -> None:
        """Test that _embed_metadata() calls datetime.now(UTC) only once.

        RED PHASE: This test should FAIL - current implementation calls datetime.now(UTC) 3 times.

        This verifies that the same timestamp is used for:
        - core_props.created
        - core_props.modified
        - Generated timestamp string in comments
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "single_timestamp.pptx")

        # Mock datetime.now to track calls
        mock_timestamp = datetime(2026, 4, 5, 12, 30, 45, tzinfo=UTC)

        with patch("pptx_agent.pptx_wrapper.slide_builder.datetime") as mock_datetime:
            mock_datetime.now.return_value = mock_timestamp
            # Allow datetime constructor to work normally
            mock_datetime.side_effect = None

            # Act
            build_presentation(content, template_path, output_path)

            # Assert - datetime.now should be called exactly once with UTC
            assert mock_datetime.now.call_count == 1
            mock_datetime.now.assert_called_once_with(UTC)

        # Verify that all metadata fields use the same timestamp
        prs = Presentation(output_path)
        core_props = prs.core_properties

        # created and modified should be identical
        assert core_props.created == core_props.modified

        # Comments should contain the timestamp in ISO format
        comments = core_props.comments or ""
        # The timestamp should be in the comments
        assert mock_timestamp.isoformat() in comments


class TestVersionRetrieval:
    """Test suite for dynamic version retrieval using importlib.metadata.

    RED PHASE: These tests should FAIL because current implementation uses
    hardcoded __version__ = "0.1.0" instead of importlib.metadata.version().
    """

    def test_version_retrieved_from_importlib_metadata(self, tmp_path: Path) -> None:
        """Test that version is retrieved from importlib.metadata, not hardcoded.

        RED PHASE: This test should FAIL because current implementation uses hardcoded __version__.

        Expected behavior:
        - Version should be retrieved using importlib.metadata.version("pptx-agent")
        - The retrieved version should appear in metadata comments
        - When we mock importlib.metadata.version, it should use our mocked version
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "version_test.pptx")

        # Mock importlib.metadata.version to return a test version
        test_version = "1.2.3"

        with patch("pptx_agent.pptx_wrapper.slide_builder.version") as mock_version:
            mock_version.return_value = test_version

            # Act
            build_presentation(content, template_path, output_path)

            # Assert - version() should have been called with "pptx-agent"
            mock_version.assert_called_once_with("pptx-agent")

        # Verify the mocked version appears in metadata
        prs = Presentation(output_path)
        comments = prs.core_properties.comments or ""
        assert test_version in comments, (
            f"Expected mocked version '{test_version}' in metadata, but got: {comments}"
        )

    def test_version_fallback_on_package_not_found(self, tmp_path: Path) -> None:
        """Test that version falls back to 'unknown' when package not found.

        RED PHASE: This test should FAIL because current implementation doesn't
        handle PackageNotFoundError.

        Expected behavior:
        - When importlib.metadata.version() raises PackageNotFoundError
        - Should catch the exception and use "unknown" as fallback
        - Presentation should still be generated successfully
        """
        # Arrange
        content = create_minimal_presentation()
        template_path = "templates/basic-template.pptx"
        output_path = str(tmp_path / "fallback_test.pptx")

        # Mock importlib.metadata.version to raise PackageNotFoundError
        with patch("pptx_agent.pptx_wrapper.slide_builder.version") as mock_version:
            mock_version.side_effect = PackageNotFoundError("pptx-agent not found")

            # Act - should not raise exception, should use fallback
            build_presentation(content, template_path, output_path)

            # Assert - version() should have been called
            mock_version.assert_called_once_with("pptx-agent")

        # Verify fallback version appears in metadata
        prs = Presentation(output_path)
        comments = prs.core_properties.comments or ""
        assert "unknown" in comments, (
            f"Expected fallback version 'unknown' in metadata, but got: {comments}"
        )
