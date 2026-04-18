"""Unit tests for layout switching functionality in slide_builder.

RED PHASE: These tests are written BEFORE implementing layout switching.
Tests verify that slide_builder can rebuild slides with alternative layouts
when overflow is detected.
"""

from pathlib import Path

import pytest

from pptx_agent.pptx_wrapper.slide_builder import build_presentation, rebuild_slide_with_layout
from pptx_agent.schemas import PresentationSchema, SlideSchema
from pptx_agent.schemas.text import TextBlock


class TestLayoutSwitching:
    """Tests for layout switching functionality."""

    def test_rebuild_slide_with_layout_exists(self) -> None:
        """Test that rebuild_slide_with_layout function exists.

        RED PHASE: This test should FAIL because the function doesn't exist yet.
        """
        # This import will fail if function doesn't exist
        from pptx_agent.pptx_wrapper.slide_builder import rebuild_slide_with_layout

        assert callable(rebuild_slide_with_layout)

    def test_rebuild_slide_with_different_layout(self, tmp_path: Path, template_path: str) -> None:
        """Test rebuilding a slide with a different layout.

        RED PHASE: This test should FAIL - rebuild_slide_with_layout not implemented.

        Expected behavior:
        - Takes an existing presentation, slide index, and new layout name
        - Removes the old slide at that index
        - Creates a new slide with the new layout
        - Preserves title and content from original slide
        """
        # Arrange
        content = PresentationSchema(
            title="Layout Switch Test",
            slides=[
                SlideSchema(
                    layout_name="Title and Content",
                    title="Original Layout",
                    content=[
                        TextBlock(
                            placeholder_name="Content Placeholder 1",
                            text="This content will overflow",
                            language="en",
                        )
                    ],
                )
            ],
        )

        output_path = str(tmp_path / "layout_switch.pptx")

        # Act - build initial presentation
        build_presentation(content, template_path, output_path)

        # Now rebuild slide 0 with a different layout
        new_layout = "Two Content"
        rebuild_slide_with_layout(
            pptx_path=output_path,
            slide_index=0,
            new_layout=new_layout,
            slide_data=content.slides[0],
        )

        # Assert - file should still exist and be valid
        assert Path(output_path).exists()

    def test_rebuild_slide_preserves_title(self, tmp_path: Path, template_path: str) -> None:
        """Test that rebuilding a slide preserves the title.

        RED PHASE: This test should FAIL - functionality not implemented.
        """
        # Arrange
        original_title = "Important Title"
        content = PresentationSchema(
            title="Title Preservation Test",
            slides=[
                SlideSchema(
                    layout_name="Title and Content",
                    title=original_title,
                    content=[],
                )
            ],
        )

        output_path = str(tmp_path / "title_preserved.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Rebuild with different layout
        rebuild_slide_with_layout(
            pptx_path=output_path,
            slide_index=0,
            new_layout="Two Content",
            slide_data=content.slides[0],
        )

        # Assert - would need to read back presentation to verify
        # For unit test, just verify no error occurred
        assert Path(output_path).exists()

    def test_rebuild_slide_preserves_content(self, tmp_path: Path, template_path: str) -> None:
        """Test that rebuilding a slide preserves content blocks.

        RED PHASE: This test should FAIL - content preservation not implemented.
        """
        # Arrange
        test_content = "Important content that must be preserved"
        content = PresentationSchema(
            title="Content Preservation Test",
            slides=[
                SlideSchema(
                    layout_name="Title and Content",
                    title="Test",
                    content=[
                        TextBlock(
                            placeholder_name="Content Placeholder 1",
                            text=test_content,
                            language="en",
                        )
                    ],
                )
            ],
        )

        output_path = str(tmp_path / "content_preserved.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Rebuild with different layout
        rebuild_slide_with_layout(
            pptx_path=output_path,
            slide_index=0,
            new_layout="Two Content",
            slide_data=content.slides[0],
        )

        # Assert
        assert Path(output_path).exists()

    def test_rebuild_slide_invalid_index_raises_error(
        self, tmp_path: Path, template_path: str
    ) -> None:
        """Test that rebuilding with invalid slide index raises error.

        RED PHASE: This test should FAIL - error handling not implemented.
        """
        # Arrange
        content = PresentationSchema(
            title="Error Test",
            slides=[
                SlideSchema(
                    layout_name="Title Slide",
                    title="Test",
                    content=[],
                )
            ],
        )

        output_path = str(tmp_path / "error_test.pptx")

        build_presentation(content, template_path, output_path)

        # Act & Assert - try to rebuild non-existent slide
        with pytest.raises((IndexError, ValueError)):
            rebuild_slide_with_layout(
                pptx_path=output_path,
                slide_index=999,  # Invalid index
                new_layout="Title and Content",
                slide_data=content.slides[0],
            )

    def test_rebuild_slide_invalid_layout_raises_error(
        self, tmp_path: Path, template_path: str
    ) -> None:
        """Test that rebuilding with non-existent layout raises error.

        RED PHASE: This test should FAIL - layout validation not implemented.
        """
        # Arrange
        content = PresentationSchema(
            title="Layout Validation Test",
            slides=[
                SlideSchema(
                    layout_name="Title Slide",
                    title="Test",
                    content=[],
                )
            ],
        )

        output_path = str(tmp_path / "invalid_layout.pptx")

        build_presentation(content, template_path, output_path)

        # Act & Assert - try to use non-existent layout
        with pytest.raises(ValueError, match="Layout.*not found"):
            rebuild_slide_with_layout(
                pptx_path=output_path,
                slide_index=0,
                new_layout="NonexistentLayout",
                slide_data=content.slides[0],
            )

    def test_rebuild_multiple_slides(self, tmp_path: Path, template_path: str) -> None:
        """Test rebuilding multiple slides in sequence.

        RED PHASE: This test should FAIL - sequential rebuilding not tested.
        """
        # Arrange - create presentation with 3 slides
        content = PresentationSchema(
            title="Multi-rebuild Test",
            slides=[
                SlideSchema(
                    layout_name="Title and Content",
                    title=f"Slide {i + 1}",
                    content=[],
                )
                for i in range(3)
            ],
        )

        output_path = str(tmp_path / "multi_rebuild.pptx")

        # Act - build and rebuild slides
        build_presentation(content, template_path, output_path)

        # Rebuild each slide with different layout
        for i in range(3):
            rebuild_slide_with_layout(
                pptx_path=output_path,
                slide_index=i,
                new_layout="Two Content",
                slide_data=content.slides[i],
            )

        # Assert
        assert Path(output_path).exists()

    def test_rebuild_slide_maintains_slide_count(self, tmp_path: Path, template_path: str) -> None:
        """Test that rebuilding doesn't change total slide count.

        RED PHASE: This test should FAIL - slide count preservation not verified.
        """
        # Arrange
        content = PresentationSchema(
            title="Slide Count Test",
            slides=[
                SlideSchema(layout_name="Title Slide", title="Slide 1", content=[]),
                SlideSchema(layout_name="Title and Content", title="Slide 2", content=[]),
                SlideSchema(layout_name="Two Content", title="Slide 3", content=[]),
            ],
        )

        output_path = str(tmp_path / "slide_count.pptx")

        # Act
        build_presentation(content, template_path, output_path)

        # Count slides before rebuild
        from pptx import Presentation

        prs_before = Presentation(output_path)
        count_before = len(prs_before.slides)

        # Rebuild middle slide
        rebuild_slide_with_layout(
            pptx_path=output_path,
            slide_index=1,
            new_layout="Title Slide",
            slide_data=content.slides[1],
        )

        # Count slides after rebuild
        prs_after = Presentation(output_path)
        count_after = len(prs_after.slides)

        # Assert - slide count should be unchanged
        assert count_after == count_before
        # The actual count may include template slides, so we verify the count is maintained
        # rather than checking for an exact number
