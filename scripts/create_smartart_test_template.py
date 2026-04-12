"""
Create a minimal SmartArt test template.

This script creates a basic template with layouts that have OBJECT placeholders
where SmartArt diagrams would be placed. This serves as a test template until
actual SmartArt diagrams can be manually created in PowerPoint.

Note: python-pptx cannot create actual SmartArt diagrams. This template provides
the layout structure with OBJECT placeholders for testing purposes.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# Template output directory
TEMPLATES_DIR = Path(__file__).parent.parent / "templates"


def _add_smartart_test_slide(prs, title: str, smartart_type: str, node_count_range: str) -> None:  # noqa: ANN001
    """Add a single SmartArt test slide to the presentation.

    Args:
        prs: Presentation object to add the slide to
        title: Title text for the slide
        smartart_type: Type of SmartArt (e.g., "Process Flow", "Hierarchy")
        node_count_range: Node count range (e.g., "3-5 nodes")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True

    # Add placeholder for SmartArt (represented as a rectangle)
    smartart_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(4),
    )
    smartart_placeholder.fill.background()
    smartart_placeholder.line.color.rgb = RGBColor(0, 0, 0)

    # Add text to indicate this is a SmartArt placeholder
    text_frame = smartart_placeholder.text_frame
    text_frame.text = f"SmartArt Placeholder: {smartart_type}\n({node_count_range})"
    text_frame.paragraphs[0].font.size = Pt(18)


def create_smartart_test_template(output_path: str) -> None:
    """Create a minimal SmartArt test template with OBJECT placeholders.

    This creates a basic template structure that can be used for testing
    the template parser and layout detection. Actual SmartArt diagrams
    must be added manually in PowerPoint.

    Args:
        output_path: Path where the template should be saved
    """
    prs = Presentation()

    # Note: We use the default layouts provided by python-pptx
    # Custom layouts cannot be created programmatically

    logger.info("Creating SmartArt test template...")

    # We'll create test slides with different layouts
    # Note: python-pptx has limitations with custom layouts, so we'll create
    # regular slides with shapes positioned where SmartArt would be

    # Layout 1: Process Flow
    logger.info("  - Creating Process Flow layout test slide...")
    _add_smartart_test_slide(prs, "Process Flow Test Layout", "Process Flow", "3-5 nodes")

    # Layout 2: Hierarchy
    logger.info("  - Creating Hierarchy layout test slide...")
    _add_smartart_test_slide(prs, "Hierarchy Test Layout", "Hierarchy", "3-7 nodes")

    # Layout 3: Cycle
    logger.info("  - Creating Cycle layout test slide...")
    _add_smartart_test_slide(prs, "Cycle Test Layout", "Cycle", "3-6 nodes")

    # Layout 4: Relationship
    logger.info("  - Creating Relationship layout test slide...")
    _add_smartart_test_slide(prs, "Relationship Test Layout", "Relationship", "2-4 nodes")

    # Save the presentation
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))

    logger.info("✓ Created test template at: %s", output_path)
    logger.info("")
    logger.info("IMPORTANT: This is a TEST template with placeholder shapes.")
    logger.info("For production use, actual SmartArt diagrams must be added manually:")
    logger.info("  1. Open this file in Microsoft PowerPoint")
    logger.info("  2. Go to View > Slide Master")
    logger.info("  3. Replace placeholder shapes with actual SmartArt diagrams")
    logger.info("  4. Follow specifications in templates/SMARTART-TEMPLATE-SPEC.md")
    logger.info("  5. Save as smartart-template.pptx")


def main() -> None:
    """CLI entry point for creating SmartArt test template."""
    logging.basicConfig(
        level=logging.INFO,
        format="%(message)s",
    )

    logger.info("=" * 70)
    logger.info("SmartArt Test Template Creator")
    logger.info("=" * 70)
    logger.info("")

    output_path = TEMPLATES_DIR / "smartart-test-template.pptx"
    create_smartart_test_template(str(output_path))

    logger.info("")
    logger.info("=" * 70)
    logger.info("Summary")
    logger.info("=" * 70)
    logger.info("Test template created: %s", output_path)
    logger.info("")
    logger.info("This template contains placeholder shapes for testing purposes.")
    logger.info("See templates/SMARTART-TEMPLATE-SPEC.md for creating production templates.")


if __name__ == "__main__":
    main()
