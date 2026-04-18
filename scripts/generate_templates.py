"""
Template validation and documentation script.

IMPORTANT TECHNICAL LIMITATION:
python-pptx cannot programmatically create custom slide layouts or modify slide masters.
This fundamental limitation means templates MUST be created manually in PowerPoint UI.

This script provides:
1. Documentation of the manual creation process
2. Validation that manually-created templates meet requirements
3. Instructions for executing Plan B (manual creation)

Design Document Reference: Section 8.3.4 - Execute Plan B

Requirements:
- FR-CG-080: Template generation script (Plan B: manual creation instructions)
- FR-CG-081: data-template.pptx (manual creation)
- FR-CG-082: smartart-template.pptx (manual creation)
- FR-CG-083: 4 SmartArt layouts (manual creation)
- FR-CG-084: template_parser compatibility (validation)
- FR-CG-085: Idempotent execution (N/A for manual creation)
"""

import logging
from pathlib import Path

from pptx import Presentation

logger = logging.getLogger(__name__)

# Template output directory
TEMPLATES_DIR = Path(__file__).parent.parent / "templates"


def generate_data_template(output_path: str) -> None:
    """Provide instructions for manual data-template.pptx creation.

    PLAN B EXECUTION:
    python-pptx cannot create custom slide layouts programmatically.
    This function provides instructions for MANUAL template creation.

    Manual Creation Steps (see DATA-TEMPLATE-SPEC.md):
    1. Open Microsoft PowerPoint
    2. Create a new blank presentation
    3. Go to View > Slide Master
    4. Create 4 custom layouts:
       - "Chart" layout: TITLE + OBJECT placeholder
       - "Table" layout: TITLE + OBJECT placeholder
       - "Data Analysis" layout: TITLE + OBJECT x2 placeholders
       - "Two Column Data" layout: TITLE + OBJECT x2 placeholders (left-right)
    5. Exit Slide Master view
    6. Save as templates/data-template.pptx

    Args:
        output_path: Path where the template should be manually saved

    Requirements:
        - FR-CG-081: data-template.pptx (Plan B: manual creation)
    """
    output_file = Path(output_path)

    if output_file.exists():
        logger.info("data-template.pptx already exists at: %s", output_path)
        logger.info("Validating existing template...")

        try:
            prs = Presentation(str(output_file))
            layout_names = [layout.name for layout in prs.slide_master.slide_layouts]

            required_layouts = ["Chart", "Table", "Data Analysis", "Two Column Data"]
            missing_layouts = [name for name in required_layouts if name not in layout_names]

            if missing_layouts:
                logger.warning(
                    "Template exists but missing required layouts: %s",
                    missing_layouts,
                )
                logger.warning("Please update template to include all required layouts.")
            else:
                logger.info("✓ All required layouts found: %s", required_layouts)
                logger.info("✓ data-template.pptx is valid")

        except Exception:
            logger.exception("Failed to validate template")
            logger.warning("Please recreate template following DATA-TEMPLATE-SPEC.md")
    else:
        logger.warning("data-template.pptx does NOT exist at: %s", output_path)
        logger.warning("")
        logger.warning("MANUAL ACTION REQUIRED:")
        logger.warning("1. Open Microsoft PowerPoint")
        logger.warning("2. Follow instructions in templates/DATA-TEMPLATE-SPEC.md")
        logger.warning("3. Save as: %s", output_path)
        logger.warning("")
        logger.warning("Required layouts:")
        logger.warning("  - Chart (TITLE + OBJECT)")
        logger.warning("  - Table (TITLE + OBJECT)")
        logger.warning("  - Data Analysis (TITLE + OBJECT x2)")
        logger.warning("  - Two Column Data (TITLE + OBJECT x2 left-right)")


def generate_smartart_template(output_path: str) -> None:
    """Provide instructions for manual smartart-template.pptx creation.

    PLAN B EXECUTION:
    SmartArt templates require both custom layouts AND SmartArt diagram insertion,
    neither of which can be done programmatically with python-pptx.

    Manual Creation Steps (see SMARTART-TEMPLATE-SPEC.md):
    1. Open Microsoft PowerPoint
    2. Create a new blank presentation
    3. Go to View > Slide Master
    4. Create 4 custom layouts with SmartArt:
       - "Process Flow" layout: TITLE + SmartArt (Basic Process type)
       - "Hierarchy" layout: TITLE + SmartArt (Organization Chart type)
       - "Cycle" layout: TITLE + SmartArt (Basic Cycle type)
       - "Relationship" layout: TITLE + SmartArt (Basic Venn type)
    5. For each layout: Insert > SmartArt, choose appropriate type
    6. Configure with placeholder text (e.g., "Node 1", "Node 2", etc.)
    7. Exit Slide Master view
    8. Save as templates/smartart-template.pptx

    Args:
        output_path: Path where the template should be manually saved

    Requirements:
        - FR-CG-082: smartart-template.pptx (Plan B: manual creation)
        - FR-CG-083: 4 SmartArt layouts (Plan B: manual creation)
    """
    output_file = Path(output_path)

    if output_file.exists():
        logger.info("smartart-template.pptx already exists at: %s", output_path)
        logger.info("Validating existing template...")

        try:
            prs = Presentation(str(output_file))
            layout_names = [layout.name for layout in prs.slide_master.slide_layouts]

            required_layouts = ["Process Flow", "Hierarchy", "Cycle", "Relationship"]
            missing_layouts = [name for name in required_layouts if name not in layout_names]

            if missing_layouts:
                logger.warning(
                    "Template exists but missing required layouts: %s",
                    missing_layouts,
                )
                logger.warning("Please update template to include all required layouts.")
            else:
                logger.info("✓ All required layouts found: %s", required_layouts)
                logger.info("✓ smartart-template.pptx is valid")

        except Exception:
            logger.exception("Failed to validate template")
            logger.warning("Please recreate template following SMARTART-TEMPLATE-SPEC.md")
    else:
        logger.warning("smartart-template.pptx does NOT exist at: %s", output_path)
        logger.warning("")
        logger.warning("MANUAL ACTION REQUIRED:")
        logger.warning("1. Open Microsoft PowerPoint")
        logger.warning("2. Follow instructions in templates/SMARTART-TEMPLATE-SPEC.md")
        logger.warning("3. Save as: %s", output_path)
        logger.warning("")
        logger.warning("Required layouts with SmartArt:")
        logger.warning("  - Process Flow (Process type SmartArt)")
        logger.warning("  - Hierarchy (Organization Chart type)")
        logger.warning("  - Cycle (Basic Cycle type)")
        logger.warning("  - Relationship (Basic Venn type)")


def main() -> None:
    """CLI entry point for template validation and manual creation instructions.

    This script validates existing templates or provides instructions for manual creation.

    PLAN B EXECUTION RATIONALE:
    python-pptx cannot programmatically create custom slide layouts or SmartArt diagrams.
    This is a fundamental library limitation, not a code issue.

    See design.md Section 8.3.2-8.3.4 for timebox decision criteria and Plan B details.

    Requirements:
        - FR-CG-080: Template generation script (Plan B: manual instructions)
    """
    logger.info("=" * 70)
    logger.info("Template Validation and Manual Creation Instructions")
    logger.info("=" * 70)
    logger.info("")
    logger.info("IMPORTANT: python-pptx cannot create custom slide layouts.")
    logger.info("Templates must be created manually in PowerPoint UI (Plan B).")
    logger.info("")
    logger.info("Design Document Reference: Section 8.3.4 - Execute Plan B")
    logger.info("")
    logger.info("=" * 70)
    logger.info("")

    # Check/validate data template
    logger.info("[1/2] Checking data-template.pptx...")
    logger.info("")
    data_template_path = TEMPLATES_DIR / "data-template.pptx"
    generate_data_template(str(data_template_path))
    logger.info("")

    # Check/validate smartart template
    logger.info("[2/2] Checking smartart-template.pptx...")
    logger.info("")
    smartart_template_path = TEMPLATES_DIR / "smartart-template.pptx"
    generate_smartart_template(str(smartart_template_path))
    logger.info("")

    logger.info("=" * 70)
    logger.info("Summary")
    logger.info("=" * 70)
    logger.info("Data template: %s", data_template_path)
    logger.info("SmartArt template: %s", smartart_template_path)
    logger.info("")
    logger.info("If templates do not exist, follow the manual creation instructions above.")
    logger.info("See templates/DATA-TEMPLATE-SPEC.md and SMARTART-TEMPLATE-SPEC.md for details.")


if __name__ == "__main__":
    # Configure logging for CLI execution
    logging.basicConfig(
        level=logging.INFO,
        format="%(message)s",
    )
    main()
