"""Slide-level wrapper for python-pptx."""

import logging
from typing import Any

from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Length

from pptx_agent.constants import EMU_PER_CHAR_WIDTH, EMU_PER_LINE_HEIGHT
from pptx_agent.pptx_wrapper.placeholder_ops import find_placeholder
from pptx_agent.pptx_wrapper.type_helpers import set_text_frame_text

logger = logging.getLogger(__name__)


class SlideWrapper:
    """Wrapper for python-pptx Slide with type hints."""

    def __init__(self, slide: Slide) -> None:
        """Initialize slide wrapper.

        Args:
            slide: python-pptx Slide object to wrap
        """
        self._slide = slide

    @property
    def slide(self) -> Slide:
        """Return the underlying python-pptx Slide object.

        Returns:
            The wrapped Slide object
        """
        return self._slide

    def set_title(self, text: str) -> None:
        """Update title placeholder.

        Args:
            text: Title text to set

        Raises:
            ValueError: If slide has no title placeholder
        """
        if not self._slide.shapes.title:
            msg = "Slide has no title placeholder"
            raise ValueError(msg)

        self._slide.shapes.title.text = text

    def add_text(
        self, placeholder_name: str, text: str, check_overflow: bool = True
    ) -> dict[str, Any]:
        """Add text to placeholder with optional overflow check.

        Args:
            placeholder_name: Name of the placeholder to update
            text: Text content to add
            check_overflow: Whether to check for text overflow

        Returns:
            Dictionary with keys:
                - success (bool): Whether operation succeeded
                - overflow (bool): Whether text overflow detected
                - warnings (list[str]): Any warnings generated
        """
        result: dict[str, Any] = {"success": False, "overflow": False, "warnings": []}

        # Find placeholder using shared utility function with fuzzy matching
        placeholder = find_placeholder(self._slide, placeholder_name, fuzzy_match=True)

        if placeholder is None:
            result["warnings"].append(f"Placeholder '{placeholder_name}' not found")
            return result

        # Add text to placeholder
        if hasattr(placeholder, "text_frame"):
            # Type narrowing: hasattr check ensures text_frame exists
            set_text_frame_text(placeholder, text)
            result["success"] = True

            # Check for overflow if requested
            if check_overflow:
                result["overflow"] = self._check_text_overflow(placeholder, text)
                if result["overflow"]:
                    result["warnings"].append(f"Text may overflow placeholder '{placeholder_name}'")
        else:
            result["warnings"].append(f"Placeholder '{placeholder_name}' cannot contain text")

        return result

    def get_placeholders(self) -> dict[str, dict[str, Any]]:
        """Return placeholder inventory.

        Returns:
            Dictionary mapping placeholder names to their properties:
                - type: Placeholder type
                - has_text_frame: Whether it can contain text
        """
        placeholders = {}

        for shape in self._slide.shapes:
            if shape.is_placeholder:
                name = shape.name
                placeholders[name] = {
                    "type": shape.placeholder_format.type.name,
                    "has_text_frame": hasattr(shape, "text_frame"),
                }

        return placeholders

    def _check_text_overflow(self, placeholder: Any, text: str) -> bool:
        """Internal method to detect text overflow.

        This is a simplified heuristic based on text length and
        placeholder dimensions. For production use, more sophisticated
        overflow detection would be needed.

        Args:
            placeholder: The placeholder shape
            text: Text content

        Returns:
            True if overflow is likely, False otherwise
        """
        if not hasattr(placeholder, "text_frame"):
            return False

        # Simple heuristic: estimate based on character count and dimensions
        try:
            width = placeholder.width
            height = placeholder.height

            # Rough estimate of capacity (chars per line * number of lines)
            chars_per_line = max(1, int(width / EMU_PER_CHAR_WIDTH))
            num_lines = max(1, int(height / EMU_PER_LINE_HEIGHT))
            capacity = chars_per_line * num_lines

            # Check if text exceeds estimated capacity
            return len(text) > capacity
        except (AttributeError, ZeroDivisionError):
            # If we can't estimate, assume no overflow
            return False

    def add_chart(  # noqa: PLR0913
        self,
        chart_type: XL_CHART_TYPE,
        x: Length,
        y: Length,
        cx: Length,
        cy: Length,
        chart_data: ChartData | CategoryChartData,
    ) -> Any:
        """Add a chart to the slide.

        Args:
            chart_type: Type of chart to create (from XL_CHART_TYPE enum)
            x: X coordinate (Length object, e.g., from Inches())
            y: Y coordinate (Length object, e.g., from Inches())
            cx: Width (Length object, e.g., from Inches())
            cy: Height (Length object, e.g., from Inches())
            chart_data: Chart data object (ChartData or CategoryChartData)

        Returns:
            Chart shape object
        """
        return self._slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)  # type: ignore[arg-type]

    def add_table(  # noqa: PLR0913
        self,
        rows: int,
        cols: int,
        x: Length,
        y: Length,
        cx: Length,
        cy: Length,
    ) -> Any:
        """Add a table to the slide.

        Args:
            rows: Number of rows
            cols: Number of columns
            x: X coordinate (Length object, e.g., from Inches())
            y: Y coordinate (Length object, e.g., from Inches())
            cx: Width (Length object, e.g., from Inches())
            cy: Height (Length object, e.g., from Inches())

        Returns:
            Table shape object
        """
        return self._slide.shapes.add_table(rows, cols, x, y, cx, cy)
