"""Helper functions to localize type: ignore comments for python-pptx operations.

These helpers encapsulate operations that require type: ignore comments due to
incomplete type stubs in the python-pptx library. By centralizing these operations,
we improve type checking reliability across the codebase.
"""

from typing import Any

from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide


def add_chart_with_bounds(  # noqa: PLR0913
    slide: Slide,
    chart_type: XL_CHART_TYPE,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    chart_data: Any,
) -> Any:
    """Add chart to slide with position bounds.

    This helper localizes type: ignore comments needed due to python-pptx's
    incomplete type stubs for position parameters and ChartData subclasses.

    Args:
        slide: Slide to add chart to
        chart_type: Type of chart to create
        left: Left position (typically from get_placeholder_bounds or Inches)
        top: Top position (typically from get_placeholder_bounds or Inches)
        width: Width (typically from get_placeholder_bounds or Inches)
        height: Height (typically from get_placeholder_bounds or Inches)
        chart_data: Chart data (CategoryChartData or other ChartData subclass)

    Returns:
        The created chart shape
    """
    # CategoryChartData is a subclass of ChartData, but type checker doesn't
    # recognize inheritance properly. Position parameters also have type issues.
    return slide.shapes.add_chart(
        chart_type,
        left,  # type: ignore[arg-type]
        top,  # type: ignore[arg-type]
        width,  # type: ignore[arg-type]
        height,  # type: ignore[arg-type]
        chart_data,  # type: ignore[arg-type]
    )


def add_table_with_bounds(  # noqa: PLR0913
    slide: Slide,
    rows: int,
    cols: int,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
) -> Any:
    """Add table to slide with position bounds.

    This helper localizes type: ignore comments needed due to python-pptx's
    incomplete type stubs for position parameters.

    Args:
        slide: Slide to add table to
        rows: Number of rows
        cols: Number of columns
        left: Left position (typically from get_placeholder_bounds or Inches)
        top: Top position (typically from get_placeholder_bounds or Inches)
        width: Width (typically from get_placeholder_bounds or Inches)
        height: Height (typically from get_placeholder_bounds or Inches)

    Returns:
        The created table object (accessed via .table property)
    """
    # Position parameters have type issues with python-pptx stubs
    return slide.shapes.add_table(
        rows,
        cols,
        left,  # type: ignore[arg-type]
        top,  # type: ignore[arg-type]
        width,  # type: ignore[arg-type]
        height,  # type: ignore[arg-type]
    ).table


def add_picture_with_bounds(  # noqa: PLR0913
    slide: Slide,
    image_path: str,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
) -> Any:
    """Add picture to slide with position bounds.

    This helper localizes type: ignore comments needed due to python-pptx's
    incomplete type stubs for position parameters.

    Args:
        slide: Slide to add picture to
        image_path: Path to image file
        left: Left position (typically from get_placeholder_bounds or Inches)
        top: Top position (typically from get_placeholder_bounds or Inches)
        width: Width (typically from get_placeholder_bounds or Inches)
        height: Height (typically from get_placeholder_bounds or Inches)

    Returns:
        The created picture shape
    """
    # Position parameters have type issues with python-pptx stubs
    return slide.shapes.add_picture(
        image_path,
        left,  # type: ignore[arg-type]
        top,  # type: ignore[arg-type]
        width=width,  # type: ignore[arg-type]
        height=height,  # type: ignore[arg-type]
    )


def set_text_frame_text(placeholder: Any, text: str) -> None:
    """Set text on a placeholder's text frame.

    This helper localizes type: ignore comments needed because hasattr() checks
    don't properly narrow types for the type checker.

    Args:
        placeholder: Placeholder shape (must have text_frame attribute)
        text: Text content to set
    """
    # hasattr check in calling code doesn't narrow type for type checker
    placeholder.text_frame.text = text  # type: ignore[attr-defined]
