"""Placeholder operations for shape wrappers.

Common placeholder search, extraction, and removal operations
shared across ChartWrapper, TableWrapper, and ImageWrapper.
"""

import logging
from typing import Any

from pptx.slide import Slide

from pptx_agent.pptx_wrapper.xml_utils import safe_get_element, safe_remove_element

logger = logging.getLogger(__name__)


def find_placeholder(slide: Slide, placeholder_name: str, fuzzy_match: bool = False) -> Any | None:
    """Find placeholder by name or type with optional fuzzy matching.

    Args:
        slide: Slide to search in
        placeholder_name: Name or type name of placeholder to find
        fuzzy_match: If True, use case-insensitive partial matching and fallback logic

    Returns:
        Placeholder shape if found, None otherwise
    """
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue

        # Try exact match first on name
        if shape.name == placeholder_name:
            return shape

        # Try placeholder type name with exception handling
        try:
            if shape.placeholder_format.type.name == placeholder_name:
                return shape
        except AttributeError as e:
            # If accessing placeholder format fails, log and continue searching
            logger.debug("Could not access placeholder_format for shape %s: %s", shape.name, e)

        # If fuzzy match enabled, try case-insensitive partial match on name
        if fuzzy_match and placeholder_name.lower() in shape.name.lower():
            return shape

    # If still not found and fuzzy match enabled with generic name, try fallback
    if fuzzy_match and placeholder_name.lower() in ["content", "body", "text"]:
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, "text_frame"):
                # Skip title placeholder
                try:
                    if shape.placeholder_format.type.name != "TITLE":
                        return shape
                except AttributeError:
                    # If we can't determine type, use it anyway as fallback
                    logger.debug(
                        "Could not determine placeholder type for shape, using as fallback"
                    )
                    return shape

    return None


def get_placeholder_bounds(placeholder: Any) -> tuple[int, int, int, int]:
    """Extract position and dimensions from placeholder.

    Args:
        placeholder: Placeholder shape

    Returns:
        Tuple of (left, top, width, height)
    """
    return placeholder.left, placeholder.top, placeholder.width, placeholder.height


def remove_placeholder_safely(placeholder: Any) -> None:
    """Safely remove placeholder using xml_utils.

    Args:
        placeholder: Placeholder shape to remove

    Raises:
        ValueError: If placeholder removal fails
    """
    sp = safe_get_element(placeholder, "element")
    if sp is None or not safe_remove_element(sp):
        msg = (
            "Placeholder removal failed: unable to remove placeholder element. "
            "Expected: valid removable placeholder."
        )
        raise ValueError(msg)
