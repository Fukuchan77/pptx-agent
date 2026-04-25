"""Presentation-level wrapper for python-pptx."""

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation as PresentationFactory

if TYPE_CHECKING:
    from pptx.presentation import Presentation

from pptx_agent.validators.file_validator import validate_template_path

from .slide import SlideWrapper


class PresentationWrapper:
    """Wrapper for python-pptx Presentation with type hints."""

    def __init__(self) -> None:
        """Initialize empty presentation wrapper."""
        self._prs: Presentation | None = None
        self._template_path: str | None = None

    def load_template(self, path: str) -> None:
        """Load PPTX template file.

        Args:
            path: Path to PPTX template file

        Raises:
            InvalidFileError: If file doesn't exist, has wrong extension, or is not valid PPTX
            SecurityValidationError: If path contains symlinks
            FileSizeLimitError: If file size exceeds limits
            CompressionRatioError: If compression ratio is suspicious
        """
        # Validate template path (includes existence, extension, symlinks, ZIP structure)
        validated_path = validate_template_path(path)

        self._prs = PresentationFactory(str(validated_path))
        self._template_path = str(validated_path)

    @property
    def is_loaded(self) -> bool:
        """Return True if a template has been loaded.

        Returns:
            True if a template is loaded, False otherwise
        """
        return self._prs is not None

    @property
    def prs(self) -> "Presentation":
        """Return the underlying python-pptx Presentation object.

        Returns:
            The underlying Presentation object

        Raises:
            ValueError: If template not loaded
        """
        if self._prs is None:
            msg = "Template must be loaded"
            raise ValueError(msg)
        return self._prs

    @property
    def core_properties(self) -> Any:
        """Return core properties of the presentation.

        Returns:
            Core properties object from the underlying presentation

        Raises:
            ValueError: If template not loaded
        """
        if self._prs is None:
            msg = "Template must be loaded"
            raise ValueError(msg)
        return self._prs.core_properties

    def add_slide(self, layout_name: str) -> SlideWrapper:
        """Add slide with specified layout.

        Args:
            layout_name: Name of the slide layout to use

        Returns:
            SlideWrapper instance for the new slide

        Raises:
            ValueError: If template not loaded or layout not found
        """
        if self._prs is None:
            msg = "Template must be loaded before adding slides"
            raise ValueError(msg)

        # Find layout by name
        layout = None
        for slide_layout in self._prs.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break

        if layout is None:
            msg = f"Layout '{layout_name}' not found in template"
            raise ValueError(msg)

        # Add slide with layout
        slide = self._prs.slides.add_slide(layout)
        return SlideWrapper(slide)

    def insert_slide(self, layout_name: str, index: int) -> SlideWrapper:
        """Insert slide at specified position with specified layout.

        Args:
            layout_name: Name of the slide layout to use
            index: Position to insert the slide (0-based)

        Returns:
            SlideWrapper instance for the new slide

        Raises:
            ValueError: If template not loaded or layout not found
            IndexError: If index is out of range
        """
        if self._prs is None:
            msg = "Template must be loaded before inserting slides"
            raise ValueError(msg)

        # Validate index
        slide_count = len(self._prs.slides)
        if index < 0 or index > slide_count:
            msg = f"Slide index {index} out of range (valid range: 0 to {slide_count})"
            raise IndexError(msg)

        # Find layout by name
        layout = None
        for slide_layout in self._prs.slide_layouts:
            if slide_layout.name == layout_name:
                layout = slide_layout
                break

        if layout is None:
            msg = f"Layout '{layout_name}' not found in template"
            raise ValueError(msg)

        # Add slide at end first
        slide = self._prs.slides.add_slide(layout)

        # If not at end, move to correct position
        if index < slide_count:
            # Access slide ID list from XML
            slide_id_list = self._prs.slides._sldIdLst  # noqa: SLF001  # type: ignore[reportPrivateUsage]
            # Get reference to the last element (newly added slide)
            last_slide_id = slide_id_list[-1]
            # Remove it from the list
            slide_id_list.remove(last_slide_id)
            # Insert at the correct position
            slide_id_list.insert(index, last_slide_id)

        return SlideWrapper(slide)

    def get_layouts(self) -> list[str]:
        """Return list of available layout names.

        Returns:
            List of layout names from the template

        Raises:
            ValueError: If template not loaded
        """
        if self._prs is None:
            msg = "Template must be loaded before getting layouts"
            raise ValueError(msg)

        return [layout.name for layout in self._prs.slide_layouts]

    def delete_slide(self, slide_index: int) -> None:
        """Delete slide at specified index.

        Args:
            slide_index: Index of slide to delete (0-based)

        Raises:
            ValueError: If template not loaded
            IndexError: If slide_index is out of range
        """
        if self._prs is None:
            msg = "Template must be loaded before deleting slides"
            raise ValueError(msg)

        slides = self._prs.slides
        if slide_index < 0 or slide_index >= len(slides):
            msg = f"Slide index {slide_index} out of range (presentation has {len(slides)} slides)"
            raise IndexError(msg)

        # Access slide ID list from XML (private member access required for slide deletion)
        slide_id_list = self._prs.slides._sldIdLst  # noqa: SLF001  # type: ignore[reportPrivateUsage]
        # Remove slide at index
        slide_id_list.remove(slide_id_list[slide_index])

    def slide_count(self) -> int:
        """Return the number of slides in the presentation.

        Returns:
            Number of slides

        Raises:
            ValueError: If template not loaded
        """
        if self._prs is None:
            msg = "Template must be loaded"
            raise ValueError(msg)

        return len(self._prs.slides)

    def save(self, output_path: str) -> None:
        """Save presentation to file.

        Args:
            output_path: Path where presentation should be saved

        Raises:
            ValueError: If template not loaded
        """
        if self._prs is None:
            msg = "Template must be loaded before saving"
            raise ValueError(msg)

        # Convert to Path and ensure parent directory exists
        output = Path(output_path).resolve()
        output.parent.mkdir(parents=True, exist_ok=True)

        # Save presentation
        self._prs.save(str(output))
