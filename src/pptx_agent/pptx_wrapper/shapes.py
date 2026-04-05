"""Shape wrappers for charts, tables, and images."""

import logging
from pathlib import Path
from typing import Any, ClassVar

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Inches

from pptx_agent.constants import (
    DEFAULT_SHAPE_HEIGHT_INCHES,
    DEFAULT_SHAPE_LEFT_INCHES,
    DEFAULT_SHAPE_TOP_INCHES,
    DEFAULT_SHAPE_WIDTH_INCHES,
)
from pptx_agent.pptx_wrapper.placeholder_ops import (
    find_placeholder,
    get_placeholder_bounds,
    remove_placeholder_safely,
)
from pptx_agent.pptx_wrapper.type_helpers import (
    add_chart_with_bounds,
    add_picture_with_bounds,
    add_table_with_bounds,
)
from pptx_agent.pptx_wrapper.xml_utils import safe_get_element

logger = logging.getLogger(__name__)


class ChartWrapper:
    """Wrapper for chart creation."""

    # Mapping of chart type strings to python-pptx enum values
    CHART_TYPES: ClassVar[dict[str, XL_CHART_TYPE]] = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "area": XL_CHART_TYPE.AREA,
    }

    @staticmethod
    def create_chart(
        slide: Slide, placeholder_name: str, chart_type: str, data: dict[str, Any]
    ) -> None:
        """Create chart in placeholder.

        Args:
            slide: Slide to add chart to
            placeholder_name: Name of placeholder for chart
            chart_type: Type of chart (bar, line, pie, column, area)
             Chart data with structure:
                {
                    "categories": ["Cat1", "Cat2", ...],
                    "series": [
                        {"name": "Series1", "values": [1, 2, ...]},
                        ...
                    ]
                }

        Raises:
            ValueError: If chart type is unsupported or data is invalid
        """
        # Validate chart type
        if chart_type not in ChartWrapper.CHART_TYPES:
            msg = (
                f"Chart type validation failed: '{chart_type}' is not supported. "
                f"Expected: one of {', '.join(ChartWrapper.CHART_TYPES.keys())}."
            )
            raise ValueError(msg)

        # Validate data
        if not data.get("categories") or not data.get("series"):
            msg = (
                "Chart data validation failed: missing required fields. "
                "Expected: 'categories' and 'series' fields in data."
            )
            raise ValueError(msg)

        categories = data["categories"]
        series_list = data["series"]

        if len(categories) == 0 or len(series_list) == 0:
            msg = (
                "Chart data validation failed: data cannot be empty. "
                "Expected: non-empty categories and series."
            )
            raise ValueError(msg)

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for series in series_list:
            series_name = series.get("name", "Series")
            series_values = series.get("values", [])
            chart_data.add_series(series_name, series_values)

        # Find placeholder and add chart
        chart_placeholder = find_placeholder(slide, placeholder_name)

        if chart_placeholder is None:
            # If placeholder not found, add chart at default position
            x, y = Inches(DEFAULT_SHAPE_LEFT_INCHES), Inches(DEFAULT_SHAPE_TOP_INCHES)
            cx, cy = Inches(DEFAULT_SHAPE_WIDTH_INCHES), Inches(DEFAULT_SHAPE_HEIGHT_INCHES)
            add_chart_with_bounds(
                slide, ChartWrapper.CHART_TYPES[chart_type], x, y, cx, cy, chart_data
            )
        else:
            # Replace placeholder with chart
            left, top, width, height = get_placeholder_bounds(chart_placeholder)

            # Remove placeholder safely
            remove_placeholder_safely(chart_placeholder)

            # Add chart in same position
            add_chart_with_bounds(
                slide,
                ChartWrapper.CHART_TYPES[chart_type],
                left,
                top,
                width,
                height,
                chart_data,
            )


class TableWrapper:
    """Wrapper for table creation."""

    @staticmethod
    def create_table(
        slide: Slide, placeholder_name: str, rows: list[list[str]], headers: list[str] | None = None
    ) -> None:
        """Create table in placeholder.

        Args:
            slide: Slide to add table to
            placeholder_name: Name of placeholder for table
            rows: List of rows, each row is a list of cell values
            headers: Optional list of header values

        Raises:
            ValueError: If rows are empty or have inconsistent column counts
        """
        # Validate rows
        if not rows or len(rows) == 0:
            msg = (
                "Table data validation failed: rows cannot be empty. "
                "Expected: at least one row of data."
            )
            raise ValueError(msg)

        # Check column count consistency
        col_counts = [len(row) for row in rows]
        if len(set(col_counts)) > 1:
            msg = (
                f"Table data validation failed: inconsistent column counts detected {col_counts}. "
                f"Expected: all rows must have the same number of columns."
            )
            raise ValueError(msg)

        # Determine table dimensions
        num_cols = len(rows[0])
        num_rows = len(rows) + (1 if headers else 0)

        # Find placeholder or use default position
        table_placeholder = find_placeholder(slide, placeholder_name)

        if table_placeholder is None:
            # Add table at default position
            x, y = Inches(DEFAULT_SHAPE_LEFT_INCHES), Inches(DEFAULT_SHAPE_TOP_INCHES)
            cx, cy = Inches(DEFAULT_SHAPE_WIDTH_INCHES), Inches(DEFAULT_SHAPE_HEIGHT_INCHES)
            table = add_table_with_bounds(slide, num_rows, num_cols, x, y, cx, cy)
        else:
            # Replace placeholder with table
            left, top, width, height = get_placeholder_bounds(table_placeholder)

            # Remove placeholder safely
            remove_placeholder_safely(table_placeholder)

            # Add table in same position
            table = add_table_with_bounds(
                slide,
                num_rows,
                num_cols,
                left,
                top,
                width,
                height,
            )

        # Populate table
        row_idx = 0

        # Add headers if provided
        if headers:
            for col_idx, header in enumerate(headers):
                table.cell(row_idx, col_idx).text = str(header)
            row_idx += 1

        # Add data rows
        for row in rows:
            for col_idx, cell_value in enumerate(row):
                table.cell(row_idx, col_idx).text = str(cell_value)
            row_idx += 1


class ImageWrapper:
    """Wrapper for image insertion."""

    @staticmethod
    def add_image(slide: Slide, placeholder_name: str, image_path: str, alt_text: str = "") -> None:
        """Add image to placeholder.

        Args:
            slide: Slide to add image to
            placeholder_name: Name of placeholder for image
            image_path: Path to image file
            alt_text: Alternative text for accessibility

        Raises:
            FileNotFoundError: If image file doesn't exist
        """
        # Verify image file exists
        img_path = Path(image_path)
        if not img_path.exists():
            msg = (
                f"Image file validation failed: file not found at '{image_path}'. "
                "Expected: valid path to existing image file."
            )
            raise FileNotFoundError(msg)

        # Find placeholder
        image_placeholder = find_placeholder(slide, placeholder_name)

        if image_placeholder is None:
            # Add image at default position
            left = Inches(DEFAULT_SHAPE_LEFT_INCHES)
            top = Inches(DEFAULT_SHAPE_TOP_INCHES)
            pic = slide.shapes.add_picture(str(img_path), left, top)
        else:
            # Insert image in placeholder
            left, top, width, height = get_placeholder_bounds(image_placeholder)

            # Remove placeholder safely
            remove_placeholder_safely(image_placeholder)

            # Add image in same position, maintaining aspect ratio
            pic = add_picture_with_bounds(
                slide,
                str(img_path),
                left,
                top,
                width,
                height,
            )

        # Set alt text if provided
        if alt_text:
            # Access the underlying XML element safely to set alt text
            element = safe_get_element(pic, "_element")
            if element is not None:
                try:
                    element.get_or_add_nvPicPr().cNvPr.set("descr", alt_text)
                except AttributeError as e:
                    # If setting alt text fails due to missing attribute, log and continue
                    logger.warning("Failed to set alt text for image: %s", e)
