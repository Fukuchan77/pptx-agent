"""Chart builder for creating charts from ChartBlock specifications.

This module provides functionality to add charts to PowerPoint slides
using python-pptx chart capabilities.
"""

import logging
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx_agent.pptx_wrapper.slide import SlideWrapper
from pptx_agent.schemas.visual_assets import ChartBlock

logger = logging.getLogger(__name__)

# Mapping of chart type strings to python-pptx chart type enums
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "area": XL_CHART_TYPE.AREA,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "radar": XL_CHART_TYPE.RADAR,
}


def add_chart_to_slide(slide: SlideWrapper, chart_block: ChartBlock) -> None:
    """Add a chart to a slide based on ChartBlock specification.

    Args:
        slide: SlideWrapper instance to add chart to
        chart_block: ChartBlock containing chart specification

    Raises:
        ValueError: If chart type is unsupported or data is invalid
    """
    # Validate chart type
    if chart_block.chart_type not in CHART_TYPE_MAP:
        msg = f"Unsupported chart type: {chart_block.chart_type}"
        raise ValueError(msg)

    # Validate data is not empty
    data = chart_block.data
    if not data.get("categories") or not data.get("series"):
        msg = "Chart data cannot be empty"
        raise ValueError(msg)

    # Get chart type enum
    chart_type = CHART_TYPE_MAP[chart_block.chart_type]

    # Create chart data
    chart_data = _create_chart_data(data)

    # Default chart position and size (can be customized later)
    x = Inches(1)
    y = Inches(2)
    cx = Inches(8)
    cy = Inches(5)

    # Add chart to slide using public method
    chart_shape = slide.add_chart(chart_type, x, y, cx, cy, chart_data)

    # Set chart title if provided
    if chart_block.title:
        chart_shape.chart.has_title = True
        chart_shape.chart.chart_title.text_frame.text = chart_block.title

    logger.info(
        "Added %s chart to slide with title: %s",
        chart_block.chart_type,
        chart_block.title,
    )


def _create_chart_data(data: dict[str, Any]) -> CategoryChartData:
    """Create CategoryChartData from chart data dictionary.

    Args:
        data: Dictionary containing 'categories' and 'series' lists

    Returns:
        CategoryChartData object for chart creation
    """
    chart_data = CategoryChartData()

    # Set categories
    categories = data["categories"]
    chart_data.categories = categories

    # Add each series
    for series_info in data["series"]:
        series_name = series_info["name"]
        series_values = series_info["values"]
        chart_data.add_series(series_name, series_values)

    return chart_data
