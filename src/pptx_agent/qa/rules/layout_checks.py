"""QA rules for layout validation.

This module implements layout-related quality checks including:
- Text overflow detection (QA-L-001)
- Empty placeholder detection (QA-L-002)
- Unpopulated placeholder detection (QA-L-003)
- Overlapping object detection (QA-L-004)
- Boundary overflow detection (QA-L-005)
- Minimum font size detection (QA-L-006)
"""

from typing import Any, Literal, cast

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.util import Pt

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.schemas import QAIssue, Severity
from pptx_agent.utils.language_detector import detect_language
from pptx_agent.utils.text_capacity import calculate_effective_capacity

# Text overflow threshold in characters
_TEXT_OVERFLOW_CHAR_THRESHOLD = 500


class TextOverflowRule:
    """QA-L-001: Detect text frame overflow with language-aware capacity checks.

    Checks if text content exceeds the bounds of text frames, which can
    cause text to be cut off or display incorrectly in presentations.
    Uses language detection to apply appropriate capacity multipliers
    (0.55x for Japanese, 1.0x for English).

    Attributes:
        rule_id: Unique identifier "QA-L-001"
        description: Human-readable description
        auto_fixable: True - can be fixed via font reduction or summarization
        language: Optional language override ("ja" or "en")
    """

    rule_id = "QA-L-001"
    description = "Detect text overflow in text frames (language-aware)"
    severity = "error"
    auto_fixable = True

    def __init__(self, language: Literal["ja", "en"] | None = None) -> None:
        """Initialize rule with optional language override.

        Args:
            language: Optional language override ("ja" or "en").
                     If None, language is auto-detected per text frame.
        """
        self.language = language

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for text overflow issues.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found (empty if no overflow detected)
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Only check shapes with text frames
                if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                    continue

                # Check if text overflows using language-aware capacity
                text_frame = shape.text_frame  # type: ignore[attr-defined]
                text = text_frame.text

                # Detect language for this text frame (use override if provided)
                if self.language is not None:
                    detected_lang = self.language
                else:
                    detected_lang = detect_language(text)

                # Cast to help type checker (detect_language returns Literal["ja", "en"])
                lang = cast("Literal['ja', 'en']", detected_lang)

                if self._has_text_overflow(shape, lang):
                    # Calculate effective capacity for better error message
                    effective_capacity = calculate_effective_capacity(
                        _TEXT_OVERFLOW_CHAR_THRESHOLD, lang
                    )

                    issues.append(
                        QAIssue(
                            rule_id=self.rule_id,
                            severity=Severity.ERROR,
                            slide_index=slide_idx,
                            shape_index=shape_idx,
                            message=(
                                f"Text overflow detected in text frame "
                                f"({len(text)} characters, "
                                f"effective capacity: {effective_capacity} "
                                f"for {lang.upper()} text)"
                            ),
                            auto_fixable=True,
                            suggested_fix="Reduce font size, switch layout, or summarize content",
                        )
                    )

        return issues

    def _has_text_overflow(self, shape: Any, language: Literal["ja", "en"]) -> bool:
        """Check if shape has text overflow using language-aware capacity.

        Args:
            shape: Shape to check
            language: Language of the text ("ja" or "en")

        Returns:
            True if text overflows, False otherwise
        """
        try:
            text_frame = shape.text_frame
            text = text_frame.text

            # Calculate effective capacity based on language
            effective_capacity = calculate_effective_capacity(
                _TEXT_OVERFLOW_CHAR_THRESHOLD, language
            )

            # Check if text exceeds language-aware capacity
            if len(text) > effective_capacity:
                return True

            # Check if shape has overflow method (custom implementation)
            if hasattr(shape, "has_overflow") and callable(shape.has_overflow):
                return bool(shape.has_overflow())
        except Exception:
            # If we can't check, assume no overflow
            return False
        else:
            return False


class EmptyPlaceholderRule:
    """QA-L-002: Detect empty title placeholders.

    Checks if title placeholders are empty, which results in unprofessional
    presentations with missing slide titles.

    Attributes:
        rule_id: Unique identifier "QA-L-002"
        description: Human-readable description
        auto_fixable: True - can be populated from outline headers
    """

    rule_id = "QA-L-002"
    description = "Detect empty title placeholders"
    severity = "error"
    auto_fixable = True

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for empty title placeholders.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found (empty if no empty titles)
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            # Check if slide has title placeholder
            if not hasattr(slide.shapes, "title") or slide.shapes.title is None:
                continue

            title_shape = slide.shapes.title

            # Check if title is empty
            if self._is_empty_title(title_shape):
                issues.append(
                    QAIssue(
                        rule_id=self.rule_id,
                        severity=Severity.ERROR,
                        slide_index=slide_idx,
                        shape_index=None,  # Title is slide-level
                        message="Empty title placeholder detected",
                        auto_fixable=True,
                        suggested_fix="Populate title from outline or slide content",
                    )
                )

        return issues

    def _is_empty_title(self, title_shape: Any) -> bool:
        """Check if title shape is empty.

        Args:
            title_shape: Title shape to check

        Returns:
            True if title is empty, False otherwise
        """
        try:
            if not hasattr(title_shape, "text_frame"):
                return False

            text = title_shape.text_frame.text.strip()
            return len(text) == 0
        except Exception:
            return False


class UnpopulatedPlaceholderRule:
    """QA-L-003: Detect unpopulated required placeholders.

    Checks if required placeholders (beyond title) are left empty,
    which can indicate incomplete content generation.

    Attributes:
        rule_id: Unique identifier "QA-L-003"
        description: Human-readable description
        auto_fixable: False - requires content decision
    """

    rule_id = "QA-L-003"
    description = "Detect unpopulated required placeholders"
    severity = "warning"
    auto_fixable = False

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for unpopulated placeholders.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape is a placeholder
                if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
                    continue

                # Skip title placeholders (handled by EmptyPlaceholderRule)
                if (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                ):
                    continue

                # Check if placeholder is empty
                if self._is_empty_placeholder(shape):
                    placeholder_type = (
                        shape.placeholder_format.type
                        if hasattr(shape, "placeholder_format")
                        else "unknown"
                    )
                    issues.append(
                        QAIssue(
                            rule_id=self.rule_id,
                            severity=Severity.WARNING,
                            slide_index=slide_idx,
                            shape_index=shape_idx,
                            message=f"Unpopulated placeholder detected (type: {placeholder_type})",
                            auto_fixable=False,
                            suggested_fix="Populate placeholder with appropriate content",
                        )
                    )

        return issues

    def _is_empty_placeholder(self, shape: Any) -> bool:
        """Check if placeholder is empty.

        Args:
            shape: Placeholder shape to check

        Returns:
            True if placeholder is empty, False otherwise
        """
        try:
            # Check text placeholders
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                return len(text) == 0

            # For non-text placeholders (images, charts, etc.),
            # we'd need more sophisticated checks
        except Exception:
            return False
        else:
            return False


class OverlappingObjectsRule:
    """QA-L-004: Detect overlapping objects.

    Checks if shapes overlap via bounding box intersection, which can
    indicate layout issues or unintended visual conflicts.

    Attributes:
        rule_id: Unique identifier "QA-L-004"
        description: Human-readable description
        auto_fixable: False - requires manual layout adjustment
    """

    rule_id = "QA-L-004"
    description = "Detect overlapping objects via bounding box intersection"
    severity = "warning"
    auto_fixable = False

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for overlapping objects.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            shapes = list(slide.shapes)

            # Check each pair of shapes for overlap
            for i, shape_i in enumerate(shapes):
                for j in range(i + 1, len(shapes)):
                    shape_j = shapes[j]
                    if self._shapes_overlap(shape_i, shape_j):
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.WARNING,
                                slide_index=slide_idx,
                                shape_index=i,  # Report first shape
                                message=f"Shapes {i} and {j} overlap",
                                auto_fixable=False,
                                suggested_fix="Adjust shape positions to eliminate overlap",
                            )
                        )

        return issues

    def _shapes_overlap(self, shape1: Any, shape2: Any) -> bool:
        """Check if two shapes overlap.

        Args:
            shape1: First shape
            shape2: Second shape

        Returns:
            True if shapes overlap, False otherwise
        """
        try:
            # Get bounding boxes
            left1 = shape1.left
            top1 = shape1.top
            right1 = left1 + shape1.width
            bottom1 = top1 + shape1.height

            left2 = shape2.left
            top2 = shape2.top
            right2 = left2 + shape2.width
            bottom2 = top2 + shape2.height

            # Check for overlap using rectangle intersection
            # Rectangles overlap if they intersect in both x and y dimensions
            x_overlap = left1 < right2 and right1 > left2
            y_overlap = top1 < bottom2 and bottom1 > top2
        except Exception:
            return False
        else:
            return x_overlap and y_overlap


class BoundaryOverflowRule:
    """QA-L-005: Detect objects extending beyond slide boundary.

    Checks if shapes extend beyond the slide boundaries, which can
    cause content to be cut off when presenting.

    Attributes:
        rule_id: Unique identifier "QA-L-005"
        description: Human-readable description
        auto_fixable: True - can clip coordinates to slide bounds
    """

    rule_id = "QA-L-005"
    description = "Detect objects extending beyond slide boundary"
    severity = "warning"
    auto_fixable = True

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for boundary overflow.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        # Get slide dimensions
        slide_width = presentation.prs.slide_width
        slide_height = presentation.prs.slide_height

        if slide_width is None or slide_height is None:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if self._extends_beyond_boundary(shape, int(slide_width), int(slide_height)):
                    issues.append(
                        QAIssue(
                            rule_id=self.rule_id,
                            severity=Severity.ERROR,
                            slide_index=slide_idx,
                            shape_index=shape_idx,
                            message="Shape extends beyond slide boundary",
                            auto_fixable=True,
                            suggested_fix="Clip shape coordinates to slide bounds",
                        )
                    )

        return issues

    def _extends_beyond_boundary(self, shape: Any, slide_width: int, slide_height: int) -> bool:
        """Check if shape extends beyond slide boundary.

        Args:
            shape: Shape to check
            slide_width: Slide width in EMUs
            slide_height: Slide height in EMUs

        Returns:
            True if shape extends beyond boundary, False otherwise
        """
        try:
            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height

            # Check if any edge extends beyond slide
        except Exception:
            return False
        else:
            return left < 0 or top < 0 or right > slide_width or bottom > slide_height


class MinimumFontSizeRule:
    """QA-L-006: Detect font sizes below minimum threshold.

    Checks if text uses font sizes below the minimum readable threshold
    (default 10pt), which can affect readability.

    Attributes:
        rule_id: Unique identifier "QA-L-006"
        description: Human-readable description
        auto_fixable: False - requires content or layout adjustment
    """

    rule_id = "QA-L-006"
    description = "Detect font sizes below minimum threshold"
    severity = "warning"
    auto_fixable = False

    def __init__(self, min_font_size: float = 10.0) -> None:
        """Initialize rule with minimum font size threshold.

        Args:
            min_font_size: Minimum font size in points (default 10.0)
        """
        self.min_font_size = min_font_size

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Validate presentation for minimum font size violations.

        Args:
            presentation: Presentation to validate

        Returns:
            List of QA issues found
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                    continue

                # Check font sizes in text frame
                if self._has_small_font(shape.text_frame):  # type: ignore[attr-defined]
                    issues.append(
                        QAIssue(
                            rule_id=self.rule_id,
                            severity=Severity.WARNING,
                            slide_index=slide_idx,
                            shape_index=shape_idx,
                            message=f"Font size below minimum threshold ({self.min_font_size}pt)",
                            auto_fixable=False,
                            suggested_fix=f"Increase font size to at least {self.min_font_size}pt",
                        )
                    )

        return issues

    def _has_small_font(self, text_frame: Any) -> bool:
        """Check if text frame has font sizes below threshold.

        Args:
            text_frame: Text frame to check

        Returns:
            True if any font is below threshold, False otherwise
        """
        try:
            min_size_emu = Pt(self.min_font_size)

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None and run.font.size < min_size_emu:
                        return True
        except Exception:
            return False
        else:
            return False


# Made with Bob
