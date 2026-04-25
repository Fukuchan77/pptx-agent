"""Content quality check rules for QA engine.

This module implements QA rules that validate content quality aspects:
- Bullet point length (QA-C-001)
- Duplicate slide titles (QA-C-002)
- Unpopulated image placeholders (QA-C-003)
- Pathological table dimensions (QA-C-004)
- Missing chart data (QA-C-005)
- Speaker notes verification (QA-C-006)
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE

from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.schemas import QAIssue, Severity


class BulletLengthRule:
    """QA-C-001: Detect bullet points exceeding configured length threshold.

    Long bullet points reduce readability and may indicate content that should
    be split across multiple slides or summarized.

    Attributes:
        rule_id: Unique identifier "QA-C-001"
        description: Human-readable description
        auto_fixable: False - requires manual content review
        max_length: Maximum characters per bullet point
    """

    rule_id = "QA-C-001"
    description = "Detect overly long bullet points"
    severity = "warning"
    auto_fixable = False

    def __init__(self, max_length: int = 120) -> None:
        """Initialize bullet length rule.

        Args:
            max_length: Maximum characters per bullet point (default: 120)
        """
        self.max_length = max_length

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for overly long bullet points.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for bullets exceeding max_length
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame  # type: ignore[attr-defined]
                for paragraph in text_frame.paragraphs:
                    # Check if paragraph is a bullet point (has level > 0 or bullet)
                    if paragraph.level > 0 or (
                        paragraph.level == 0 and len(text_frame.paragraphs) > 1
                    ):
                        text = paragraph.text.strip()
                        if len(text) > self.max_length:
                            issues.append(
                                QAIssue(
                                    rule_id=self.rule_id,
                                    severity=Severity.WARNING,
                                    slide_index=slide_idx,
                                    shape_index=shape_idx,
                                    message=(
                                        f"Bullet point exceeds {self.max_length} characters"
                                        f" ({len(text)} chars): '{text[:50]}...'"
                                    ),
                                    auto_fixable=False,
                                    suggested_fix=(
                                        "Split long bullet into multiple points"
                                        " or summarize content"
                                    ),
                                )
                            )

        return issues


class DuplicateTitleRule:
    """QA-C-002: Detect duplicate slide titles.

    Duplicate titles may indicate copy-paste errors or poor slide organization.

    Attributes:
        rule_id: Unique identifier "QA-C-002"
        description: Human-readable description
        auto_fixable: False - requires manual review
    """

    rule_id = "QA-C-002"
    description = "Detect duplicate slide titles"
    severity = "info"
    auto_fixable = False

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for duplicate slide titles.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for duplicate titles
        """
        issues: list[QAIssue] = []
        title_map: dict[str, list[int]] = {}

        if not presentation.is_loaded:
            return issues

        # Collect all titles and their slide indices
        for slide_idx, slide in enumerate(presentation.prs.slides):
            title_text: str | None = None

            # Try to find title placeholder
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    # Check for TITLE placeholder type (type 1)
                    if ph.type == PP_PLACEHOLDER_TYPE.TITLE and (
                        hasattr(shape, "text_frame") and shape.has_text_frame
                    ):
                        text_frame = shape.text_frame  # type: ignore[attr-defined]
                        title_text = text_frame.text.strip()
                        break

            if title_text:
                if title_text not in title_map:
                    title_map[title_text] = []
                title_map[title_text].append(slide_idx)

        # Report duplicates
        for title, slide_indices in title_map.items():
            if len(slide_indices) > 1:
                issues.extend(
                    QAIssue(
                        rule_id=self.rule_id,
                        severity=Severity.INFO,
                        slide_index=slide_idx,
                        shape_index=None,
                        message=(f"Duplicate title '{title}' found on slides {slide_indices}"),
                        auto_fixable=False,
                        suggested_fix=(
                            "Consider using unique titles or section headers"
                            " to differentiate slides"
                        ),
                    )
                    for slide_idx in slide_indices
                )

        return issues


class UnpopulatedImagePlaceholderRule:
    """QA-C-003: Detect unpopulated image placeholders.

    Empty image placeholders indicate missing visual content.

    Attributes:
        rule_id: Unique identifier "QA-C-003"
        description: Human-readable description
        auto_fixable: True - can insert default placeholder image
    """

    rule_id = "QA-C-003"
    description = "Detect unpopulated image placeholders"
    severity = "warning"
    auto_fixable = True

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for empty image placeholders.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for unpopulated image placeholders
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape is a picture placeholder
                if shape.is_placeholder:
                    ph = shape.placeholder_format
                    # Check for PICTURE placeholder type (type 18)
                    # If shape_type is PLACEHOLDER, the picture placeholder is empty
                    if (
                        ph.type == PP_PLACEHOLDER_TYPE.PICTURE
                        and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                    ):
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.WARNING,
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                message="Image placeholder is not populated",
                                auto_fixable=True,
                                suggested_fix="Insert an image or remove the placeholder",
                            )
                        )

        return issues


class PathologicalTableDimensionRule:
    """QA-C-004: Detect pathological table dimensions.

    Tables with unusual dimensions (1x1, very large, mostly empty) may indicate
    generation errors or poor content structure.

    Attributes:
        rule_id: Unique identifier "QA-C-004"
        description: Human-readable description
        auto_fixable: False - requires manual review
    """

    rule_id = "QA-C-004"
    description = "Detect pathological table dimensions"
    severity = "warning"
    auto_fixable = False

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for pathological table dimensions.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for problematic tables
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape is a table
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    if not hasattr(shape, "table"):
                        continue

                    table = shape.table  # type: ignore[attr-defined]
                    rows = len(table.rows)
                    cols = len(table.columns)

                    # Check for 1x1 tables
                    if rows == 1 and cols == 1:
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.WARNING,
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                message="Table has pathological dimensions: 1x1",
                                auto_fixable=False,
                                suggested_fix="Consider using text instead of a single-cell table",
                            )
                        )

                    # Check for very large tables (>10 rows or >8 columns)
                    elif rows > 10 or cols > 8:  # noqa: PLR2004
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.WARNING,
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                message=(
                                    f"Table has large dimensions: {rows}x{cols}"
                                    " (may be hard to read)"
                                ),
                                auto_fixable=False,
                                suggested_fix=(
                                    "Consider splitting into multiple slides or summarizing data"
                                ),
                            )
                        )

                    # Check for mostly empty tables
                    else:
                        empty_cells = 0
                        total_cells = rows * cols
                        for row in table.rows:
                            for cell in row.cells:
                                if not cell.text.strip():
                                    empty_cells += 1

                        if total_cells > 4 and empty_cells / total_cells > 0.5:  # noqa: PLR2004
                            issues.append(
                                QAIssue(
                                    rule_id=self.rule_id,
                                    severity=Severity.INFO,
                                    slide_index=slide_idx,
                                    shape_index=shape_idx,
                                    message=(
                                        f"Table is mostly empty ({empty_cells}/{total_cells} cells)"
                                    ),
                                    auto_fixable=False,
                                    suggested_fix=(
                                        "Consider removing empty cells or using a smaller table"
                                    ),
                                )
                            )

        return issues


class MissingChartDataRule:
    """QA-C-005: Detect missing chart data values.

    Charts with missing or zero data may indicate generation errors.

    Attributes:
        rule_id: Unique identifier "QA-C-005"
        description: Human-readable description
        auto_fixable: True - can populate with zero values or remove chart
    """

    rule_id = "QA-C-005"
    description = "Detect missing chart data"
    severity = "error"
    auto_fixable = True

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for charts with missing data.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for charts with missing data
        """
        issues: list[QAIssue] = []

        if not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape is a chart
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    if not hasattr(shape, "chart"):
                        continue

                    try:
                        chart = shape.chart  # type: ignore[attr-defined]

                        # Check if chart has plots
                        if not hasattr(chart, "plots") or len(chart.plots) == 0:
                            issues.append(
                                QAIssue(
                                    rule_id=self.rule_id,
                                    severity=Severity.ERROR,
                                    slide_index=slide_idx,
                                    shape_index=shape_idx,
                                    message="Chart has no data plots",
                                    auto_fixable=True,
                                    suggested_fix="Populate chart with data or remove chart",
                                )
                            )
                            continue

                        # Check if all series have zero values
                        if hasattr(chart, "series") and len(chart.series) > 0:
                            all_zero = True
                            for series in chart.series:
                                if hasattr(series, "values"):
                                    for value in series.values:
                                        if value and value != 0:
                                            all_zero = False
                                            break
                                if not all_zero:
                                    break

                            if all_zero:
                                issues.append(
                                    QAIssue(
                                        rule_id=self.rule_id,
                                        severity=Severity.WARNING,
                                        slide_index=slide_idx,
                                        shape_index=shape_idx,
                                        message="Chart has all zero values",
                                        auto_fixable=False,
                                        suggested_fix="Verify chart data is correct",
                                    )
                                )
                    except Exception:
                        # If we can't access chart data, report as potential issue
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.WARNING,
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                message="Unable to validate chart data",
                                auto_fixable=False,
                                suggested_fix="Manually verify chart data is correct",
                            )
                        )

        return issues


class SpeakerNotesVerificationRule:
    """QA-C-006: Verify speaker notes attachment when enabled.

    Ensures speaker notes are present when they should be.

    Attributes:
        rule_id: Unique identifier "QA-C-006"
        description: Human-readable description
        auto_fixable: True - can add placeholder notes
        require_notes: Whether speaker notes are required
    """

    rule_id = "QA-C-006"
    description = "Verify speaker notes attachment"
    severity = "info"
    auto_fixable = True

    def __init__(self, require_notes: bool = False) -> None:
        """Initialize speaker notes verification rule.

        Args:
            require_notes: Whether speaker notes are required (default: False)
        """
        self.require_notes = require_notes

    def validate(self, presentation: PresentationWrapper) -> list[QAIssue]:
        """Check for missing speaker notes.

        Args:
            presentation: Presentation to check

        Returns:
            List of QA issues for missing speaker notes
        """
        issues: list[QAIssue] = []

        if not self.require_notes or not presentation.is_loaded:
            return issues

        for slide_idx, slide in enumerate(presentation.prs.slides):
            # Check if slide has notes
            if not slide.has_notes_slide:
                issues.append(
                    QAIssue(
                        rule_id=self.rule_id,
                        severity=Severity.INFO,
                        slide_index=slide_idx,
                        shape_index=None,
                        message="Slide is missing speaker notes",
                        auto_fixable=True,
                        suggested_fix="Add speaker notes to provide context",
                    )
                )
            else:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if not notes_text:
                        issues.append(
                            QAIssue(
                                rule_id=self.rule_id,
                                severity=Severity.INFO,
                                slide_index=slide_idx,
                                shape_index=None,
                                message="Speaker notes are empty",
                                auto_fixable=True,
                                suggested_fix="Add speaker notes to provide context",
                            )
                        )

        return issues


# Made with Bob
