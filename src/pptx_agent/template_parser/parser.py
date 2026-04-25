"""Template parser for extracting metadata from PowerPoint templates.

This module provides functionality to parse PPTX templates and extract
information about layouts, placeholders, and their characteristics.
"""

import argparse
import json
import logging
import sys
from pathlib import Path
from typing import Any, ClassVar

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.slide import SlideLayout

from pptx_agent.cache import CacheManager
from pptx_agent.constants import (
    CHAR_WIDTH_INCHES,
    DEFAULT_PLACEHOLDER_CHARS,
    EMU_PER_INCH,
    LINE_HEIGHT_INCHES,
    MIN_PLACEHOLDER_CHARS,
)
from pptx_agent.validators.file_validator import validate_template_path

from .models import LayoutMetadata, PlaceholderMetadata, TemplateMetadata

logger = logging.getLogger(__name__)


class TemplateParser:
    """Parser for extracting metadata from PowerPoint templates.

    Extracts information about slide layouts, placeholders, and their
    properties from PPTX template files. Supports caching for improved
    performance with repeated template use.

    Attributes:
        cache_manager: Optional cache manager for template manifest caching
    """

    # Mapping of python-pptx placeholder types to string names
    PLACEHOLDER_TYPE_NAMES: ClassVar[dict[PP_PLACEHOLDER_TYPE, str]] = {
        PP_PLACEHOLDER_TYPE.TITLE: "TITLE",
        PP_PLACEHOLDER_TYPE.BODY: "BODY",
        PP_PLACEHOLDER_TYPE.CENTER_TITLE: "CENTER_TITLE",
        PP_PLACEHOLDER_TYPE.SUBTITLE: "SUBTITLE",
        PP_PLACEHOLDER_TYPE.DATE: "DATE",
        PP_PLACEHOLDER_TYPE.FOOTER: "FOOTER",
        PP_PLACEHOLDER_TYPE.HEADER: "HEADER",
        PP_PLACEHOLDER_TYPE.OBJECT: "OBJECT",
        PP_PLACEHOLDER_TYPE.CHART: "CHART",
        PP_PLACEHOLDER_TYPE.TABLE: "TABLE",
        PP_PLACEHOLDER_TYPE.ORG_CHART: "ORG_CHART",
        PP_PLACEHOLDER_TYPE.MEDIA_CLIP: "MEDIA_CLIP",
        PP_PLACEHOLDER_TYPE.SLIDE_NUMBER: "SLIDE_NUMBER",
        PP_PLACEHOLDER_TYPE.PICTURE: "PICTURE",
        PP_PLACEHOLDER_TYPE.BITMAP: "BITMAP",
        PP_PLACEHOLDER_TYPE.SLIDE_IMAGE: "SLIDE_IMAGE",
        PP_PLACEHOLDER_TYPE.VERTICAL_BODY: "VERTICAL_BODY",
        PP_PLACEHOLDER_TYPE.VERTICAL_TITLE: "VERTICAL_TITLE",
        PP_PLACEHOLDER_TYPE.VERTICAL_OBJECT: "VERTICAL_OBJECT",
        PP_PLACEHOLDER_TYPE.MIXED: "MIXED",
    }

    def __init__(self, cache_manager: CacheManager | None = None) -> None:
        """Initialize template parser.

        Args:
            cache_manager: Optional cache manager for manifest caching.
                          If None, caching is disabled.
        """
        self.cache_manager = cache_manager

    def parse_template(
        self,
        template_path: str,
        use_cache: bool = True,
    ) -> TemplateMetadata:
        """Parse a PowerPoint template and extract metadata.

        Args:
            template_path: Path to the PPTX template file

        Returns:
            TemplateMetadata containing all extracted information

        Raises:
            FileNotFoundError: If template file does not exist
            InvalidFileError: If file is not a valid PPTX file
            FileSizeLimitError: If file size exceeds limits
            CompressionRatioError: If compression ratio is suspicious


        """
        # Validate template path (includes existence, extension, symlinks, ZIP structure)
        validated_path = validate_template_path(template_path)

        # Load presentation using validated path
        prs = Presentation(str(validated_path))

        # Extract layouts
        layouts = []
        for slide_layout in prs.slide_layouts:
            layout_metadata = self._parse_layout(slide_layout)
            layouts.append(layout_metadata)

        return TemplateMetadata(
            template_path=str(validated_path),
            layouts=layouts,
        )

    def _parse_layout(self, slide_layout: SlideLayout) -> LayoutMetadata:
        """Parse a slide layout and extract its metadata.

        Args:
            slide_layout: python-pptx SlideLayout object

        Returns:
            LayoutMetadata with layout information
        """
        layout_name = slide_layout.name

        # Extract placeholders
        placeholders = []
        has_smartart = False

        for shape in slide_layout.shapes:
            if shape.is_placeholder:
                placeholder_metadata = self._parse_placeholder(shape)
                placeholders.append(placeholder_metadata)

            # Check if this shape is a SmartArt diagram
            if self._is_smartart_shape(shape):
                has_smartart = True

        return LayoutMetadata(
            name=layout_name,
            placeholders=placeholders,
            has_smartart=has_smartart,
        )

    def _parse_placeholder(self, shape: Any) -> PlaceholderMetadata:
        """Parse a placeholder shape and extract its metadata.

        Args:
            shape: python-pptx Shape object with is_placeholder=True

        Returns:
            PlaceholderMetadata with placeholder information
        """
        # Get placeholder name
        placeholder_name = shape.name

        # Get placeholder type
        placeholder_type = self._get_placeholder_type(shape)

        # Estimate max characters based on placeholder size
        max_chars = self._estimate_max_chars(shape)

        return PlaceholderMetadata(
            name=placeholder_name,
            type=placeholder_type,
            max_chars=max_chars,
        )

    def _get_placeholder_type(self, shape: Any) -> str:
        """Get the type of a placeholder shape.

        Args:
            shape: python-pptx Shape object with is_placeholder=True

        Returns:
            String representation of placeholder type
        """
        try:
            placeholder = shape.placeholder_format
            ph_type = placeholder.type
            return self.PLACEHOLDER_TYPE_NAMES.get(ph_type, "UNKNOWN")
        except (AttributeError, KeyError) as e:
            logger.debug(
                "Failed to get placeholder type: %s. Returning UNKNOWN.",
                e,
            )
            return "UNKNOWN"

    def _estimate_max_chars(self, shape: Any) -> int:
        """Estimate maximum character capacity of a placeholder.

        This is a rough estimate based on placeholder dimensions.
        Assumes average character width and line spacing.

        Args:
            shape: python-pptx Shape object

        Returns:
            Estimated maximum number of characters
        """
        try:
            # Get shape dimensions (in EMU - English Metric Units)
            width_emu = shape.width
            height_emu = shape.height

            # Convert to inches
            width_inches = width_emu / EMU_PER_INCH
            height_inches = height_emu / EMU_PER_INCH

            # Rough estimate using average character/line dimensions
            chars_per_line = int(width_inches / CHAR_WIDTH_INCHES)
            lines = int(height_inches / LINE_HEIGHT_INCHES)

            # Estimate total characters
            max_chars = chars_per_line * lines

            # Ensure minimum character capacity
            return max(max_chars, MIN_PLACEHOLDER_CHARS)

        except (AttributeError, TypeError, ValueError, ZeroDivisionError) as e:
            logger.debug(
                "Failed to estimate max chars: %s. Returning default value.",
                e,
            )
            # Default estimate if dimensions unavailable
            return DEFAULT_PLACEHOLDER_CHARS

    def _is_smartart_shape(self, shape: Any) -> bool:
        """Check if a shape is a SmartArt diagram.

        SmartArt diagrams are represented as graphicFrame elements in PowerPoint XML.
        This method detects them by checking the element tag.

        Args:
            shape: python-pptx Shape object

        Returns:
            True if shape is a SmartArt diagram, False otherwise
        """
        try:
            # SmartArt is typically a graphicFrame element
            if hasattr(shape, "element") and hasattr(shape.element, "tag"):
                return shape.element.tag.endswith("}graphicFrame")
            return False  # noqa: TRY300
        except Exception:
            return False


def main() -> None:
    """CLI entry point for TemplateParser."""
    parser = argparse.ArgumentParser(
        description="Extract layout and placeholder metadata from a PowerPoint template."
    )
    parser.add_argument(
        "--template",
        required=True,
        help="Path to the input .pptx template file",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Path to the output JSON manifest file",
    )

    args = parser.parse_args()

    try:
        template_parser = TemplateParser()
        metadata = template_parser.parse_template(args.template)

        # Serialize to JSON and write to output file
        output_path = Path(args.output)
        # Using parents=True as a convenience for CLI users to auto-create output directories
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with output_path.open("w", encoding="utf-8") as f:
            json.dump(metadata.model_dump(), f, indent=2, ensure_ascii=False)

        logger.info("Successfully generated template manifest at: %s", args.output)

    except Exception:
        logger.exception("Error parsing template")
        sys.exit(1)


if __name__ == "__main__":
    main()
