"""Fix strategies for placeholder population issues."""

from typing import Any

from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from pptx_agent.fixer.schemas import FixResult, FixStatus
from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.schemas import QAIssue


class PlaceholderPopulationStrategy:
    """Fix empty placeholders by populating from outline data.

    This strategy populates empty title and body placeholders using
    content from the presentation outline or slide schema.

    Attributes:
        rule_id: QA rule this strategy addresses
        description: Human-readable description
    """

    rule_id = "QA-L-002"
    description = "Populate empty placeholders from outline"

    def apply(
        self,
        issue: QAIssue,
        presentation: PresentationWrapper,
        outline: Any | None = None,
    ) -> FixResult:
        """Apply placeholder population fix.

        Args:
            issue: QA issue describing the empty placeholder
            presentation: Presentation wrapper with the issue
            outline: Optional outline data with slide content

        Returns:
            Fix result describing the outcome
        """
        try:
            # Check if outline data is available
            if outline is None:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="No outline data available for population",
                )

            # Validate indices
            if issue.slide_index >= len(presentation.prs.slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid slide index: {issue.slide_index}",
                )

            slide = presentation.prs.slides[issue.slide_index]

            if issue.shape_index is None or issue.shape_index >= len(slide.shapes):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid shape index: {issue.shape_index}",
                )

            shape = slide.shapes[issue.shape_index]

            # Check if shape is a placeholder
            if not hasattr(shape, "is_placeholder") or not shape.is_placeholder:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Shape is not a placeholder",
                )

            # Check if placeholder has text frame
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Placeholder has no text frame",
                )

            text_frame = getattr(shape, "text_frame")

            # Check if already has content
            if text_frame.text.strip():
                return FixResult(
                    issue=issue,
                    status=FixStatus.SKIPPED,
                    message="Placeholder already has content",
                )

            # Get placeholder type
            placeholder_type = None
            if hasattr(shape, "placeholder_format"):
                placeholder_type = shape.placeholder_format.type

            # Get content from outline
            outline_slides = getattr(outline, "slides", [])
            if issue.slide_index >= len(outline_slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="No outline data for this slide",
                )

            slide_data = outline_slides[issue.slide_index]
            changes_made: list[str] = []

            # Populate based on placeholder type
            if placeholder_type == PP_PLACEHOLDER_TYPE.TITLE:
                title = slide_data.get("title", "")
                if title:
                    text_frame.text = title
                    changes_made.append(f"Populated title: {title}")
                else:
                    return FixResult(
                        issue=issue,
                        status=FixStatus.FAILED,
                        message="No title available in outline",
                    )

            elif placeholder_type == PP_PLACEHOLDER_TYPE.BODY:
                bullets = slide_data.get("bullets", [])
                if bullets:
                    # Clear existing content and add bullet points
                    if hasattr(text_frame, "clear"):
                        text_frame.clear()
                        # Add bullet points
                        for bullet in bullets:
                            if hasattr(text_frame, "add_paragraph"):
                                p = text_frame.add_paragraph()
                                p.text = bullet
                                p.level = 0
                        changes_made.append(f"Populated {len(bullets)} bullet points")
                    else:
                        # Fallback: just set text with newlines
                        text_frame.text = "\n".join(f"• {bullet}" for bullet in bullets)
                        changes_made.append(f"Populated {len(bullets)} bullet points")
                else:
                    return FixResult(
                        issue=issue,
                        status=FixStatus.FAILED,
                        message="No bullet points available in outline",
                    )

            else:
                # Generic placeholder - try to populate with any available text
                content = slide_data.get("content", "") or slide_data.get("title", "")
                if content:
                    text_frame.text = content
                    changes_made.append("Populated placeholder with content")
                else:
                    return FixResult(
                        issue=issue,
                        status=FixStatus.FAILED,
                        message="No content available for placeholder",
                    )

            if not changes_made:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Could not populate placeholder",
                )

            return FixResult(
                issue=issue,
                status=FixStatus.SUCCESS,
                message="Placeholder populated successfully",
                changes_made=changes_made,
            )

        except Exception as exc:
            return FixResult(
                issue=issue,
                status=FixStatus.FAILED,
                message=f"Placeholder population failed: {exc}",
            )


# Made with Bob
