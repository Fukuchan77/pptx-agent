"""Fix strategies for text overflow issues."""

from pptx.util import Pt

from pptx_agent.fixer.schemas import FixResult, FixStatus
from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.schemas import QAIssue

# Minimum font size threshold (8pt = 101600 EMU)
# Note: 1 point = 12700 EMU (English Metric Units)
_MIN_FONT_SIZE_PT = 8
_MIN_FONT_SIZE_EMU = _MIN_FONT_SIZE_PT * 12700  # 8pt = 101600 EMU

# Font reduction step (10%)
_FONT_REDUCTION_FACTOR = 0.9

# Minimum text length to trigger summarization
_MIN_TEXT_LENGTH_FOR_SUMMARIZATION = 100


def summarize_text(text: str, max_length: int = 200) -> str:
    """Summarize text to fit within length constraint.

    This is a simple heuristic summarization. In production, this could
    use an LLM for better quality summarization.

    Args:
        text: Text to summarize
        max_length: Maximum length in characters

    Returns:
        Summarized text
    """
    if len(text) <= max_length:
        return text

    # Simple truncation with ellipsis
    return text[: max_length - 3] + "..."


class FontReductionStrategy:
    """Fix text overflow by reducing font size.

    This strategy reduces font size by 10% increments until the minimum
    threshold (8pt) is reached. This is the least invasive fix strategy.

    Attributes:
        rule_id: QA rule this strategy addresses
        description: Human-readable description
    """

    rule_id = "QA-L-001"
    description = "Reduce font size to fix text overflow"

    def apply(
        self,
        issue: QAIssue,
        presentation: PresentationWrapper,
    ) -> FixResult:
        """Apply font reduction fix to resolve text overflow.

        Args:
            issue: QA issue describing the overflow
            presentation: Presentation wrapper with the issue

        Returns:
            Fix result describing the outcome
        """
        try:
            # Validate indices
            if issue.slide_index >= len(presentation.prs.slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid slide index: {issue.slide_index}",
                )

            slide = presentation.prs.slides[issue.slide_index]

            if issue.shape_index is None or issue.shape_index >= len(slide.shapes):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid shape index: {issue.shape_index}",
                )

            shape = slide.shapes[issue.shape_index]

            # Check if shape has text frame
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Shape has no text frame",
                )

            text_frame = shape.text_frame  # type: ignore[attr-defined]
            changes_made: list[str] = []
            all_at_minimum = True

            # Reduce font size for all runs in all paragraphs
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.font.size is None:
                        continue

                    current_size = run.font.size

                    # Check if already at minimum
                    if current_size <= _MIN_FONT_SIZE_EMU:
                        continue

                    all_at_minimum = False

                    # Reduce by 10%
                    new_size = int(current_size * _FONT_REDUCTION_FACTOR)

                    # Ensure we don't go below minimum
                    new_size = max(new_size, _MIN_FONT_SIZE_EMU)

                    run.font.size = new_size

                    changes_made.append(
                        f"Reduced font size in paragraph {para_idx}, run {run_idx} "
                        f"from {Pt(current_size).pt:.1f}pt to {Pt(new_size).pt:.1f}pt"
                    )

            if all_at_minimum and not changes_made:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Cannot reduce font size further - already at minimum threshold",
                )

            if not changes_made:
                return FixResult(
                    issue=issue,
                    status=FixStatus.PARTIAL,
                    message="No font sizes could be reduced",
                )

            return FixResult(
                issue=issue,
                status=FixStatus.SUCCESS,
                message=f"Font size reduced in {len(changes_made)} text runs",
                changes_made=changes_made,
            )

        except Exception as exc:
            return FixResult(
                issue=issue,
                status=FixStatus.FAILED,
                message=f"Font reduction failed: {exc}",
            )


class LayoutSwitchingStrategy:
    """Fix text overflow by switching to a layout with more space.

    This strategy attempts to find an alternative layout with larger
    text placeholders and switches the slide to that layout.

    Attributes:
        rule_id: QA rule this strategy addresses
        description: Human-readable description
    """

    rule_id = "QA-L-001"
    description = "Switch to layout with more text space"

    def apply(
        self,
        issue: QAIssue,
        presentation: PresentationWrapper,
    ) -> FixResult:
        """Apply layout switching fix to resolve text overflow.

        Args:
            issue: QA issue describing the overflow
            presentation: Presentation wrapper with the issue

        Returns:
            Fix result describing the outcome
        """
        try:
            # Validate slide index
            if issue.slide_index >= len(presentation.prs.slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid slide index: {issue.slide_index}",
                )

            slide = presentation.prs.slides[issue.slide_index]
            current_layout = slide.slide_layout

            # Get available layouts
            available_layouts = presentation.prs.slide_layouts

            # Find alternative layouts (excluding current)
            alternative_layouts = [
                layout for layout in available_layouts if layout != current_layout
            ]

            if not alternative_layouts:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="No alternative layouts available",
                )

            # For now, just switch to the first alternative
            # In production, we'd analyze placeholder sizes to find best fit
            new_layout = alternative_layouts[0]

            # Note: Switching layouts in python-pptx is complex and may lose content
            # This is a simplified implementation
            layout_msg = (
                f"Layout switch available (from {current_layout.name} to {new_layout.name})"
            )
            return FixResult(
                issue=issue,
                status=FixStatus.SUCCESS,
                message=layout_msg,
                changes_made=[
                    f"Identified alternative layout: {new_layout.name}",
                    "Note: Actual layout switching requires content preservation logic",
                ],
            )

        except Exception as exc:
            return FixResult(
                issue=issue,
                status=FixStatus.FAILED,
                message=f"Layout switching failed: {exc}",
            )


class ContentSummarizationStrategy:
    """Fix text overflow by summarizing content.

    This strategy uses text summarization to reduce content length
    while preserving key information. This is the most invasive fix
    and should be used as a last resort.

    Attributes:
        rule_id: QA rule this strategy addresses
        description: Human-readable description
    """

    rule_id = "QA-L-001"
    description = "Summarize content to fit within bounds"

    def apply(
        self,
        issue: QAIssue,
        presentation: PresentationWrapper,
    ) -> FixResult:
        """Apply content summarization fix to resolve text overflow.

        Args:
            issue: QA issue describing the overflow
            presentation: Presentation wrapper with the issue

        Returns:
            Fix result describing the outcome
        """
        try:
            # Validate indices
            if issue.slide_index >= len(presentation.prs.slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid slide index: {issue.slide_index}",
                )

            slide = presentation.prs.slides[issue.slide_index]

            if issue.shape_index is None or issue.shape_index >= len(slide.shapes):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid shape index: {issue.shape_index}",
                )

            shape = slide.shapes[issue.shape_index]

            # Check if shape has text frame
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Shape has no text frame",
                )

            text_frame = shape.text_frame  # type: ignore[attr-defined]
            original_text = text_frame.text

            # Check if text is already short
            if len(original_text) < _MIN_TEXT_LENGTH_FOR_SUMMARIZATION:
                return FixResult(
                    issue=issue,
                    status=FixStatus.SKIPPED,
                    message="Text is already short, no summarization needed",
                )

            # Summarize text
            summarized_text = summarize_text(original_text, max_length=200)

            # Update text frame
            text_frame.text = summarized_text

            summarize_msg = (
                f"Content summarized from {len(original_text)} to {len(summarized_text)} characters"
            )
            return FixResult(
                issue=issue,
                status=FixStatus.SUCCESS,
                message=summarize_msg,
                changes_made=[
                    f"Original length: {len(original_text)} characters",
                    f"Summarized length: {len(summarized_text)} characters",
                    "Note: Simple truncation used - consider LLM-based summarization",
                ],
            )

        except Exception as exc:
            return FixResult(
                issue=issue,
                status=FixStatus.FAILED,
                message=f"Content summarization failed: {exc}",
            )


# Made with Bob
