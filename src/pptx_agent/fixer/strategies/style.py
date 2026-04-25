"""Fix strategies for style violation issues."""

from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt

from pptx_agent.fixer.schemas import FixResult, FixStatus
from pptx_agent.pptx_wrapper.presentation import PresentationWrapper
from pptx_agent.qa.schemas import QAIssue


class StyleResetStrategy:
    """Fix style violations by resetting to master slide defaults.

    This strategy resets font properties, colors, and bullet indentation
    to match the master slide layout defaults.

    Attributes:
        rule_id: QA rule this strategy addresses
        description: Human-readable description
    """

    rule_id = "QA-S-001"
    description = "Reset style to master slide defaults"

    def apply(
        self,
        issue: QAIssue,
        presentation: PresentationWrapper,
        _outline: None = None,
    ) -> FixResult:
        """Apply style reset fix.

        Args:
            issue: QA issue describing the style violation
            presentation: Presentation wrapper with the issue
            outline: Not used for style fixes

        Returns:
            Fix result describing the outcome
        """
        try:
            # Validate indices
            if issue.slide_index >= len(presentation.prs.slides):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid slide index: {issue.slide_index}",
                )

            slide = presentation.prs.slides[issue.slide_index]

            if issue.shape_index is None or issue.shape_index >= len(slide.shapes):
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message=f"Invalid shape index: {issue.shape_index}",
                )

            shape = slide.shapes[issue.shape_index]

            # Check if shape has text frame
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                return FixResult(
                    issue=issue,
                    status=FixStatus.FAILED,
                    message="Shape has no text frame",
                )

            text_frame = shape.text_frame  # type: ignore[attr-defined]
            changes_made: list[str] = []

            # Get master slide defaults
            slide_layout = slide.slide_layout
            master = slide_layout.slide_master

            # Default font properties from master
            default_font_name = "Calibri"  # PowerPoint default
            default_font_size = Pt(18)  # Common default
            default_color_rgb = (0, 0, 0)  # Black

            # Try to get actual defaults from master
            if master and hasattr(master, "shapes"):
                try:
                    for master_shape in master.shapes:
                        if hasattr(master_shape, "text_frame") and master_shape.has_text_frame:
                            master_tf = master_shape.text_frame  # type: ignore[attr-defined]
                            if hasattr(master_tf, "paragraphs") and master_tf.paragraphs:
                                first_para = master_tf.paragraphs[0]
                                if hasattr(first_para, "runs") and first_para.runs:
                                    first_run = first_para.runs[0]
                                    if hasattr(first_run, "font"):
                                        if hasattr(first_run.font, "name") and first_run.font.name:
                                            default_font_name = first_run.font.name
                                        if hasattr(first_run.font, "size") and first_run.font.size:
                                            default_font_size = first_run.font.size
                                        if (
                                            hasattr(first_run.font, "color")
                                            and hasattr(first_run.font.color, "rgb")
                                            and first_run.font.color.rgb
                                        ):
                                            default_color_rgb = first_run.font.color.rgb
                                    break
                except (TypeError, AttributeError):
                    # If master.shapes is not iterable or has issues, use defaults
                    pass

            # Reset all paragraphs
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                # Reset paragraph-level properties
                if (
                    hasattr(paragraph, "alignment")
                    and paragraph.alignment != PP_PARAGRAPH_ALIGNMENT.LEFT
                ):
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    changes_made.append(f"Reset paragraph {para_idx} alignment to left")

                # Reset bullet indentation
                if hasattr(paragraph, "level") and paragraph.level != 0:
                    original_level = paragraph.level
                    paragraph.level = 0
                    changes_made.append(
                        f"Reset paragraph {para_idx} indent from level {original_level} to 0"
                    )

                # Reset run-level properties
                if not hasattr(paragraph, "runs"):
                    continue

                for run_idx, run in enumerate(paragraph.runs):
                    run_changes: list[str] = []

                    # Reset font name
                    if run.font.name != default_font_name:
                        run.font.name = default_font_name
                        run_changes.append(f"font={default_font_name}")

                    # Reset font size
                    if run.font.size != default_font_size:
                        run.font.size = default_font_size
                        run_changes.append(f"size={default_font_size.pt}pt")

                    # Reset font color
                    if (
                        hasattr(run.font.color, "rgb")
                        and run.font.color.rgb
                        and run.font.color.rgb != default_color_rgb
                    ):
                        run.font.color.rgb = default_color_rgb
                        run_changes.append(f"color=RGB{default_color_rgb}")

                    # Reset bold
                    if run.font.bold:
                        run.font.bold = False
                        run_changes.append("bold=False")

                    # Reset italic
                    if run.font.italic:
                        run.font.italic = False
                        run_changes.append("italic=False")

                    # Reset underline
                    if run.font.underline:
                        run.font.underline = False
                        run_changes.append("underline=False")

                    if run_changes:
                        changes_made.append(
                            f"Reset paragraph {para_idx} run {run_idx}: {', '.join(run_changes)}"
                        )

            if not changes_made:
                return FixResult(
                    issue=issue,
                    status=FixStatus.SKIPPED,
                    message="No style violations found to reset",
                )

            return FixResult(
                issue=issue,
                status=FixStatus.SUCCESS,
                message="Style reset to master defaults",
                changes_made=changes_made,
            )

        except Exception as exc:
            return FixResult(
                issue=issue,
                status=FixStatus.FAILED,
                message=f"Style reset failed: {exc}",
            )


# Made with Bob
